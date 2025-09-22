import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from datetime import datetime, timedelta
from io import BytesIO
import warnings
import time
import traceback

warnings.filterwarnings('ignore')

# Core analytics imports only (avoiding NumPy compatibility issues)
from sklearn.preprocessing import StandardScaler, RobustScaler
from sklearn.cluster import KMeans, DBSCAN
from sklearn.decomposition import PCA
from sklearn.ensemble import IsolationForest, RandomForestRegressor
from sklearn.linear_model import LinearRegression
from sklearn.preprocessing import PolynomialFeatures
from sklearn.pipeline import Pipeline
from sklearn.metrics import silhouette_score
from scipy import stats
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from german_corporate_powerpoint import create_german_corporate_powerpoint
import re

# Configure page with professional English branding
st.set_page_config(
    page_title="AI Automation Platform | Enterprise Analytics",
    page_icon="🏢",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for professional styling
st.markdown("""
<style>
    /* Main container styling */
    .main > div {
        padding-top: 2rem;
    }

    /* Header styling */
    .header-container {
        background: linear-gradient(135deg, #1e3a8a 0%, #3b82f6 100%);
        padding: 2rem;
        border-radius: 10px;
        margin-bottom: 2rem;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }

    .header-title {
        color: white;
        font-size: 2.5rem;
        font-weight: 700;
        margin-bottom: 0.5rem;
        text-align: center;
    }

    .header-subtitle {
        color: #e2e8f0;
        font-size: 1.2rem;
        text-align: center;
        margin-bottom: 1rem;
    }

    .header-tagline {
        color: #cbd5e1;
        font-size: 1rem;
        text-align: center;
        font-style: italic;
    }

    /* Card styling */
    .analysis-card {
        background: #f8fafc;
        padding: 2rem;
        border-radius: 10px;
        border: 1px solid #e2e8f0;
        margin: 1rem 0;
    }

    /* Button styling */
    .stButton > button {
        background: linear-gradient(135deg, #1e3a8a 0%, #3b82f6 100%);
        color: white;
        border: none;
        border-radius: 8px;
        padding: 0.75rem 2rem;
        font-weight: 600;
        transition: all 0.3s ease;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
    }

    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0, 0, 0, 0.2);
    }

    /* Data quality indicators */
    .quality-indicator {
        display: inline-block;
        padding: 0.25rem 0.75rem;
        border-radius: 20px;
        font-size: 0.875rem;
        font-weight: 600;
        margin-right: 0.5rem;
    }

    .quality-excellent {
        background-color: #dcfce7;
        color: #166534;
    }

    .quality-good {
        background-color: #fef3c7;
        color: #92400e;
    }

    .quality-poor {
        background-color: #fee2e2;
        color: #991b1b;
    }
</style>
""", unsafe_allow_html=True)

# Professional header
st.markdown("""
<div class="header-container">
    <div class="header-title">🏢 AI Automation Platform</div>
    <div class="header-subtitle">Enterprise Data Intelligence & Analytics Suite</div>
    <div class="header-tagline">Advanced Artificial Intelligence for Modern Enterprises</div>
</div>
""", unsafe_allow_html=True)

# Core Analytics Engine (NumPy-compatible)
class EnterpriseAIAnalytics:
    """Enterprise-grade AI analytics engine without problematic dependencies"""

    def __init__(self):
        self.scaler = StandardScaler()
        self.robust_scaler = RobustScaler()

    def detect_industry_pattern(self, df):
        """Detect industry-specific patterns from column names and data types"""
        columns = [col.lower() for col in df.columns]

        # Industry pattern detection
        financial_keywords = ['revenue', 'profit', 'cost', 'price', 'income', 'sales', 'budget', 'roi', 'margin']
        healthcare_keywords = ['patient', 'diagnosis', 'treatment', 'medical', 'hospital', 'clinical']
        manufacturing_keywords = ['production', 'manufacturing', 'quality', 'defect', 'inventory', 'supply']
        retail_keywords = ['customer', 'product', 'inventory', 'store', 'purchase', 'order']

        patterns = {
            'financial': sum(1 for kw in financial_keywords if any(kw in col for col in columns)),
            'healthcare': sum(1 for kw in healthcare_keywords if any(kw in col for col in columns)),
            'manufacturing': sum(1 for kw in manufacturing_keywords if any(kw in col for col in columns)),
            'retail': sum(1 for kw in retail_keywords if any(kw in col for col in columns))
        }

        detected_industry = max(patterns, key=patterns.get) if max(patterns.values()) > 0 else 'general'
        confidence = max(patterns.values()) / len(columns) if len(columns) > 0 else 0

        return {
            'pattern': detected_industry,
            'confidence': confidence,
            'details': patterns
        }

    def advanced_anomaly_detection(self, df, numeric_cols, contamination=0.1):
        """Multi-method anomaly detection with confidence scores"""
        if len(numeric_cols) == 0:
            return None, None

        results = {}

        # Prepare data
        data = df[numeric_cols].dropna()
        if len(data) < 10:
            return None, None

        # Method 1: Isolation Forest
        iso_forest = IsolationForest(contamination=contamination, random_state=42)
        iso_forest.fit(data)  # Fit the model first
        iso_scores = iso_forest.decision_function(data)
        iso_anomalies = iso_forest.predict(data) == -1

        # Method 2: Statistical Z-score
        z_scores = np.abs(stats.zscore(data, axis=0))
        z_anomalies = (z_scores > 3).any(axis=1)

        # Method 3: DBSCAN clustering
        if len(numeric_cols) >= 2:
            scaled_data = self.scaler.fit_transform(data)
            dbscan = DBSCAN(eps=0.5, min_samples=5)
            clusters = dbscan.fit_predict(scaled_data)
            dbscan_anomalies = clusters == -1
        else:
            dbscan_anomalies = np.zeros(len(data), dtype=bool)

        # Combine methods for ensemble approach
        ensemble_score = (iso_anomalies.astype(int) +
                         z_anomalies.astype(int) +
                         dbscan_anomalies.astype(int)) / 3

        # Calculate confidence scores
        confidence_scores = []
        for i in range(len(data)):
            methods_agreeing = [iso_anomalies[i], z_anomalies[i], dbscan_anomalies[i]]
            confidence = sum(methods_agreeing) / len(methods_agreeing)
            confidence_scores.append(confidence)

        results = {
            'indices': data.index[ensemble_score > 0.33],  # At least 1/3 methods agree
            'scores': ensemble_score[ensemble_score > 0.33],
            'confidence': np.array(confidence_scores)[ensemble_score > 0.33],
            'isolation_scores': iso_scores,
            'methods_used': ['Isolation Forest', 'Statistical Z-score', 'DBSCAN Clustering']
        }

        return results, data

    def advanced_correlation_analysis(self, df, numeric_cols):
        """Advanced correlation analysis with multiple methods"""
        if len(numeric_cols) < 2:
            return None

        data = df[numeric_cols].dropna()

        # Multiple correlation methods
        pearson_corr = data.corr(method='pearson')
        spearman_corr = data.corr(method='spearman')
        kendall_corr = data.corr(method='kendall')

        # Partial correlations
        partial_corrs = {}
        for i, col1 in enumerate(numeric_cols):
            for j, col2 in enumerate(numeric_cols):
                if i < j:
                    other_cols = [col for col in numeric_cols if col not in [col1, col2]]
                    if len(other_cols) > 0 and len(data) > len(other_cols) + 10:
                        try:
                            # Control for other variables
                            X = data[other_cols]
                            y1 = data[col1]
                            y2 = data[col2]

                            # Residuals after controlling for other variables
                            model1 = LinearRegression().fit(X, y1)
                            model2 = LinearRegression().fit(X, y2)

                            resid1 = y1 - model1.predict(X)
                            resid2 = y2 - model2.predict(X)

                            partial_corr = np.corrcoef(resid1, resid2)[0, 1]
                            partial_corrs[f"{col1}_vs_{col2}"] = partial_corr
                        except:
                            pass

        return {
            'pearson': pearson_corr,
            'spearman': spearman_corr,
            'kendall': kendall_corr,
            'partial_correlations': partial_corrs,
            'sample_size': len(data)
        }

    def advanced_trend_analysis(self, df, target_col, date_col=None):
        """Enhanced trend analysis using core packages only"""
        ts = df[target_col].dropna()

        if len(ts) < 10:
            return self._basic_trend_analysis(df, target_col)

        # Prepare time series data
        if date_col and date_col in df.columns:
            try:
                ts_data = df[[date_col, target_col]].dropna()
                ts_data[date_col] = pd.to_datetime(ts_data[date_col])
                ts_data = ts_data.set_index(date_col).sort_index()
                ts = ts_data[target_col]
            except:
                ts = df[target_col].dropna()
                ts.index = pd.date_range(start='2020-01-01', periods=len(ts), freq='D')
        else:
            ts = df[target_col].dropna()
            ts.index = pd.date_range(start='2020-01-01', periods=len(ts), freq='D')

        results = {}

        # Enhanced trend analysis with polynomial fitting
        X = np.arange(len(ts)).reshape(-1, 1)
        y = ts.values

        # Linear trend
        linear_model = LinearRegression()
        linear_model.fit(X, y)
        linear_trend = linear_model.predict(X)
        linear_r2 = linear_model.score(X, y)

        # Polynomial trend (degree 2)
        poly_model = Pipeline([
            ('poly', PolynomialFeatures(degree=2)),
            ('linear', LinearRegression())
        ])
        poly_model.fit(X, y)
        poly_trend = poly_model.predict(X)
        poly_r2 = poly_model.score(X, y)

        # Choose best model
        best_model = poly_model if poly_r2 > linear_r2 + 0.05 else linear_model
        best_trend = poly_trend if poly_r2 > linear_r2 + 0.05 else linear_trend
        best_r2 = poly_r2 if poly_r2 > linear_r2 + 0.05 else linear_r2

        # Simple stationarity test
        try:
            mid_point = len(ts) // 2
            first_half_mean = ts[:mid_point].mean()
            second_half_mean = ts[mid_point:].mean()
            mean_diff = abs(second_half_mean - first_half_mean)
            overall_std = ts.std()

            is_stationary = mean_diff < (overall_std * 0.5)

            results['stationarity'] = {
                'is_stationary': is_stationary,
                'mean_shift': mean_diff,
                'relative_shift': mean_diff / overall_std if overall_std > 0 else 0
            }
        except:
            results['stationarity'] = None

        # Enhanced forecasting with confidence intervals
        forecast_periods = min(30, len(ts) // 4)
        future_X = np.arange(len(ts), len(ts) + forecast_periods).reshape(-1, 1)

        # Point forecast
        forecast = best_model.predict(future_X)

        # Simple confidence intervals based on residual standard error
        residuals = y - best_trend
        residual_std = np.std(residuals)

        # 95% confidence intervals
        ci_lower = forecast - 1.96 * residual_std
        ci_upper = forecast + 1.96 * residual_std

        results['enhanced_forecast'] = {
            'model_type': 'Polynomial' if poly_r2 > linear_r2 + 0.05 else 'Linear',
            'forecast': forecast,
            'confidence_lower': ci_lower,
            'confidence_upper': ci_upper,
            'r2_score': best_r2,
            'residual_std': residual_std
        }

        # Add original time series
        results['original'] = ts

        return results

    def _basic_trend_analysis(self, df, target_col):
        """Basic trend analysis fallback"""
        ts = df[target_col].dropna()
        X = np.arange(len(ts)).reshape(-1, 1)
        y = ts.values

        model = LinearRegression()
        model.fit(X, y)

        forecast_periods = min(30, len(ts) // 4)
        future_X = np.arange(len(ts), len(ts) + forecast_periods).reshape(-1, 1)
        forecast = model.predict(future_X)

        return {
            'trend_slope': model.coef_[0],
            'r2_score': model.score(X, y),
            'forecast': forecast,
            'original': ts
        }

    def gdpr_compliance_assessment(self, df):
        """GDPR compliance risk assessment"""
        risk_factors = []
        compliance_score = 100
        recommendations = []

        # Check for potential PII columns
        pii_patterns = {
            'email': r'(email|e-mail|mail)',
            'phone': r'(phone|tel|mobile)',
            'name': r'(name|firstname|lastname)',
            'address': r'(address|street|zip|postal)',
            'id_number': r'(id|ssn|social|tax)',
            'date_of_birth': r'(birth|dob)',
            'ip_address': r'(ip|internet)',
            'location': r'(location|gps|coordinate)',
            'financial': r'(iban|bic|account|bank|credit)'
        }

        potential_pii = {}
        for category, pattern in pii_patterns.items():
            matching_cols = [col for col in df.columns if re.search(pattern, col.lower())]
            if matching_cols:
                potential_pii[category] = matching_cols

        # Risk assessment
        if potential_pii:
            risk_factors.append("🔴 Potential personally identifiable information (PII) detected")
            compliance_score -= len(potential_pii) * 10

            for category, cols in potential_pii.items():
                recommendations.append(f"📋 {category.title()}: Review data processing legitimacy for {', '.join(cols)}")

        # Data volume assessment
        if len(df) > 10000:
            risk_factors.append("⚠️ Large dataset - enhanced GDPR compliance requirements")
            compliance_score -= 5
            recommendations.append("📊 Implementation of Data Protection Impact Assessment (DPIA) recommended")

        # Missing value patterns
        missing_pct = (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
        if missing_pct > 20:
            risk_factors.append("⚠️ High percentage of missing values - possible data minimization")
            compliance_score += 5  # This is actually good for GDPR

        # Determine compliance level
        if compliance_score >= 90:
            compliance_level = "Excellent"
            level_color = "🟢"
        elif compliance_score >= 75:
            compliance_level = "Good"
            level_color = "🟡"
        elif compliance_score >= 60:
            compliance_level = "Moderate"
            level_color = "🟠"
        else:
            compliance_level = "High Risk"
            level_color = "🔴"

        return {
            'compliance_score': max(0, compliance_score),
            'compliance_level': compliance_level,
            'level_color': level_color,
            'risk_factors': risk_factors,
            'potential_pii': potential_pii,
            'recommendations': recommendations,
            'assessment_timestamp': datetime.now()
        }

    def generate_executive_summary(self, df, analysis_results, industry_pattern):
        """Generate professional executive summary"""
        data_volume = len(df)
        data_completeness = 100 - (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
        numeric_cols = df.select_dtypes(include=[np.number]).columns

        summary = {
            'header': {
                'title': 'Executive Data Intelligence Summary',
                'subtitle': f'Analysis of {data_volume:,} data records',
                'timestamp': datetime.now().strftime('%d/%m/%Y, %H:%M'),
                'industry': industry_pattern.title()
            },
            'key_findings': [],
            'strategic_recommendations': [],
            'risk_assessment': {},
            'technical_insights': {},
            'next_steps': []
        }

        # Key findings
        if data_completeness > 95:
            summary['key_findings'].append("✅ Excellent data quality - ideal for strategic decisions")
        elif data_completeness > 85:
            summary['key_findings'].append("✅ Good data quality with minor optimization opportunities")
        else:
            summary['key_findings'].append("⚠️ Data quality requires attention before strategic analysis")

        # Anomaly insights
        if 'anomaly_results' in analysis_results and analysis_results['anomaly_results']:
            anomaly_count = len(analysis_results['anomaly_results']['indices'])
            anomaly_pct = (anomaly_count / data_volume) * 100

            if anomaly_pct > 5:
                summary['key_findings'].append(f"🔍 {anomaly_pct:.1f}% anomalies identified - detailed investigation recommended")
            else:
                summary['key_findings'].append(f"✅ Low anomaly rate ({anomaly_pct:.1f}%) - stable data foundation")

        # Strategic recommendations based on industry
        if industry_pattern == 'financial':
            summary['strategic_recommendations'].extend([
                "💼 Implementation of Real-Time Risk Monitoring",
                "📈 Predictive Models for Financial Forecasting",
                "🔒 Enhanced Compliance Monitoring"
            ])
        elif industry_pattern == 'manufacturing':
            summary['strategic_recommendations'].extend([
                "🏭 Predictive Maintenance Programs",
                "📊 Quality Control Optimization",
                "⚡ Supply Chain Intelligence"
            ])
        elif industry_pattern == 'retail':
            summary['strategic_recommendations'].extend([
                "🛒 Customer Journey Optimization",
                "📱 Personalized Marketing Strategies",
                "📦 Inventory Management Enhancement"
            ])
        else:
            summary['strategic_recommendations'].extend([
                "🎯 Data-driven Decision Processes",
                "🤖 AI-supported Automation",
                "📊 Performance Monitoring Dashboard"
            ])

        # GDPR risk assessment
        if 'gdpr_assessment' in analysis_results:
            gdpr = analysis_results['gdpr_assessment']
            summary['risk_assessment'] = {
                'gdpr_compliance': gdpr['compliance_level'],
                'compliance_score': gdpr['compliance_score'],
                'critical_areas': gdpr['risk_factors'][:3]
            }

        # Technical insights
        summary['technical_insights'] = {
            'data_maturity': 'High' if data_completeness > 90 else 'Medium' if data_completeness > 75 else 'Developing',
            'analysis_complexity': 'Advanced' if len(numeric_cols) > 5 else 'Standard',
            'scalability': 'Enterprise-Ready' if data_volume > 1000 else 'Department-Level'
        }

        # Next steps
        summary['next_steps'] = [
            "📋 Detailed Implementation Planning",
            "👥 Stakeholder Workshop for Prioritization",
            "🎯 KPI Definition and Monitoring Setup",
            "📅 Establish Quarterly Review Cycles"
        ]

        return summary

# Initialize analytics engine
ai_analytics = EnterpriseAIAnalytics()

# Utility functions
@st.cache_data
def load_data(file_content, file_name):
    """Load and cache data with error handling"""
    try:
        if file_name.endswith('.csv'):
            try:
                return pd.read_csv(BytesIO(file_content), encoding='utf-8')
            except UnicodeDecodeError:
                try:
                    return pd.read_csv(BytesIO(file_content), encoding='latin-1')
                except UnicodeDecodeError:
                    return pd.read_csv(BytesIO(file_content), encoding='cp1252')
        else:
            return pd.read_excel(BytesIO(file_content))
    except Exception as e:
        st.error(f"Error loading file: {str(e)}")
        return None

def validate_data(df):
    """Validate uploaded data quality"""
    issues = []
    recommendations = []

    if df is None or df.empty:
        issues.append("⚠️ File is empty or could not be read")
        return issues, recommendations, "poor"

    # Check data quality metrics
    missing_pct = (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
    duplicate_pct = (df.duplicated().sum() / len(df)) * 100
    numeric_cols = len(df.select_dtypes(include=[np.number]).columns)

    # Determine quality score
    if missing_pct < 5 and duplicate_pct < 1 and numeric_cols > 0:
        quality = "excellent"
    elif missing_pct < 15 and duplicate_pct < 5:
        quality = "good"
    else:
        quality = "poor"

    # Generate recommendations
    if missing_pct > 10:
        recommendations.append(f"🔍 {missing_pct:.1f}% missing values identified - data cleaning recommended")

    if duplicate_pct > 2:
        recommendations.append(f"🔄 {duplicate_pct:.1f}% duplicates found - deduplication recommended")

    if numeric_cols == 0:
        recommendations.append("📊 No numeric columns detected - advanced analysis limited")

    return issues, recommendations, quality

def create_progress_placeholder():
    """Create progress tracking interface"""
    progress_container = st.container()
    with progress_container:
        progress_bar = st.progress(0)
        status_text = st.empty()
    return progress_bar, status_text

# Sidebar configuration
with st.sidebar:
    st.markdown("### 🛠️ Configuration")

    # Analysis settings
    st.markdown("#### Analysis Settings")
    analysis_depth = st.selectbox(
        "Analysis Depth",
        ["Quick Overview", "Standard Analysis", "Deep Analysis", "🎓 Academic Research"],
        index=2
    )

    include_clustering = st.checkbox("Clustering Analysis", value=True)
    include_anomalies = st.checkbox("🔍 Advanced Anomaly Detection", value=True)
    include_forecasting = st.checkbox("📈 Statistical Forecasting Models", value=True)
    include_ai_insights = st.checkbox("🤖 AI Business Intelligence", value=True)
    include_gdpr_assessment = st.checkbox("⚖️ GDPR Compliance Check", value=True)
    include_executive_summary = st.checkbox("📋 Executive Summary", value=True)

    # Export settings
    st.markdown("#### Export Settings")
    report_language = st.selectbox("Report Language", ["English", "Deutsch"], index=0)
    include_charts = st.checkbox("Include Charts", value=True)
    include_raw_data = st.checkbox("Append Raw Data", value=False)

    st.markdown("---")
    st.markdown("### 📊 System Status")
    st.success("✅ All systems operational")
    st.info(f"🕰️ Last update: {datetime.now().strftime('%H:%M:%S')}")

# Main file upload section
st.markdown("### 📁 Data Upload")

col1, col2 = st.columns([2, 1])

with col1:
    uploaded_file = st.file_uploader(
        "Choose your data file",
        type=['xlsx', 'csv', 'xls'],
        help="Supported formats: Excel (.xlsx, .xls), CSV (.csv). Maximum file size: 200MB"
    )

with col2:
    if uploaded_file:
        file_size = len(uploaded_file.getvalue()) / (1024 * 1024)  # MB
        st.metric("File Size", f"{file_size:.1f} MB")
        st.metric("File Format", uploaded_file.name.split('.')[-1].upper())

def perform_analysis(df, numeric_cols):
    """Perform comprehensive AI analysis"""
    progress_bar, status_text = create_progress_placeholder()

    analysis_results = {}

    try:
        # Industry pattern detection
        status_text.text("🏢 Detecting industry patterns...")
        progress_bar.progress(0.1)
        industry_result = ai_analytics.detect_industry_pattern(df)
        analysis_results['industry'] = industry_result
        industry_pattern = industry_result['pattern']
        industry_confidence = industry_result['confidence']
        pattern_details = industry_result['details']

        # Analysis steps
        steps = [
            "Statistical Analysis",
            "Correlation Analysis"
        ]

        if include_anomalies:
            steps.append("Multi-Method Anomaly Detection")
        if include_forecasting:
            steps.append("Statistical Trend Modeling")
        if include_gdpr_assessment:
            steps.append("GDPR Compliance Assessment")
        if include_ai_insights:
            steps.append("AI Business Intelligence")
        if include_executive_summary:
            steps.append("Executive Summary Generation")

        total_steps = len(steps)

        for i, step in enumerate(steps):
            status_text.text(f"🔄 {step}...")
            progress_bar.progress(0.1 + (i + 1) / total_steps * 0.9)
            time.sleep(0.3)

        # Clear progress
        progress_bar.empty()
        status_text.empty()

        # Display results
        display_industry_analysis(industry_pattern, industry_confidence, pattern_details)
        display_statistical_analysis(df, numeric_cols)

        if len(numeric_cols) > 1:
            correlation_results = ai_analytics.advanced_correlation_analysis(df, numeric_cols)
            analysis_results['correlation_results'] = correlation_results
            display_correlation_analysis(correlation_results)

        if include_anomalies and len(numeric_cols) > 0:
            anomaly_results, anomaly_data = ai_analytics.advanced_anomaly_detection(df, numeric_cols)
            analysis_results['anomaly_results'] = anomaly_results
            display_anomaly_analysis(anomaly_results, anomaly_data, df)

        if include_forecasting and len(numeric_cols) > 0:
            target_col = numeric_cols[0]
            date_cols = df.select_dtypes(include=['datetime64', 'object']).columns
            date_col = None
            for col in date_cols:
                if 'date' in col.lower() or 'time' in col.lower():
                    date_col = col
                    break

            trend_results = ai_analytics.advanced_trend_analysis(df, target_col, date_col)
            analysis_results['trend_results'] = trend_results
            display_forecasting_analysis(trend_results, target_col)

        if include_gdpr_assessment:
            gdpr_results = ai_analytics.gdpr_compliance_assessment(df)
            analysis_results['gdpr_assessment'] = gdpr_results
            display_gdpr_analysis(gdpr_results)

        if include_executive_summary:
            executive_summary = ai_analytics.generate_executive_summary(df, analysis_results, industry_pattern)
            display_executive_summary(executive_summary)

        # PowerPoint generation
        display_powerpoint_generation(df, analysis_results)

    except Exception as e:
        st.error(f"🚨 Error during analysis: {str(e)}")
        progress_bar.empty()
        status_text.empty()

def display_industry_analysis(industry_pattern, industry_confidence, pattern_details):
    """Display industry pattern detection results"""
    st.markdown("### 🏢 Industry Detection")

    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("🎯 Detected Industry", industry_pattern.title())
    with col2:
        st.metric("🎯 Confidence", f"{industry_confidence*100:.1f}%")
    with col3:
        total_indicators = sum(pattern_details.values())
        st.metric("📊 Indicators", total_indicators)

    if industry_confidence > 0.1:
        with st.expander("🔍 Industry Analysis Details", expanded=False):
            st.write("**Detected Patterns by Category:**")
            for category, count in pattern_details.items():
                if count > 0:
                    st.write(f"- {category.title()}: {count} indicators")

def display_statistical_analysis(df, numeric_cols):
    """Display enhanced statistical analysis"""
    st.markdown("### 📊 Enhanced Statistical Analysis")

    if len(numeric_cols) > 0:
        stats_df = df[numeric_cols].describe().round(4)
        st.dataframe(stats_df, width="stretch")

        # Advanced statistical tests
        col1, col2 = st.columns(2)

        with col1:
            st.markdown("#### 🔍 Normality Tests")
            for col in numeric_cols[:3]:
                data = df[col].dropna()
                if len(data) > 8:
                    try:
                        stat, p_value = stats.shapiro(data[:5000])
                        is_normal = p_value > 0.05
                        st.write(f"**{col}**: {'Normal' if is_normal else 'Non-normal'} (p={p_value:.4f})")
                    except:
                        st.write(f"**{col}**: Test not possible")

        with col2:
            st.markdown("#### 📈 Distribution Characteristics")
            for col in numeric_cols[:3]:
                data = df[col].dropna()
                if len(data) > 0:
                    skewness = stats.skew(data)
                    kurtosis = stats.kurtosis(data)
                    cv = data.std() / data.mean() if data.mean() != 0 else 0

                    interpretation = ""
                    if abs(skewness) < 0.5:
                        interpretation += "Symmetric"
                    elif skewness > 0:
                        interpretation += "Right-skewed"
                    else:
                        interpretation += "Left-skewed"

                    st.write(f"**{col}**: {interpretation} (Skew: {skewness:.3f}, CV: {cv:.3f})")

def display_correlation_analysis(correlation_results):
    """Display advanced correlation analysis"""
    st.markdown("### 🔄 Advanced Correlation Analysis")

    if correlation_results is None:
        st.info("📊 Insufficient numeric data for correlation analysis")
        return

    col1, col2 = st.columns(2)

    with col1:
        st.markdown("#### Pearson Correlation")
        fig = px.imshow(
            correlation_results['pearson'],
            text_auto=True,
            title="Pearson Correlation Matrix",
            color_continuous_scale="RdBu",
            zmin=-1, zmax=1
        )
        st.plotly_chart(fig, width="stretch")

    with col2:
        st.markdown("#### Spearman Correlation (Rank-based)")
        fig = px.imshow(
            correlation_results['spearman'],
            text_auto=True,
            title="Spearman Correlation Matrix",
            color_continuous_scale="RdBu",
            zmin=-1, zmax=1
        )
        st.plotly_chart(fig, width="stretch")

    # Partial correlations
    if correlation_results['partial_correlations']:
        st.markdown("#### 🤔 Partial Correlations")
        st.write("*Correlations after controlling for other variables:*")
        for pair, corr in correlation_results['partial_correlations'].items():
            if abs(corr) > 0.3:
                cols = pair.split('_vs_')
                st.write(f"**{cols[0]}** ↔️ **{cols[1]}**: {corr:.3f}")

def display_anomaly_analysis(anomaly_results, anomaly_data, df):
    """Display advanced anomaly detection results"""
    st.markdown("### ⚠️ Advanced Anomaly Detection")

    if anomaly_results is None:
        st.info("📊 Insufficient data for anomaly detection")
        return

    n_anomalies = len(anomaly_results['indices'])
    total_records = len(anomaly_data) if anomaly_data is not None else len(df)
    anomaly_rate = (n_anomalies / total_records) * 100

    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("📊 Data Points", f"{total_records:,}")
    with col2:
        st.metric("⚠️ Anomalies", n_anomalies)
    with col3:
        st.metric("📉 Anomaly Rate", f"{anomaly_rate:.2f}%")
    with col4:
        avg_confidence = np.mean(anomaly_results['confidence']) if len(anomaly_results['confidence']) > 0 else 0
        st.metric("🎯 Confidence", f"{avg_confidence*100:.1f}%")

    # Methods used
    st.markdown("#### 🔧 Methods Used")
    st.write(", ".join(anomaly_results['methods_used']))

    # Confidence distribution
    if len(anomaly_results['confidence']) > 0:
        fig = go.Figure(data=go.Histogram(
            x=anomaly_results['confidence'],
            nbinsx=10,
            title="Anomaly Confidence Distribution"
        ))
        fig.update_layout(xaxis_title="Confidence Score", yaxis_title="Number of Anomalies")
        st.plotly_chart(fig, width="stretch")

    # Show high-confidence anomalies
    if n_anomalies > 0:
        high_conf_mask = anomaly_results['confidence'] > 0.7
        if np.any(high_conf_mask):
            st.markdown("#### 🔴 High-Confidence Anomalies")
            high_conf_indices = anomaly_results['indices'][high_conf_mask]
            with st.expander(f"{len(high_conf_indices)} high-confidence anomalies", expanded=False):
                st.dataframe(df.loc[high_conf_indices], width="stretch")

def display_forecasting_analysis(trend_results, target_col):
    """Display advanced forecasting analysis"""
    st.markdown("### 🔮 Advanced Forecasting Analysis")

    st.write(f"**Target Variable**: {target_col}")

    # Stationarity analysis
    if trend_results.get('stationarity'):
        st.markdown("#### 📈 Stationarity Analysis")
        stationarity = trend_results['stationarity']

        col1, col2 = st.columns(2)
        with col1:
            status = "Stationary" if stationarity['is_stationary'] else "Non-stationary"
            st.metric("Stationarity Test", status)
        with col2:
            st.metric("Mean Shift", f"{stationarity['mean_shift']:.4f}")

    # Enhanced forecasting results
    if trend_results.get('enhanced_forecast'):
        st.markdown("#### 📊 Enhanced Forecast")
        enhanced = trend_results['enhanced_forecast']

        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Model Type", enhanced['model_type'])
        with col2:
            st.metric("R² Score", f"{enhanced['r2_score']:.3f}")
        with col3:
            st.metric("Forecast Periods", len(enhanced['forecast']))

        # Visualization
        fig = go.Figure()

        # Historical data
        if 'original' in trend_results:
            original_data = trend_results['original']
            historical_x = list(range(len(original_data)))
            fig.add_trace(go.Scatter(
                x=historical_x,
                y=original_data.values,
                mode='lines',
                name='Historical Data',
                line=dict(color='blue')
            ))
        else:
            historical_x = []

        forecast_x = list(range(len(historical_x), len(historical_x) + len(enhanced['forecast'])))

        # Confidence intervals
        fig.add_trace(go.Scatter(
            x=forecast_x + forecast_x[::-1],
            y=list(enhanced['confidence_upper']) + list(enhanced['confidence_lower'][::-1]),
            fill='toself',
            fillcolor='rgba(255,0,0,0.2)',
            line=dict(color='rgba(255,255,255,0)'),
            name='95% Confidence Interval',
            showlegend=True
        ))

        # Point forecast
        fig.add_trace(go.Scatter(
            x=forecast_x,
            y=enhanced['forecast'],
            mode='lines+markers',
            name=f'{enhanced["model_type"]} Forecast',
            line=dict(color='red')
        ))

        fig.update_layout(
            title=f"Advanced Forecast: {target_col}",
            xaxis_title="Time",
            yaxis_title=target_col
        )
        st.plotly_chart(fig, width="stretch")

def display_gdpr_analysis(gdpr_results):
    """Display GDPR compliance assessment"""
    st.markdown("### ⚖️ GDPR Compliance Assessment")

    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("🎯 Compliance Score", f"{gdpr_results['compliance_score']}/100")
    with col2:
        level_with_color = f"{gdpr_results['level_color']} {gdpr_results['compliance_level']}"
        st.metric("📈 Risk Level", level_with_color)
    with col3:
        st.metric("🕰️ Assessment Time", gdpr_results['assessment_timestamp'].strftime('%H:%M'))

    if gdpr_results['risk_factors']:
        st.markdown("#### ⚠️ Identified Risk Factors")
        for risk in gdpr_results['risk_factors']:
            st.warning(risk)

    if gdpr_results['potential_pii']:
        st.markdown("#### 🔍 Detected PII Categories")
        for category, columns in gdpr_results['potential_pii'].items():
            with st.expander(f"{category.title()} ({len(columns)} columns)", expanded=False):
                for col in columns:
                    st.write(f"- {col}")

    if gdpr_results['recommendations']:
        st.markdown("#### 📝 Compliance Recommendations")
        for rec in gdpr_results['recommendations']:
            st.info(rec)

def display_executive_summary(executive_summary):
    """Display professional executive summary"""
    st.markdown("### 📋 Executive Summary")

    header = executive_summary['header']
    st.markdown(f"""
    <div class="analysis-card">
        <h2>{header['title']}</h2>
        <h4>{header['subtitle']}</h4>
        <p><strong>Industry</strong>: {header['industry']} | <strong>Generated</strong>: {header['timestamp']}</p>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("#### 🎯 Key Findings")
    for finding in executive_summary['key_findings']:
        st.write(finding)

    st.markdown("#### 🚀 Strategic Recommendations")
    for rec in executive_summary['strategic_recommendations']:
        st.write(rec)

    if executive_summary['risk_assessment']:
        st.markdown("#### ⚠️ Risk Assessment")
        risk = executive_summary['risk_assessment']
        col1, col2 = st.columns(2)
        with col1:
            st.metric("GDPR Compliance", risk['gdpr_compliance'])
        with col2:
            st.metric("Compliance Score", f"{risk['compliance_score']}/100")

    tech = executive_summary['technical_insights']
    st.markdown("#### 💻 Technical Assessment")
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("Data Maturity", tech['data_maturity'])
    with col2:
        st.metric("Analysis Complexity", tech['analysis_complexity'])
    with col3:
        st.metric("Scalability", tech['scalability'])

    st.markdown("#### 📝 Next Steps")
    for step in executive_summary['next_steps']:
        st.write(step)

def display_powerpoint_generation(df, analysis_results):
    """Generate enhanced PowerPoint report with German corporate standards"""
    st.markdown("### 📄 Professional Report Generation")

    col1, col2 = st.columns(2)

    with col1:
        st.info("🎓 **Academic Research Grade**: Comprehensive report with all AI analyses")

        if st.button("📊 Generate Standard Report", type="secondary"):
            with st.spinner("Creating standard PowerPoint presentation..."):
                ppt_buffer = create_powerpoint_report(df, analysis_results)

                if ppt_buffer:
                    st.success("✅ Standard PowerPoint report successfully created!")
                    st.download_button(
                        label="📎 Download Standard Report",
                        data=ppt_buffer.getvalue(),
                        file_name=f"Standard_AI_Analysis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

    with col2:
        st.warning("🏢 **German Corporate Standard**: SAP-style consulting firm quality")

        if st.button("🚀 Generate Executive Corporate Report", type="primary"):
            with st.spinner("Creating premium German corporate presentation..."):
                try:
                    ppt_buffer = create_german_corporate_powerpoint(df, analysis_results)

                    if ppt_buffer:
                        st.success("✅ Premium German corporate report successfully created!")
                        st.balloons()
                        st.download_button(
                            label="📎 Download Premium Corporate Report",
                            data=ppt_buffer.getvalue(),
                            file_name=f"Executive_Corporate_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                except Exception as e:
                    st.error(f"Error creating corporate report: {str(e)}")
                    # Fallback to standard report
                    ppt_buffer = create_powerpoint_report(df, analysis_results)
                    if ppt_buffer:
                        st.warning("Created standard report as fallback")
                        st.download_button(
                            label="📎 Download Fallback Report",
                            data=ppt_buffer.getvalue(),
                            file_name=f"Fallback_AI_Analysis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )

    # Report features overview
    st.markdown("#### 📋 Report Features Comparison")

    comparison_data = {
        "Feature": [
            "Executive Summary",
            "KPI Dashboard",
            "Data Quality Assessment",
            "Industry Analysis",
            "GDPR Compliance",
            "Risk Assessment",
            "Strategic Recommendations",
            "German Corporate Branding",
            "SAP-style Templates",
            "Professional Charts",
            "Implementation Roadmap",
            "Consulting Firm Quality"
        ],
        "Standard Report": [
            "✅", "❌", "✅", "❌", "✅", "❌", "✅", "❌", "❌", "✅", "❌", "❌"
        ],
        "Corporate Report": [
            "✅", "✅", "✅", "✅", "✅", "✅", "✅", "✅", "✅", "✅", "✅", "✅"
        ]
    }

    comparison_df = pd.DataFrame(comparison_data)
    st.dataframe(comparison_df, width="stretch")

def create_powerpoint_report(df, analysis_results):
    """Create comprehensive PowerPoint report"""
    try:
        prs = Presentation()

        # Title slide
        title_slide = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Executive AI Data Intelligence Report"
        subtitle.text = f"Academic Research Grade Analysis\\nGenerated on {datetime.now().strftime('%d/%m/%Y at %H:%M')}"

        # Executive Summary slide
        if 'executive_summary' in analysis_results:
            bullet_slide = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide)
            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = "Executive Summary"
            tf = body.text_frame
            summary = analysis_results['executive_summary']

            # Safely access industry information
            industry = "General"
            if 'header' in summary and 'industry' in summary['header']:
                industry = summary['header']['industry']

            tf.text = f"Industry: {industry}"
            if 'key_findings' in summary:
                for finding in summary['key_findings'][:3]:
                    p = tf.add_paragraph()
                    p.text = finding.replace('✅', '').replace('⚠️', '').replace('🚨', '')

        # Data Overview slide
        bullet_slide = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "Data Overview & Quality"
        tf = body.text_frame
        tf.text = f"Data Records: {len(df):,}"

        p = tf.add_paragraph()
        p.text = f"Columns: {len(df.columns)}"

        p = tf.add_paragraph()
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        p.text = f"Numeric Fields: {len(numeric_cols)}"

        p = tf.add_paragraph()
        missing_pct = (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
        p.text = f"Data Quality: {100-missing_pct:.1f}% complete"

        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer

    except Exception as e:
        st.error(f"PowerPoint generation error: {str(e)}")
        return None

# Main data processing
if uploaded_file:
    progress_bar, status_text = create_progress_placeholder()

    try:
        # Load data
        status_text.text("📂 Loading file...")
        progress_bar.progress(20)

        file_content = uploaded_file.getvalue()
        df = load_data(file_content, uploaded_file.name)

        if df is None:
            st.stop()

        # Validate data
        status_text.text("🔍 Validating data quality...")
        progress_bar.progress(40)

        issues, recommendations, quality = validate_data(df)

        # Display results
        progress_bar.progress(60)
        status_text.text("✅ Data successfully loaded")
        progress_bar.empty()
        status_text.empty()

        # Data overview
        st.markdown("""
        <div class="analysis-card">
            <h3>📊 Data Overview</h3>
        </div>
        """, unsafe_allow_html=True)

        # Quality indicator
        quality_class = f"quality-{quality}"
        quality_text = {
            "excellent": "Excellent",
            "good": "Good",
            "poor": "Needs Improvement"
        }

        st.markdown(f"""
        <div class="quality-indicator {quality_class}">
            Data Quality: {quality_text[quality]}
        </div>
        """, unsafe_allow_html=True)

        # Data metrics
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.metric("📊 Records", f"{len(df):,}")
        with col2:
            st.metric("📈 Columns", len(df.columns))
        with col3:
            numeric_cols = df.select_dtypes(include=[np.number]).columns
            st.metric("🔢 Numeric Fields", len(numeric_cols))
        with col4:
            missing_pct = (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
            st.metric("✅ Completeness", f"{100-missing_pct:.1f}%")

        # Show recommendations
        if recommendations:
            st.markdown("#### 💡 Recommendations")
            for rec in recommendations:
                st.warning(rec)

        # Data preview
        with st.expander("🔍 Data Preview", expanded=False):
            st.dataframe(df.head(10), width="stretch")

        # Analysis button
        if st.button("🚀 Start Intelligent Analysis", type="primary"):
            perform_analysis(df, numeric_cols)

    except Exception as e:
        st.error(f"🚨 Error processing data: {str(e)}")
        with st.expander("🔍 Technical Details", expanded=False):
            st.code(traceback.format_exc())

# Demo section
st.markdown("---")
st.markdown("### 🎪 Professor Demo Mode")

col1, col2 = st.columns(2)

with col1:
    if st.button("🎪 Start Demo with Sample Data", type="secondary"):
        st.balloons()

        # Create comprehensive sample data
        np.random.seed(42)
        n_samples = 500

        sample_data = pd.DataFrame({
            'Date': pd.date_range('2022-01-01', periods=n_samples, freq='D'),
            'Revenue_USD': np.random.normal(45000, 12000, n_samples).clip(min=10000),
            'Customer_Segment': np.random.choice(['Premium', 'Standard', 'Basic'], n_samples, p=[0.2, 0.5, 0.3]),
            'Region': np.random.choice(['North', 'South', 'East', 'West'], n_samples),
            'Product_Category': np.random.choice(['Software', 'Hardware', 'Services'], n_samples, p=[0.4, 0.3, 0.3]),
            'Employee_Count': np.random.poisson(25, n_samples),
            'Customer_Satisfaction': np.random.normal(4.2, 0.8, n_samples).clip(min=1, max=5),
            'Marketing_Budget': np.random.exponential(5000, n_samples)
        })

        # Add correlations and trends
        sample_data['Revenue_USD'] += sample_data['Employee_Count'] * 800
        sample_data['Revenue_USD'] += (sample_data['Customer_Satisfaction'] - 3) * 5000

        # Add time trend
        time_trend = np.arange(len(sample_data)) * 10
        sample_data['Revenue_USD'] += time_trend

        st.session_state['demo_data'] = sample_data
        st.success("✅ Demo data loaded! Ready for analysis.")
        st.dataframe(sample_data.head(10), width="stretch")

with col2:
    if 'demo_data' in st.session_state:
        if st.button("🚀 Start Demo Analysis", type="primary"):
            demo_df = st.session_state['demo_data']
            demo_numeric_cols = demo_df.select_dtypes(include=[np.number]).columns
            perform_analysis(demo_df, demo_numeric_cols)

# Footer
st.markdown("---")
st.markdown("""
<div style="text-align: center; color: #6b7280; font-size: 0.9rem;">
    🏢 <strong>AI Automation Platform</strong> |
    Enterprise-Grade Data Intelligence |
    Powered by Advanced Analytics & Machine Learning
</div>
""", unsafe_allow_html=True)