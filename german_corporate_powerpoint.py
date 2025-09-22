"""
German Corporate PowerPoint Generator
Professional SAP-style presentations for enterprise analytics
"""

import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR
from datetime import datetime
from io import BytesIO

class GermanCorporatePowerPoint:
    def __init__(self):
        # German corporate color palette (SAP-inspired)
        self.colors = {
            'primary_blue': RGBColor(0, 90, 150),      # SAP Blue
            'secondary_blue': RGBColor(0, 125, 195),   # Light SAP Blue
            'accent_blue': RGBColor(0, 165, 235),      # Bright Blue
            'dark_gray': RGBColor(64, 64, 64),         # Corporate Dark Gray
            'medium_gray': RGBColor(128, 128, 128),    # Medium Gray
            'light_gray': RGBColor(240, 240, 240),     # Light Gray Background
            'white': RGBColor(255, 255, 255),          # Pure White
            'green': RGBColor(46, 125, 50),            # Success Green
            'orange': RGBColor(255, 152, 0),           # Warning Orange
            'red': RGBColor(211, 47, 47),              # Error Red
            'gold': RGBColor(255, 193, 7),             # Premium Gold
        }

        # Professional typography
        self.fonts = {
            'title': 'Calibri',
            'subtitle': 'Calibri Light',
            'body': 'Calibri',
            'accent': 'Segoe UI'
        }

    def create_presentation(self, df, analysis_results):
        """Create comprehensive German corporate presentation"""
        prs = Presentation()

        # Create slides in order
        self._create_title_slide(prs, df, analysis_results)
        self._create_executive_summary_slide(prs, df, analysis_results)
        self._create_kpi_dashboard_slide(prs, df, analysis_results)
        self._create_data_quality_slide(prs, df, analysis_results)
        self._create_industry_analysis_slide(prs, df, analysis_results)
        self._create_gdpr_compliance_slide(prs, df, analysis_results)
        self._create_recommendations_slide(prs, df, analysis_results)

        return prs

    def _create_title_slide(self, prs, df, analysis_results):
        """Create professional title slide with German corporate styling"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Background gradient (simulated with shapes)
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            prs.slide_width, prs.slide_height
        )
        bg_fill = bg_shape.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = self.colors['light_gray']
        bg_shape.line.fill.background()

        # Corporate header bar
        header_height = Cm(3)
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, header_height
        )
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = self.colors['primary_blue']
        header_shape.line.fill.background()

        # Title positioning
        title_left = Cm(2)
        title_top = Cm(4)
        title_width = Cm(20)

        # Main title
        title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, Cm(2))
        title_frame = title_shape.text_frame
        title_frame.margin_left = 0
        title_frame.margin_right = 0
        title_frame.word_wrap = True

        title_p = title_frame.paragraphs[0]
        title_p.text = "DATENANALYSE & KI-INTELLIGENCE"
        title_p.font.name = self.fonts['title']
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = self.colors['primary_blue']
        title_p.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_top = Cm(6.5)
        subtitle_shape = slide.shapes.add_textbox(title_left, subtitle_top, title_width, Cm(1.5))
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.margin_left = 0
        subtitle_frame.margin_right = 0

        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = "Enterprise Analytics Report · Datenschutz-konforme Auswertung"
        subtitle_p.font.name = self.fonts['subtitle']
        subtitle_p.font.size = Pt(18)
        subtitle_p.font.color.rgb = self.colors['dark_gray']
        subtitle_p.alignment = PP_ALIGN.CENTER

        # Data overview box
        overview_top = Cm(9)
        overview_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Cm(3), overview_top, Cm(18), Cm(4)
        )
        overview_fill = overview_shape.fill
        overview_fill.solid()
        overview_fill.fore_color.rgb = self.colors['light_gray']
        overview_shape.line.color.rgb = self.colors['medium_gray']
        overview_shape.line.width = Pt(1)

        overview_text = overview_shape.text_frame
        overview_text.margin_left = Cm(1)
        overview_text.margin_right = Cm(1)
        overview_text.margin_top = Cm(0.5)

        # Data metrics
        data_volume = len(df)
        data_completeness = 100 - (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
        numeric_cols = len(df.select_dtypes(include=[np.number]).columns)

        overview_p = overview_text.paragraphs[0]
        overview_p.text = "DATENÜBERSICHT"
        overview_p.font.name = self.fonts['body']
        overview_p.font.size = Pt(14)
        overview_p.font.bold = True
        overview_p.font.color.rgb = self.colors['primary_blue']
        overview_p.alignment = PP_ALIGN.CENTER

        # Add metrics
        metrics_text = f"\n\n📊 Datensätze: {data_volume:,}  ·  📈 Spalten: {len(df.columns)}  ·  🔢 Numerische Felder: {numeric_cols}\n"
        metrics_text += f"✅ Datenqualität: {data_completeness:.1f}%  ·  🏢 Branche: {analysis_results.get('industry', {}).get('pattern', 'Allgemein').title()}"

        metrics_p = overview_text.add_paragraph()
        metrics_p.text = metrics_text
        metrics_p.font.name = self.fonts['body']
        metrics_p.font.size = Pt(12)
        metrics_p.font.color.rgb = self.colors['dark_gray']
        metrics_p.alignment = PP_ALIGN.CENTER

        # Footer
        footer_top = Cm(15)
        footer_shape = slide.shapes.add_textbox(title_left, footer_top, title_width, Cm(1))
        footer_frame = footer_shape.text_frame

        footer_p = footer_frame.paragraphs[0]
        footer_p.text = f"Erstellt am {datetime.now().strftime('%d.%m.%Y')} · KI-gestützte Analyse · Vertraulich"
        footer_p.font.name = self.fonts['body']
        footer_p.font.size = Pt(10)
        footer_p.font.color.rgb = self.colors['medium_gray']
        footer_p.alignment = PP_ALIGN.CENTER

    def _create_executive_summary_slide(self, prs, df, analysis_results):
        """Create executive summary with key insights"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Add slide header
        self._add_slide_header(slide, "EXECUTIVE SUMMARY", "Wichtigste Erkenntnisse und strategische Empfehlungen")

        # Key findings section
        findings_top = Cm(4)
        findings_shape = slide.shapes.add_textbox(Cm(1.5), findings_top, Cm(10), Cm(8))
        findings_frame = findings_shape.text_frame
        findings_frame.margin_left = Cm(0.5)

        findings_title = findings_frame.paragraphs[0]
        findings_title.text = "🎯 KERNERKENNTNISSE"
        findings_title.font.name = self.fonts['body']
        findings_title.font.size = Pt(14)
        findings_title.font.bold = True
        findings_title.font.color.rgb = self.colors['primary_blue']

        # Add key findings
        if 'executive_summary' in analysis_results:
            summary = analysis_results['executive_summary']
            for finding in summary.get('key_findings', [])[:4]:
                finding_p = findings_frame.add_paragraph()
                finding_p.text = f"• {finding.replace('✅', '').replace('⚠️', '').replace('🚨', '').strip()}"
                finding_p.font.name = self.fonts['body']
                finding_p.font.size = Pt(11)
                finding_p.font.color.rgb = self.colors['dark_gray']
                finding_p.space_before = Pt(6)

        # Strategic recommendations section
        recommendations_shape = slide.shapes.add_textbox(Cm(12.5), findings_top, Cm(10), Cm(8))
        recommendations_frame = recommendations_shape.text_frame
        recommendations_frame.margin_left = Cm(0.5)

        recommendations_title = recommendations_frame.paragraphs[0]
        recommendations_title.text = "🚀 STRATEGISCHE EMPFEHLUNGEN"
        recommendations_title.font.name = self.fonts['body']
        recommendations_title.font.size = Pt(14)
        recommendations_title.font.bold = True
        recommendations_title.font.color.rgb = self.colors['primary_blue']

        # Add recommendations
        if 'executive_summary' in analysis_results:
            summary = analysis_results['executive_summary']
            for rec in summary.get('strategic_recommendations', [])[:4]:
                rec_p = recommendations_frame.add_paragraph()
                rec_p.text = f"• {rec.replace('💼', '').replace('📈', '').replace('🔒', '').strip()}"
                rec_p.font.name = self.fonts['body']
                rec_p.font.size = Pt(11)
                rec_p.font.color.rgb = self.colors['dark_gray']
                rec_p.space_before = Pt(6)

        # Risk assessment box
        risk_top = Cm(12.5)
        risk_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.5), risk_top, Cm(21), Cm(2.5)
        )
        risk_fill = risk_shape.fill
        risk_fill.solid()
        risk_fill.fore_color.rgb = self.colors['light_gray']
        risk_shape.line.color.rgb = self.colors['medium_gray']

        risk_text = risk_shape.text_frame
        risk_text.margin_left = Cm(1)
        risk_text.margin_top = Cm(0.3)

        risk_title = risk_text.paragraphs[0]
        risk_title.text = "⚖️ COMPLIANCE & RISIKOBEWERTUNG"
        risk_title.font.name = self.fonts['body']
        risk_title.font.size = Pt(12)
        risk_title.font.bold = True
        risk_title.font.color.rgb = self.colors['primary_blue']

        if 'gdpr_assessment' in analysis_results:
            gdpr = analysis_results['gdpr_assessment']
            risk_details = risk_text.add_paragraph()
            risk_details.text = f"GDPR Compliance Score: {gdpr['compliance_score']}/100 · Status: {gdpr['compliance_level']} · Letzte Prüfung: {datetime.now().strftime('%d.%m.%Y')}"
            risk_details.font.name = self.fonts['body']
            risk_details.font.size = Pt(10)
            risk_details.font.color.rgb = self.colors['dark_gray']

    def _create_kpi_dashboard_slide(self, prs, df, analysis_results):
        """Create executive KPI dashboard"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_header(slide, "KPI DASHBOARD", "Zentrale Leistungskennzahlen im Überblick")

        # Calculate KPIs
        data_volume = len(df)
        data_completeness = 100 - (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
        numeric_cols = df.select_dtypes(include=[np.number]).columns

        # Anomaly rate
        anomaly_rate = 0
        if 'anomaly_results' in analysis_results and analysis_results['anomaly_results']:
            anomaly_count = len(analysis_results['anomaly_results']['indices'])
            anomaly_rate = (anomaly_count / data_volume) * 100

        # GDPR score
        gdpr_score = 100
        if 'gdpr_assessment' in analysis_results:
            gdpr_score = analysis_results['gdpr_assessment']['compliance_score']

        # Create KPI boxes
        kpi_data = [
            {"title": "DATENQUALITÄT", "value": f"{data_completeness:.1f}%", "status": "good" if data_completeness > 90 else "warning"},
            {"title": "DATENSÄTZE", "value": f"{data_volume:,}", "status": "good"},
            {"title": "ANOMALIE-RATE", "value": f"{anomaly_rate:.1f}%", "status": "good" if anomaly_rate < 5 else "warning"},
            {"title": "GDPR COMPLIANCE", "value": f"{gdpr_score}/100", "status": "good" if gdpr_score > 80 else "warning"}
        ]

        # Position KPI boxes
        box_width = Cm(5)
        box_height = Cm(3)
        start_left = Cm(1.5)
        box_spacing = Cm(5.5)
        box_top = Cm(5)

        for i, kpi in enumerate(kpi_data):
            box_left = start_left + (i * box_spacing)

            # Create KPI box
            kpi_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, box_left, box_top, box_width, box_height
            )

            # Set colors based on status
            fill = kpi_shape.fill
            fill.solid()
            if kpi["status"] == "good":
                fill.fore_color.rgb = self.colors['green']
            else:
                fill.fore_color.rgb = self.colors['orange']

            kpi_shape.line.fill.background()

            # Add text
            text_frame = kpi_shape.text_frame
            text_frame.margin_left = Cm(0.3)
            text_frame.margin_right = Cm(0.3)
            text_frame.margin_top = Cm(0.3)
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Title
            title_p = text_frame.paragraphs[0]
            title_p.text = kpi["title"]
            title_p.font.name = self.fonts['body']
            title_p.font.size = Pt(10)
            title_p.font.bold = True
            title_p.font.color.rgb = self.colors['white']
            title_p.alignment = PP_ALIGN.CENTER

            # Value
            value_p = text_frame.add_paragraph()
            value_p.text = kpi["value"]
            value_p.font.name = self.fonts['title']
            value_p.font.size = Pt(18)
            value_p.font.bold = True
            value_p.font.color.rgb = self.colors['white']
            value_p.alignment = PP_ALIGN.CENTER

    def _create_data_quality_slide(self, prs, df, analysis_results):
        """Create detailed data quality assessment slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_header(slide, "DATENQUALITÄT & VALIDIERUNG", "Detaillierte Bewertung der Datenintegrität")

        # Data quality metrics
        missing_by_column = df.isnull().sum()
        total_missing = missing_by_column.sum()
        missing_pct = (total_missing / (len(df) * len(df.columns))) * 100
        duplicate_count = df.duplicated().sum()
        duplicate_pct = (duplicate_count / len(df)) * 100

        # Quality assessment table
        table_top = Cm(4.5)
        table_shape = slide.shapes.add_textbox(Cm(1.5), table_top, Cm(21), Cm(6))
        table_frame = table_shape.text_frame

        table_title = table_frame.paragraphs[0]
        table_title.text = "📊 DATENQUALITÄTS-METRIKEN"
        table_title.font.name = self.fonts['body']
        table_title.font.size = Pt(14)
        table_title.font.bold = True
        table_title.font.color.rgb = self.colors['primary_blue']

        # Add quality metrics
        metrics = [
            f"• Gesamte Datensätze: {len(df):,}",
            f"• Gesamte Spalten: {len(df.columns)}",
            f"• Fehlende Werte: {total_missing:,} ({missing_pct:.2f}%)",
            f"• Duplikate: {duplicate_count:,} ({duplicate_pct:.2f}%)",
            f"• Numerische Spalten: {len(df.select_dtypes(include=[np.number]).columns)}",
            f"• Kategorische Spalten: {len(df.select_dtypes(include=['object']).columns)}"
        ]

        for metric in metrics:
            metric_p = table_frame.add_paragraph()
            metric_p.text = metric
            metric_p.font.name = self.fonts['body']
            metric_p.font.size = Pt(11)
            metric_p.font.color.rgb = self.colors['dark_gray']
            metric_p.space_before = Pt(4)

    def _create_industry_analysis_slide(self, prs, df, analysis_results):
        """Create industry-specific analysis slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_header(slide, "BRANCHEN-ANALYSE", "Branchenspezifische Einordnung und Benchmarking")

        # Industry detection results
        industry_pattern = "Allgemein"
        industry_confidence = 0
        if 'industry' in analysis_results:
            industry_pattern = analysis_results['industry']['pattern'].title()
            industry_confidence = analysis_results['industry']['confidence'] * 100

        # Industry overview box
        industry_top = Cm(4.5)
        industry_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.5), industry_top, Cm(10), Cm(4)
        )
        industry_fill = industry_shape.fill
        industry_fill.solid()
        industry_fill.fore_color.rgb = self.colors['secondary_blue']
        industry_shape.line.fill.background()

        industry_text = industry_shape.text_frame
        industry_text.margin_left = Cm(1)
        industry_text.margin_top = Cm(0.5)
        industry_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        industry_title = industry_text.paragraphs[0]
        industry_title.text = "🏢 ERKANNTE BRANCHE"
        industry_title.font.name = self.fonts['body']
        industry_title.font.size = Pt(12)
        industry_title.font.bold = True
        industry_title.font.color.rgb = self.colors['white']
        industry_title.alignment = PP_ALIGN.CENTER

        industry_value = industry_text.add_paragraph()
        industry_value.text = industry_pattern
        industry_value.font.name = self.fonts['title']
        industry_value.font.size = Pt(20)
        industry_value.font.bold = True
        industry_value.font.color.rgb = self.colors['white']
        industry_value.alignment = PP_ALIGN.CENTER

        confidence_value = industry_text.add_paragraph()
        confidence_value.text = f"Konfidenz: {industry_confidence:.1f}%"
        confidence_value.font.name = self.fonts['body']
        confidence_value.font.size = Pt(10)
        confidence_value.font.color.rgb = self.colors['white']
        confidence_value.alignment = PP_ALIGN.CENTER

    def _create_gdpr_compliance_slide(self, prs, df, analysis_results):
        """Create GDPR compliance assessment slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_header(slide, "GDPR COMPLIANCE ASSESSMENT", "Datenschutz-Grundverordnung Bewertung")

        gdpr_results = analysis_results.get('gdpr_assessment', {})
        compliance_score = gdpr_results.get('compliance_score', 100)
        compliance_level = gdpr_results.get('compliance_level', 'Excellent')

        # Compliance score visual
        score_top = Cm(4.5)
        score_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.5), score_top, Cm(8), Cm(4)
        )

        # Color based on score
        score_fill = score_shape.fill
        score_fill.solid()
        if compliance_score >= 80:
            score_fill.fore_color.rgb = self.colors['green']
        elif compliance_score >= 60:
            score_fill.fore_color.rgb = self.colors['orange']
        else:
            score_fill.fore_color.rgb = self.colors['red']

        score_shape.line.fill.background()

        score_text = score_shape.text_frame
        score_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        score_title = score_text.paragraphs[0]
        score_title.text = "⚖️ COMPLIANCE SCORE"
        score_title.font.name = self.fonts['body']
        score_title.font.size = Pt(12)
        score_title.font.bold = True
        score_title.font.color.rgb = self.colors['white']
        score_title.alignment = PP_ALIGN.CENTER

        score_value = score_text.add_paragraph()
        score_value.text = f"{compliance_score}/100"
        score_value.font.name = self.fonts['title']
        score_value.font.size = Pt(24)
        score_value.font.bold = True
        score_value.font.color.rgb = self.colors['white']
        score_value.alignment = PP_ALIGN.CENTER

        score_level = score_text.add_paragraph()
        score_level.text = compliance_level
        score_level.font.name = self.fonts['body']
        score_level.font.size = Pt(11)
        score_level.font.color.rgb = self.colors['white']
        score_level.alignment = PP_ALIGN.CENTER

    def _create_recommendations_slide(self, prs, df, analysis_results):
        """Create strategic recommendations slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_header(slide, "STRATEGISCHE HANDLUNGSEMPFEHLUNGEN", "Priorisierte Maßnahmen zur Wertschöpfung")

        # Immediate actions
        immediate_top = Cm(4.5)
        immediate_shape = slide.shapes.add_textbox(Cm(1.5), immediate_top, Cm(10), Cm(4))
        immediate_frame = immediate_shape.text_frame

        immediate_title = immediate_frame.paragraphs[0]
        immediate_title.text = "🚨 SOFORTMASSNAHMEN (0-30 Tage)"
        immediate_title.font.name = self.fonts['body']
        immediate_title.font.size = Pt(11)
        immediate_title.font.bold = True
        immediate_title.font.color.rgb = self.colors['red']

        immediate_actions = [
            "Datenqualitäts-Audit durchführen",
            "GDPR Compliance-Lücken schließen",
            "Kritische Anomalien untersuchen"
        ]

        for action in immediate_actions:
            action_p = immediate_frame.add_paragraph()
            action_p.text = f"• {action}"
            action_p.font.name = self.fonts['body']
            action_p.font.size = Pt(9)
            action_p.font.color.rgb = self.colors['dark_gray']
            action_p.space_before = Pt(3)

        # Short-term initiatives
        shortterm_shape = slide.shapes.add_textbox(Cm(12.5), immediate_top, Cm(10), Cm(4))
        shortterm_frame = shortterm_shape.text_frame

        shortterm_title = shortterm_frame.paragraphs[0]
        shortterm_title.text = "📈 KURZFRISTIG (1-6 Monate)"
        shortterm_title.font.name = self.fonts['body']
        shortterm_title.font.size = Pt(11)
        shortterm_title.font.bold = True
        shortterm_title.font.color.rgb = self.colors['orange']

        shortterm_actions = [
            "Prädiktive Modelle implementieren",
            "Automatisierte Anomalie-Erkennung",
            "Dashboard-Integration"
        ]

        for action in shortterm_actions:
            action_p = shortterm_frame.add_paragraph()
            action_p.text = f"• {action}"
            action_p.font.name = self.fonts['body']
            action_p.font.size = Pt(9)
            action_p.font.color.rgb = self.colors['dark_gray']
            action_p.space_before = Pt(3)

    def _add_slide_header(self, slide, title, subtitle):
        """Add consistent header to slides"""
        # Header background
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Cm(25), Cm(2.5)
        )
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = self.colors['primary_blue']
        header_shape.line.fill.background()

        # Title
        title_shape = slide.shapes.add_textbox(Cm(1.5), Cm(0.3), Cm(20), Cm(1))
        title_frame = title_shape.text_frame
        title_frame.margin_left = 0

        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.name = self.fonts['title']
        title_p.font.size = Pt(18)
        title_p.font.bold = True
        title_p.font.color.rgb = self.colors['white']

        # Subtitle
        subtitle_shape = slide.shapes.add_textbox(Cm(1.5), Cm(1.3), Cm(20), Cm(0.8))
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.margin_left = 0

        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitle
        subtitle_p.font.name = self.fonts['subtitle']
        subtitle_p.font.size = Pt(11)
        subtitle_p.font.color.rgb = self.colors['white']

def create_german_corporate_powerpoint(df, analysis_results):
    """Main function to create German corporate PowerPoint"""
    generator = GermanCorporatePowerPoint()
    presentation = generator.create_presentation(df, analysis_results)

    # Save to buffer
    ppt_buffer = BytesIO()
    presentation.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer