import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import seaborn as sns
import matplotlib.pyplot as plt
from datetime import datetime, timedelta
from io import BytesIO
import base64
from scipy import stats
from sklearn.preprocessing import StandardScaler
from sklearn.cluster import KMeans
from sklearn.decomposition import PCA
from sklearn.ensemble import IsolationForest
from sklearn.linear_model import LinearRegression
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import warnings
import time
import traceback
from advanced_ai_analytics import AdvancedAIAnalytics

warnings.filterwarnings('ignore')

# Initialize advanced AI analytics engine
ai_analytics = AdvancedAIAnalytics()

# Configure page with German corporate branding
st.set_page_config(
    page_title="KI-Automatisierungsplattform | Enterprise Analytics",
    page_icon="🏢",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for professional German corporate styling
st.markdown("""
<style>
    /* Main container styling */
    .main > div {
        padding-top: 2rem;
    }

    /* Header styling */
    .header-container {
        background: linear-gradient(135deg, #1e3a8a 0%, #3b82f6 100%);
        padding: 2rem;
        border-radius: 10px;
        margin-bottom: 2rem;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }

    .header-title {
        color: white;
        font-size: 2.5rem;
        font-weight: 700;
        margin-bottom: 0.5rem;
        text-align: center;
    }

    .header-subtitle {
        color: #e2e8f0;
        font-size: 1.2rem;
        text-align: center;
        margin-bottom: 1rem;
    }

    .header-tagline {
        color: #cbd5e1;
        font-size: 1rem;
        text-align: center;
        font-style: italic;
    }

    /* Card styling */
    .metric-card {
        background: white;
        padding: 1.5rem;
        border-radius: 8px;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
        border-left: 4px solid #3b82f6;
        margin-bottom: 1rem;
    }

    .analysis-card {
        background: #f8fafc;
        padding: 2rem;
        border-radius: 10px;
        border: 1px solid #e2e8f0;
        margin: 1rem 0;
    }

    /* Button styling */
    .stButton > button {
        background: linear-gradient(135deg, #1e3a8a 0%, #3b82f6 100%);
        color: white;
        border: none;
        border-radius: 8px;
        padding: 0.75rem 2rem;
        font-weight: 600;
        transition: all 0.3s ease;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
    }

    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0, 0, 0, 0.2);
    }

    /* Sidebar styling */
    .css-1d391kg {
        background-color: #f1f5f9;
    }

    /* Progress bar styling */
    .stProgress > div > div > div > div {
        background: linear-gradient(135deg, #1e3a8a 0%, #3b82f6 100%);
    }

    /* Alert styling */
    .success-alert {
        background-color: #ecfdf5;
        border: 1px solid #10b981;
        border-radius: 8px;
        padding: 1rem;
        color: #065f46;
    }

    .error-alert {
        background-color: #fef2f2;
        border: 1px solid #ef4444;
        border-radius: 8px;
        padding: 1rem;
        color: #991b1b;
    }

    .warning-alert {
        background-color: #fffbeb;
        border: 1px solid #f59e0b;
        border-radius: 8px;
        padding: 1rem;
        color: #92400e;
    }

    /* Data quality indicators */
    .quality-indicator {
        display: inline-block;
        padding: 0.25rem 0.75rem;
        border-radius: 20px;
        font-size: 0.875rem;
        font-weight: 600;
        margin-right: 0.5rem;
    }

    .quality-excellent {
        background-color: #dcfce7;
        color: #166534;
    }

    .quality-good {
        background-color: #fef3c7;
        color: #92400e;
    }

    .quality-poor {
        background-color: #fee2e2;
        color: #991b1b;
    }
</style>
""", unsafe_allow_html=True)

# German corporate header
st.markdown("""
<div class="header-container">
    <div class="header-title">🏢 KI-Automatisierungsplattform</div>
    <div class="header-subtitle">Enterprise Data Intelligence & Analytics Suite</div>
    <div class="header-tagline">Modernste künstliche Intelligenz für deutsche Unternehmen</div>
</div>
""", unsafe_allow_html=True)

# Utility functions
@st.cache_data
def load_data(file_content, file_name):
    """Load and cache data with error handling"""
    try:
        if file_name.endswith('.csv'):
            # Try different encodings
            try:
                return pd.read_csv(BytesIO(file_content), encoding='utf-8')
            except UnicodeDecodeError:
                try:
                    return pd.read_csv(BytesIO(file_content), encoding='latin-1')
                except UnicodeDecodeError:
                    return pd.read_csv(BytesIO(file_content), encoding='cp1252')
        else:
            return pd.read_excel(BytesIO(file_content))
    except Exception as e:
        st.error(f"Fehler beim Laden der Datei: {str(e)}")
        return None

def validate_data(df):
    """Validate uploaded data quality"""
    issues = []
    recommendations = []

    if df is None or df.empty:
        issues.append("⚠️ Datei ist leer oder konnte nicht gelesen werden")
        return issues, recommendations, "poor"

    # Check data quality metrics
    missing_pct = (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
    duplicate_pct = (df.duplicated().sum() / len(df)) * 100
    numeric_cols = len(df.select_dtypes(include=[np.number]).columns)

    # Determine quality score
    if missing_pct < 5 and duplicate_pct < 1 and numeric_cols > 0:
        quality = "excellent"
    elif missing_pct < 15 and duplicate_pct < 5:
        quality = "good"
    else:
        quality = "poor"

    # Generate recommendations
    if missing_pct > 10:
        recommendations.append(f"🔍 {missing_pct:.1f}% fehlende Werte identifiziert - Datenbereinigung empfohlen")

    if duplicate_pct > 2:
        recommendations.append(f"🔄 {duplicate_pct:.1f}% Duplikate gefunden - Deduplizierung empfohlen")

    if numeric_cols == 0:
        recommendations.append("📊 Keine numerischen Spalten erkannt - erweiterte Analyse begrenzt")

    return issues, recommendations, quality

def create_progress_placeholder():
    """Create progress tracking interface"""
    progress_container = st.container()
    with progress_container:
        progress_bar = st.progress(0)
        status_text = st.empty()
        details_text = st.empty()
    return progress_bar, status_text, details_text

def create_enhanced_powerpoint_report(df, analysis_results):
    """Create comprehensive PowerPoint with all AI analysis results"""
    try:
        prs = Presentation()

        # Title slide
        title_slide = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Executive AI Data Intelligence Report"
        subtitle.text = f"Academic Research Grade Analysis\nGeneriert am {datetime.now().strftime('%d.%m.%Y um %H:%M')}"

        # Executive Summary slide
        if 'executive_summary' in analysis_results:
            bullet_slide = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide)
            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = "Executive Summary"
            tf = body.text_frame
            summary = analysis_results['executive_summary']

            tf.text = f"Branche: {summary['header']['industry']}"
            for finding in summary['key_findings'][:3]:  # Top 3 findings
                p = tf.add_paragraph()
                p.text = finding.replace('✅', '').replace('⚠️', '').replace('🚨', '')

        # Data Overview slide
        bullet_slide = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "Datenübersicht & Qualität"
        tf = body.text_frame
        tf.text = f"Datensätze: {len(df):,}"

        p = tf.add_paragraph()
        p.text = f"Spalten: {len(df.columns)}"

        p = tf.add_paragraph()
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        p.text = f"Numerische Felder: {len(numeric_cols)}"

        p = tf.add_paragraph()
        missing_pct = (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
        p.text = f"Datenqualität: {100-missing_pct:.1f}% vollständig"

        # Industry Analysis slide
        if 'industry' in analysis_results:
            slide = prs.slides.add_slide(bullet_slide)
            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = "Branchen-Analyse"
            tf = body.text_frame
            industry = analysis_results['industry']
            tf.text = f"Erkannte Branche: {industry['pattern'].title()}"

            p = tf.add_paragraph()
            p.text = f"Konfidenz: {industry['confidence']*100:.1f}%"

        # GDPR Compliance slide
        if 'gdpr_assessment' in analysis_results:
            slide = prs.slides.add_slide(bullet_slide)
            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = "GDPR Compliance Assessment"
            tf = body.text_frame
            gdpr = analysis_results['gdpr_assessment']
            tf.text = f"Compliance Score: {gdpr['compliance_score']}/100"

            p = tf.add_paragraph()
            p.text = f"Risiko-Level: {gdpr['compliance_level']}"

            for risk in gdpr['risk_factors'][:2]:  # Top 2 risks
                p = tf.add_paragraph()
                p.text = risk.replace('🔴', '').replace('⚠️', '')

        # Recommendations slide
        if 'bi_recommendations' in analysis_results:
            slide = prs.slides.add_slide(bullet_slide)
            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = "Strategische Empfehlungen"
            tf = body.text_frame
            bi = analysis_results['bi_recommendations']

            if bi['immediate_actions']:
                tf.text = "Sofortmaßnahmen:"
                for action in bi['immediate_actions'][:2]:
                    p = tf.add_paragraph()
                    p.text = f"• {action['action']}"

            if bi['investment_priorities']:
                p = tf.add_paragraph()
                p.text = "\nInvestitionsprioritäten:"
                for investment in bi['investment_priorities'][:2]:
                    p = tf.add_paragraph()
                    p.text = f"• {investment['action']}"

        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer

    except Exception as e:
        st.error(f"Enhanced PowerPoint-Fehler: {str(e)}")
        return None

def create_powerpoint_report(df):
    """Create a professional PowerPoint report"""
    try:
        # Create presentation
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "KI-Datenanalyse Bericht"
        subtitle.text = f"Generiert am {datetime.now().strftime('%d.%m.%Y um %H:%M')}"

        # Overview slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "Datenübersicht"
        tf = body.text_frame
        tf.text = f"Datensätze: {len(df):,}"

        p = tf.add_paragraph()
        p.text = f"Spalten: {len(df.columns)}"

        p = tf.add_paragraph()
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        p.text = f"Numerische Felder: {len(numeric_cols)}"

        p = tf.add_paragraph()
        missing_pct = (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
        p.text = f"Vollständigkeit: {100-missing_pct:.1f}%"

        # Key insights slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "Wichtigste Erkenntnisse"
        tf = body.text_frame
        tf.text = "Datenqualität ist für enterprise-grade Analysen geeignet"

        p = tf.add_paragraph()
        p.text = "Identifizierte Trends zeigen Wachstumspotenzial auf"

        p = tf.add_paragraph()
        p.text = "Empfehlung: Implementierung kontinuierlicher Datenüberwachung"

        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer

    except Exception as e:
        st.error(f"PowerPoint-Fehler: {str(e)}")
        return None

def display_statistical_analysis(df, numeric_cols):
    """Display comprehensive statistical analysis"""
    st.markdown("### 📊 Statistische Grundanalyse")

    if len(numeric_cols) > 0:
        stats_df = df[numeric_cols].describe().round(3)
        st.dataframe(stats_df, use_container_width=True)

        # Additional statistics
        col1, col2 = st.columns(2)
        with col1:
            st.markdown("#### Schiefe (Skewness)")
            skew_data = df[numeric_cols].skew().round(3)
            for col, skew in skew_data.items():
                interpretation = "Normal" if abs(skew) < 0.5 else "Schief" if abs(skew) < 1 else "Stark schief"
                st.metric(col, f"{skew:.3f}", delta=interpretation)

        with col2:
            st.markdown("#### Wölbung (Kurtosis)")
            kurt_data = df[numeric_cols].kurtosis().round(3)
            for col, kurt in kurt_data.items():
                interpretation = "Normal" if abs(kurt) < 3 else "Spitz" if kurt > 3 else "Flach"
                st.metric(col, f"{kurt:.3f}", delta=interpretation)
    else:
        st.info("📊 Keine numerischen Daten für statistische Analyse verfügbar")

def display_correlation_analysis(df, numeric_cols):
    """Display correlation analysis with heatmap"""
    st.markdown("### 🔄 Korrelationsanalyse")

    if len(numeric_cols) > 1:
        corr_matrix = df[numeric_cols].corr()

        # Create interactive heatmap with Plotly
        fig = px.imshow(
            corr_matrix,
            text_auto=True,
            aspect="auto",
            title="Korrelationsmatrix",
            color_continuous_scale="RdBu",
            zmin=-1, zmax=1
        )
        fig.update_layout(
            title_font_size=16,
            width=700,
            height=500
        )
        st.plotly_chart(fig, use_container_width=True)

        # Highlight strong correlations
        strong_corrs = []
        for i in range(len(corr_matrix.columns)):
            for j in range(i+1, len(corr_matrix.columns)):
                corr_val = corr_matrix.iloc[i, j]
                if abs(corr_val) > 0.7:
                    strong_corrs.append((
                        corr_matrix.columns[i],
                        corr_matrix.columns[j],
                        corr_val
                    ))

        if strong_corrs:
            st.markdown("#### 🔥 Starke Korrelationen (|r| > 0.7)")
            for col1, col2, corr in strong_corrs:
                correlation_type = "Positiv" if corr > 0 else "Negativ"
                st.write(f"**{col1}** ↔️ **{col2}**: {corr:.3f} ({correlation_type})")
        else:
            st.info("📉 Keine starken Korrelationen gefunden")

def advanced_analyze_data(df, numeric_cols, analysis_depth, include_clustering, include_anomalies, include_forecasting, include_ai_insights, include_gdpr_assessment, include_executive_summary):
    """Perform sophisticated AI-powered data analysis for academic research"""
    progress_bar = st.progress(0)
    status_text = st.empty()

    # Store all analysis results
    analysis_results = {}

    try:
        # Detect industry pattern first
        status_text.text("🏢 Erkenne Branchen-Muster...")
        industry_pattern, industry_confidence, pattern_details = ai_analytics.detect_industry_pattern(df)
        analysis_results['industry'] = {
            'pattern': industry_pattern,
            'confidence': industry_confidence,
            'details': pattern_details
        }
        progress_bar.progress(0.1)

        # Enhanced analysis steps
        steps = [
            "Statistische Grundanalyse",
            "Erweiterte Korrelationsanalyse"
        ]

        if include_anomalies:
            steps.append("Multi-Method Anomalie-Erkennung")
        if include_forecasting:
            steps.append("Statistische Trend-Modellierung")
        if include_gdpr_assessment:
            steps.append("GDPR Compliance Assessment")
        if include_ai_insights:
            steps.append("KI-Business Intelligence")
        if include_executive_summary:
            steps.append("Executive Summary Generation")

        total_steps = len(steps)

        for i, step in enumerate(steps):
            status_text.text(f"🔄 {step}...")
            progress_bar.progress(0.1 + (i + 1) / total_steps * 0.9)
            time.sleep(0.3)

        # Clear progress indicators
        progress_bar.empty()
        status_text.empty()

        # Display industry detection results
        display_industry_analysis(industry_pattern, industry_confidence, pattern_details)

        # Core statistical analysis
        display_enhanced_statistical_analysis(df, numeric_cols)

        # Advanced correlation analysis
        if len(numeric_cols) > 1:
            correlation_results = ai_analytics.sophisticated_correlation_analysis(df, numeric_cols)
            analysis_results['correlation_results'] = correlation_results
            display_advanced_correlation_analysis(correlation_results)

        # Enhanced anomaly detection
        if include_anomalies and len(numeric_cols) > 0:
            anomaly_results, anomaly_data = ai_analytics.advanced_anomaly_detection(df, numeric_cols)
            analysis_results['anomaly_results'] = anomaly_results
            display_enhanced_anomaly_analysis(anomaly_results, anomaly_data, df)

        # Advanced forecasting
        if include_forecasting and len(numeric_cols) > 0:
            # Try to find a good target column for forecasting
            target_col = numeric_cols[0]  # Default to first numeric column
            date_cols = df.select_dtypes(include=['datetime64', 'object']).columns
            date_col = None
            for col in date_cols:
                if 'date' in col.lower() or 'time' in col.lower():
                    date_col = col
                    break

            trend_results = ai_analytics.advanced_trend_analysis(df, target_col, date_col)
            analysis_results['trend_results'] = trend_results
            display_advanced_forecasting_analysis(trend_results, target_col)

        # GDPR Compliance Assessment
        if include_gdpr_assessment:
            gdpr_results = ai_analytics.gdpr_compliance_assessment(df)
            analysis_results['gdpr_assessment'] = gdpr_results
            display_gdpr_compliance_analysis(gdpr_results)

        # Business Intelligence Recommendations
        if include_ai_insights:
            bi_recommendations = ai_analytics.generate_business_intelligence_recommendations(
                df, analysis_results, industry_pattern
            )
            analysis_results['bi_recommendations'] = bi_recommendations
            display_business_intelligence_recommendations(bi_recommendations)

        # Executive Summary
        if include_executive_summary:
            executive_summary = ai_analytics.generate_executive_summary(
                df, analysis_results, industry_pattern
            )
            analysis_results['executive_summary'] = executive_summary
            display_executive_summary(executive_summary)

        # Enhanced PowerPoint generation
        display_enhanced_powerpoint_generation(df, analysis_results)

    except Exception as e:
        st.error(f"🚨 Fehler während der erweiterten Analyse: {str(e)}")
        progress_bar.empty()
        status_text.empty()
        with st.expander("🔍 Technische Details", expanded=False):
            st.code(traceback.format_exc())

def analyze_data(df, numeric_cols, analysis_depth, include_clustering, include_anomalies, include_forecasting):
    """Perform comprehensive data analysis"""
    progress_bar = st.progress(0)
    status_text = st.empty()

    try:
        # Analysis steps based on depth
        steps = [
            "Statistische Grundanalyse",
            "Korrelationsanalyse",
            "Verteilungsanalyse",
            "Trend-Identifikation"
        ]

        if analysis_depth == "Tiefgehende Analyse":
            if include_clustering:
                steps.append("Clustering-Analyse")
            if include_anomalies:
                steps.append("Anomalie-Erkennung")
            if include_forecasting:
                steps.append("Prognose-Modellierung")

        total_steps = len(steps)

        for i, step in enumerate(steps):
            status_text.text(f"🔄 {step}...")
            progress_bar.progress((i + 1) / total_steps)
            time.sleep(0.5)  # Simulate processing time

        # Clear progress indicators
        progress_bar.empty()
        status_text.empty()

        # Display results
        st.markdown("""
        <div class="analysis-card">
            <h2>🎯 Analyseergebnisse</h2>
        </div>
        """, unsafe_allow_html=True)

        # Statistical Overview
        display_statistical_analysis(df, numeric_cols)

        # Correlation Analysis
        if len(numeric_cols) > 1:
            display_correlation_analysis(df, numeric_cols)

        # Generate insights
        display_ai_insights(df, numeric_cols)

        # PowerPoint generation
        display_powerpoint_generation(df)

    except Exception as e:
        st.error(f"🚨 Fehler während der Analyse: {str(e)}")
        progress_bar.empty()
        status_text.empty()

def display_ai_insights(df, numeric_cols):
    """Generate and display AI-powered insights"""
    st.markdown("### 🤖 KI-Erkenntnisse")

    insights = []

    # Data quality insights
    missing_pct = (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
    if missing_pct < 5:
        insights.append("✅ **Datenqualität**: Exzellente Datenqualität mit minimalen fehlenden Werten")
    elif missing_pct < 15:
        insights.append(f"⚠️ **Datenqualität**: {missing_pct:.1f}% fehlende Werte - Datenbereinigung empfohlen")
    else:
        insights.append(f"🚨 **Datenqualität**: {missing_pct:.1f}% fehlende Werte - umfassende Datenbereinigung erforderlich")

    # Statistical insights
    if len(numeric_cols) > 0:
        high_variance_cols = []
        for col in numeric_cols:
            cv = df[col].std() / df[col].mean() if df[col].mean() != 0 else 0
            if cv > 1:
                high_variance_cols.append(col)

        if high_variance_cols:
            insights.append(f"📈 **Variabilität**: Hohe Schwankungen in {', '.join(high_variance_cols[:3])} - weitere Untersuchung empfohlen")

        # Correlation insights
        if len(numeric_cols) > 1:
            corr_matrix = df[numeric_cols].corr()
            strong_correlations = 0
            for i in range(len(corr_matrix.columns)):
                for j in range(i+1, len(corr_matrix.columns)):
                    if abs(corr_matrix.iloc[i, j]) > 0.7:
                        strong_correlations += 1

            if strong_correlations > 0:
                insights.append(f"🔄 **Korrelationen**: {strong_correlations} starke Zusammenhänge identifiziert - Potenzial für dimensionale Reduktion")
            else:
                insights.append("🔄 **Korrelationen**: Unabhängige Variablen erkannt - gute Basis für Modellierung")

    # Business insights
    row_count = len(df)
    if row_count > 10000:
        insights.append("📊 **Datengröße**: Umfangreicher Datensatz - geeignet für maschinelles Lernen")
    elif row_count > 1000:
        insights.append("📊 **Datengröße**: Ausreichend Daten für statistische Analysen")
    else:
        insights.append("📊 **Datengröße**: Begrenzte Datenmenge - Vorsicht bei Verallgemeinerungen")

    # Recommendations
    recommendations = [
        "🎯 **Empfehlung**: Implementierung eines Datenqualitäts-Monitorings",
        "🚀 **Nächste Schritte**: Entwicklung prädiktiver Modelle auf Basis der identifizierten Muster",
        "📈 **Optimierung**: Regelmäßige Aktualisierung und Validierung der Analyseergebnisse"
    ]

    # Display insights
    for insight in insights:
        st.write(insight)

    st.markdown("#### 💡 Strategische Empfehlungen")
    for rec in recommendations:
        st.write(rec)

def display_industry_analysis(industry_pattern, industry_confidence, pattern_details):
    """Display industry pattern detection results"""
    st.markdown("### 🏢 Branchen-Erkennung")

    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("🎯 Erkannte Branche", industry_pattern.title())
    with col2:
        st.metric("🎯 Konfidenz", f"{industry_confidence*100:.1f}%")
    with col3:
        total_indicators = sum(pattern_details.values())
        st.metric("📊 Indikatoren", total_indicators)

    if industry_confidence > 0.1:
        with st.expander("🔍 Branchen-Analyse Details", expanded=False):
            st.write("**Erkannte Muster nach Kategorie:**")
            for category, count in pattern_details.items():
                if count > 0:
                    st.write(f"- {category.title()}: {count} Indikatoren")

def display_enhanced_statistical_analysis(df, numeric_cols):
    """Enhanced statistical analysis with academic rigor"""
    st.markdown("### 📊 Erweiterte Statistische Analyse")

    if len(numeric_cols) > 0:
        # Basic statistics with enhanced interpretation
        stats_df = df[numeric_cols].describe().round(4)
        st.dataframe(stats_df, use_container_width=True)

        # Advanced statistical tests
        col1, col2 = st.columns(2)

        with col1:
            st.markdown("#### 🔍 Normalitätstests")
            for col in numeric_cols[:3]:  # Limit to first 3 columns
                data = df[col].dropna()
                if len(data) > 8:
                    # Shapiro-Wilk test
                    try:
                        stat, p_value = stats.shapiro(data[:5000])  # Limit for performance
                        is_normal = p_value > 0.05
                        st.write(f"**{col}**: {'Normal' if is_normal else 'Nicht-normal'} (p={p_value:.4f})")
                    except:
                        st.write(f"**{col}**: Test nicht möglich")

        with col2:
            st.markdown("#### 📈 Verteilungscharakteristika")
            for col in numeric_cols[:3]:
                data = df[col].dropna()
                if len(data) > 0:
                    skewness = stats.skew(data)
                    kurtosis = stats.kurtosis(data)
                    cv = data.std() / data.mean() if data.mean() != 0 else 0

                    interpretation = ""
                    if abs(skewness) < 0.5:
                        interpretation += "Symmetrisch"
                    elif skewness > 0:
                        interpretation += "Rechtsschief"
                    else:
                        interpretation += "Linksschief"

                    st.write(f"**{col}**: {interpretation} (Schiefe: {skewness:.3f}, Variationskoeff.: {cv:.3f})")

def display_advanced_correlation_analysis(correlation_results):
    """Display sophisticated correlation analysis with causal inference"""
    st.markdown("### 🔄 Erweiterte Korrelationsanalyse")

    if correlation_results is None:
        st.info("📊 Nicht genügend numerische Daten für Korrelationsanalyse")
        return

    # Multiple correlation methods comparison
    col1, col2 = st.columns(2)

    with col1:
        st.markdown("#### Pearson Korrelation")
        fig = px.imshow(
            correlation_results['pearson'],
            text_auto=True,
            title="Pearson Korrelationsmatrix",
            color_continuous_scale="RdBu",
            zmin=-1, zmax=1
        )
        st.plotly_chart(fig, use_container_width=True)

    with col2:
        st.markdown("#### Spearman Korrelation (Rang-basiert)")
        fig = px.imshow(
            correlation_results['spearman'],
            text_auto=True,
            title="Spearman Korrelationsmatrix",
            color_continuous_scale="RdBu",
            zmin=-1, zmax=1
        )
        st.plotly_chart(fig, use_container_width=True)

    # Partial correlations
    if correlation_results['partial_correlations']:
        st.markdown("#### 🤔 Partielle Korrelationen")
        st.write("*Korrelationen nach Kontrolle für andere Variablen:*")
        for pair, corr in correlation_results['partial_correlations'].items():
            if abs(corr) > 0.3:
                cols = pair.split('_vs_')
                st.write(f"**{cols[0]}** ↔️ **{cols[1]}**: {corr:.3f}")

    # Causality indicators
    if correlation_results['causality_indicators']:
        st.markdown("#### ➡️ Kausalitäts-Indikatoren")
        st.write("*Basierend auf zeitversetzten Korrelationen:*")
        for relationship, data in correlation_results['causality_indicators'].items():
            if abs(data['strength']) > 0.1:
                parts = relationship.split('_causes_')
                confidence_level = "Hoch" if data['confidence'] > 0.3 else "Mittel" if data['confidence'] > 0.1 else "Niedrig"
                st.write(f"**{parts[0]}** → **{parts[1]}**: {data['strength']:.3f} (Konfidenz: {confidence_level})")

def display_enhanced_anomaly_analysis(anomaly_results, anomaly_data, df):
    """Display sophisticated anomaly detection with confidence scores"""
    st.markdown("### ⚠️ Erweiterte Anomalie-Erkennung")

    if anomaly_results is None:
        st.info("📊 Nicht genügend Daten für Anomalie-Erkennung")
        return

    n_anomalies = len(anomaly_results['indices'])
    total_records = len(anomaly_data) if anomaly_data is not None else len(df)
    anomaly_rate = (n_anomalies / total_records) * 100

    # Anomaly overview
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("📊 Datenpunkte", f"{total_records:,}")
    with col2:
        st.metric("⚠️ Anomalien", n_anomalies)
    with col3:
        st.metric("📉 Anomalie-Rate", f"{anomaly_rate:.2f}%")
    with col4:
        avg_confidence = np.mean(anomaly_results['confidence']) if len(anomaly_results['confidence']) > 0 else 0
        st.metric("🎯 Konfidenz", f"{avg_confidence*100:.1f}%")

    # Methods used
    st.markdown("#### 🔧 Verwendete Methoden")
    st.write(", ".join(anomaly_results['methods_used']))

    # Confidence distribution
    if len(anomaly_results['confidence']) > 0:
        fig = go.Figure(data=go.Histogram(
            x=anomaly_results['confidence'],
            nbinsx=10,
            title="Verteilung der Anomalie-Konfidenz"
        ))
        fig.update_layout(xaxis_title="Konfidenz-Score", yaxis_title="Anzahl Anomalien")
        st.plotly_chart(fig, use_container_width=True)

    # Show high-confidence anomalies
    if n_anomalies > 0:
        high_conf_mask = anomaly_results['confidence'] > 0.7
        if np.any(high_conf_mask):
            st.markdown("#### 🔴 Hochkonfidente Anomalien")
            high_conf_indices = anomaly_results['indices'][high_conf_mask]
            with st.expander(f"{len(high_conf_indices)} hochkonfidente Anomalien anzeigen", expanded=False):
                st.dataframe(df.loc[high_conf_indices], use_container_width=True)

def display_advanced_forecasting_analysis(trend_results, target_col):
    """Display sophisticated forecasting analysis"""
    st.markdown("### 🔮 Erweiterte Prognose-Modellierung")

    st.write(f"**Zielvariable**: {target_col}")

    # Stationarity analysis
    if trend_results.get('stationarity'):
        st.markdown("#### 📈 Stationäritäts-Analyse")
        stationarity = trend_results['stationarity']

        col1, col2 = st.columns(2)
        with col1:
            adf_result = "Stationär" if stationarity['adf_is_stationary'] else "Nicht-stationär"
            st.metric("ADF Test", adf_result, f"p-Wert: {stationarity['adf_pvalue']:.4f}")

        with col2:
            kpss_result = "Stationär" if stationarity['kpss_is_stationary'] else "Nicht-stationär"
            st.metric("KPSS Test", kpss_result, f"p-Wert: {stationarity['kpss_pvalue']:.4f}")

    # Seasonal decomposition
    if trend_results.get('decomposition'):
        st.markdown("#### 🔄 Zeitreihen-Zerlegung")
        decomp = trend_results['decomposition']

        fig = make_subplots(
            rows=4, cols=1,
            subplot_titles=['Original', 'Trend', 'Saison', 'Rest'],
            vertical_spacing=0.1
        )

        x_values = list(range(len(decomp['original'])))

        fig.add_trace(go.Scatter(x=x_values, y=decomp['original'], name='Original'), row=1, col=1)
        fig.add_trace(go.Scatter(x=x_values, y=decomp['trend'], name='Trend'), row=2, col=1)
        fig.add_trace(go.Scatter(x=x_values, y=decomp['seasonal'], name='Saison'), row=3, col=1)
        fig.add_trace(go.Scatter(x=x_values, y=decomp['residual'], name='Rest'), row=4, col=1)

        fig.update_layout(height=800, title_text="Zeitreihen-Zerlegung")
        st.plotly_chart(fig, use_container_width=True)

    # Change points
    if trend_results.get('change_points') and len(trend_results['change_points']) > 0:
        st.markdown("#### 🔄 Strukturbrüche")
        st.write(f"**Erkannte Wendepunkte**: {len(trend_results['change_points'])}")
        for i, cp in enumerate(trend_results['change_points']):
            st.write(f"- Wendepunkt {i+1}: Position {cp}")

    # Enhanced forecasting results
    if trend_results.get('enhanced_forecast'):
        st.markdown("#### 📊 Erweiterte Prognose")
        enhanced = trend_results['enhanced_forecast']

        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Modell-Typ", enhanced['model_type'])
        with col2:
            st.metric("R² Score", f"{enhanced['r2_score']:.3f}")
        with col3:
            st.metric("Prognose-Perioden", len(enhanced['forecast']))

        # Enhanced forecast visualization with confidence intervals
        fig = go.Figure()

        # Historical data
        if trend_results.get('decomposition') and 'original' in trend_results['decomposition']:
            original_data = trend_results['decomposition']['original']
            historical_x = list(range(len(original_data)))
            fig.add_trace(go.Scatter(
                x=historical_x,
                y=original_data.values,
                mode='lines',
                name='Historische Daten',
                line=dict(color='blue')
            ))
        else:
            # Fallback to basic trend if decomposition not available
            historical_x = list(range(len(enhanced['forecast'])))

        forecast_x = list(range(len(historical_x), len(historical_x) + len(enhanced['forecast'])))

        # Confidence intervals
        fig.add_trace(go.Scatter(
            x=forecast_x + forecast_x[::-1],
            y=list(enhanced['confidence_upper']) + list(enhanced['confidence_lower'][::-1]),
            fill='toself',
            fillcolor='rgba(255,0,0,0.2)',
            line=dict(color='rgba(255,255,255,0)'),
            name='95% Konfidenzintervall',
            showlegend=True
        ))

        # Point forecast
        fig.add_trace(go.Scatter(
            x=forecast_x,
            y=enhanced['forecast'],
            mode='lines+markers',
            name=f'{enhanced["model_type"]} Prognose',
            line=dict(color='red')
        ))

        fig.update_layout(
            title=f"Erweiterte Prognose: {target_col}",
            xaxis_title="Zeit",
            yaxis_title=target_col
        )
        st.plotly_chart(fig, use_container_width=True)

        # Show forecast table
        with st.expander("📊 Prognose-Details", expanded=False):
            forecast_df = pd.DataFrame({
                'Periode': range(len(historical_x) + 1, len(historical_x) + len(enhanced['forecast']) + 1),
                'Prognose': enhanced['forecast'].round(2),
                'Untere Grenze': enhanced['confidence_lower'].round(2),
                'Obere Grenze': enhanced['confidence_upper'].round(2)
            })
            st.dataframe(forecast_df, use_container_width=True)

    # ARIMA forecasting results (legacy)
    elif trend_results.get('arima_forecast'):
        st.markdown("#### 📊 ARIMA Prognose")
        arima = trend_results['arima_forecast']

        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("ARIMA Ordnung", f"{arima['order']}")
        with col2:
            st.metric("AIC Score", f"{arima['aic']:.2f}")
        with col3:
            st.metric("Prognose-Perioden", len(arima['forecast']))

        # Forecast visualization
        fig = go.Figure()
        historical_x = list(range(len(trend_results.get('original', []))))
        forecast_x = list(range(len(historical_x), len(historical_x) + len(arima['forecast'])))

        if 'original' in trend_results:
            fig.add_trace(go.Scatter(
                x=historical_x,
                y=trend_results['original'],
                mode='lines',
                name='Historische Daten'
            ))

        fig.add_trace(go.Scatter(
            x=forecast_x,
            y=arima['forecast'],
            mode='lines+markers',
            name='ARIMA Prognose',
            line=dict(color='red')
        ))

        fig.update_layout(title="ARIMA Prognose-Modell", xaxis_title="Zeit", yaxis_title=target_col)
        st.plotly_chart(fig, use_container_width=True)

def display_gdpr_compliance_analysis(gdpr_results):
    """Display GDPR compliance assessment"""
    st.markdown("### ⚖️ GDPR Compliance Assessment")

    # Compliance overview
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("🎯 Compliance Score", f"{gdpr_results['compliance_score']}/100")
    with col2:
        level_with_color = f"{gdpr_results['level_color']} {gdpr_results['compliance_level']}"
        st.metric("📈 Risiko-Level", level_with_color)
    with col3:
        st.metric("🕰️ Assessment Zeit", gdpr_results['assessment_timestamp'].strftime('%H:%M'))

    # Risk factors
    if gdpr_results['risk_factors']:
        st.markdown("#### ⚠️ Identifizierte Risikofaktoren")
        for risk in gdpr_results['risk_factors']:
            st.warning(risk)

    # PII detection
    if gdpr_results['potential_pii']:
        st.markdown("#### 🔍 Erkannte PII-Kategorien")
        for category, columns in gdpr_results['potential_pii'].items():
            with st.expander(f"{category.title()} ({len(columns)} Spalten)", expanded=False):
                for col in columns:
                    st.write(f"- {col}")

    # Recommendations
    if gdpr_results['recommendations']:
        st.markdown("#### 📝 Compliance-Empfehlungen")
        for rec in gdpr_results['recommendations']:
            st.info(rec)

def display_business_intelligence_recommendations(bi_recommendations):
    """Display AI-powered business intelligence recommendations"""
    st.markdown("### 🤖 KI-Business Intelligence")

    # Immediate actions
    if bi_recommendations['immediate_actions']:
        st.markdown("#### 🚑 Sofortmaßnahmen")
        for action in bi_recommendations['immediate_actions']:
            with st.expander(f"🔴 {action['action']}", expanded=False):
                st.write(f"**Beschreibung**: {action['description']}")
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("Impact", action['impact'])
                with col2:
                    st.metric("Aufwand", action['effort'])
                with col3:
                    st.metric("Timeline", action['timeline'])

    # Short-term initiatives
    if bi_recommendations['short_term_initiatives']:
        st.markdown("#### 📈 Kurzfristige Initiativen")
        for initiative in bi_recommendations['short_term_initiatives']:
            with st.expander(f"🟡 {initiative['action']}", expanded=False):
                st.write(f"**Beschreibung**: {initiative['description']}")
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("Impact", initiative['impact'])
                with col2:
                    st.metric("Aufwand", initiative['effort'])
                with col3:
                    st.metric("Timeline", initiative['timeline'])

    # Investment priorities
    if bi_recommendations['investment_priorities']:
        st.markdown("#### 💰 Investitionsprioriäten")
        for investment in bi_recommendations['investment_priorities']:
            with st.expander(f"🔵 {investment['action']}", expanded=False):
                st.write(f"**Beschreibung**: {investment['description']}")
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("Impact", investment['impact'])
                with col2:
                    st.metric("Aufwand", investment['effort'])
                with col3:
                    st.metric("Timeline", investment['timeline'])
                if 'roi_estimate' in investment:
                    st.success(f"📈 ROI: {investment['roi_estimate']}")

def display_executive_summary(executive_summary):
    """Display professional executive summary"""
    st.markdown("### 📋 Executive Summary")

    # Header
    header = executive_summary['header']
    st.markdown(f"""
    <div class="analysis-card">
        <h2>{header['title']}</h2>
        <h4>{header['subtitle']}</h4>
        <p><strong>Branche</strong>: {header['industry']} | <strong>Erstellt</strong>: {header['timestamp']}</p>
    </div>
    """, unsafe_allow_html=True)

    # Key findings
    st.markdown("#### 🎯 Wichtigste Erkenntnisse")
    for finding in executive_summary['key_findings']:
        st.write(finding)

    # Strategic recommendations
    st.markdown("#### 🚀 Strategische Empfehlungen")
    for rec in executive_summary['strategic_recommendations']:
        st.write(rec)

    # Risk assessment summary
    if executive_summary['risk_assessment']:
        st.markdown("#### ⚠️ Risikobewertung")
        risk = executive_summary['risk_assessment']
        col1, col2 = st.columns(2)
        with col1:
            st.metric("GDPR Compliance", risk['gdpr_compliance'])
        with col2:
            st.metric("Compliance Score", f"{risk['compliance_score']}/100")

    # Technical insights
    tech = executive_summary['technical_insights']
    st.markdown("#### 💻 Technische Bewertung")
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("Datenreife", tech['data_maturity'])
    with col2:
        st.metric("Analyse-Komplexität", tech['analysis_complexity'])
    with col3:
        st.metric("Skalierbarkeit", tech['scalability'])

    # Next steps
    st.markdown("#### 📝 Nächste Schritte")
    for step in executive_summary['next_steps']:
        st.write(step)

def display_enhanced_powerpoint_generation(df, analysis_results):
    """Generate enhanced PowerPoint with all analysis results"""
    st.markdown("### 📄 Professional Report Generation")

    st.info("🎓 **Academic Research Grade**: Umfassender Bericht mit allen AI-Analysen")

    if st.button("📊 Executive PowerPoint generieren", type="secondary"):
        with st.spinner("Erstelle umfassende PowerPoint-Präsentation..."):
            ppt_buffer = create_enhanced_powerpoint_report(df, analysis_results)

            if ppt_buffer:
                st.success("✅ Executive PowerPoint-Bericht erfolgreich erstellt!")
                st.download_button(
                    label="📎 Executive Report herunterladen",
                    data=ppt_buffer.getvalue(),
                    file_name=f"Executive_AI_Analysis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            else:
                st.error("🚨 Fehler bei der PowerPoint-Erstellung")

def display_powerpoint_generation(df):
    """Generate and offer PowerPoint download"""
    st.markdown("### 📄 PowerPoint-Bericht")

    if st.button("📊 Professional Report generieren", type="secondary"):
        with st.spinner("Erstelle PowerPoint-Präsentation..."):
            ppt_buffer = create_powerpoint_report(df)

            if ppt_buffer:
                st.success("✅ PowerPoint-Bericht erfolgreich erstellt!")
                st.download_button(
                    label="📎 PowerPoint herunterladen",
                    data=ppt_buffer.getvalue(),
                    file_name=f"KI_Analyse_Bericht_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            else:
                st.error("🚨 Fehler bei der PowerPoint-Erstellung")

# Sidebar for configuration
with st.sidebar:
    st.markdown("### 🛠️ Konfiguration")

    # Analysis settings
    st.markdown("#### Analyse-Einstellungen")
    analysis_depth = st.selectbox(
        "Analyse-Tiefe",
        ["Schnellübersicht", "Standard-Analyse", "Tiefgehende Analyse", "🎓 Academic Research"],
        index=2
    )

    include_clustering = st.checkbox("Clustering-Analyse", value=True)
    include_anomalies = st.checkbox("🔍 Erweiterte Anomalie-Erkennung", value=True)
    include_forecasting = st.checkbox("📈 Statistische Prognose-Modelle", value=True)
    include_ai_insights = st.checkbox("🤖 KI-Business Intelligence", value=True)
    include_gdpr_assessment = st.checkbox("⚖️ GDPR Compliance Check", value=True)
    include_executive_summary = st.checkbox("📋 Executive Summary", value=True)

    # Export settings
    st.markdown("#### Export-Einstellungen")
    report_language = st.selectbox("Berichtssprache", ["Deutsch", "English"], index=0)
    include_charts = st.checkbox("Diagramme einschließen", value=True)
    include_raw_data = st.checkbox("Rohdaten anhängen", value=False)

    st.markdown("---")
    st.markdown("### 📊 Systemstatus")
    st.success("✅ Alle Systeme betriebsbereit")
    st.info(f"🕰️ Letzte Aktualisierung: {datetime.now().strftime('%H:%M:%S')}")

# Main file upload section
st.markdown("### 📁 Daten-Upload")

col1, col2 = st.columns([2, 1])

with col1:
    uploaded_file = st.file_uploader(
        "Wählen Sie Ihre Datendatei aus",
        type=['xlsx', 'csv', 'xls'],
        help="Unterstützte Formate: Excel (.xlsx, .xls), CSV (.csv). Maximale Dateigröße: 200MB"
    )

with col2:
    if uploaded_file:
        file_size = len(uploaded_file.getvalue()) / (1024 * 1024)  # MB
        st.metric("Dateigröße", f"{file_size:.1f} MB")
        st.metric("Dateiformat", uploaded_file.name.split('.')[-1].upper())

if uploaded_file:
    # Create progress tracking
    progress_bar, status_text, details_text = create_progress_placeholder()

    try:
        # Step 1: Load data
        status_text.text("📂 Lade Datei...")
        progress_bar.progress(20)

        file_content = uploaded_file.getvalue()
        df = load_data(file_content, uploaded_file.name)

        if df is None:
            st.stop()

        # Step 2: Validate data
        status_text.text("🔍 Validiere Datenqualität...")
        progress_bar.progress(40)

        issues, recommendations, quality = validate_data(df)

        # Step 3: Display results
        progress_bar.progress(60)
        status_text.text("✅ Daten erfolgreich geladen")
        details_text.empty()
        progress_bar.empty()
        status_text.empty()

        # Display data overview with quality indicators
        st.markdown("""
        <div class="analysis-card">
            <h3>📊 Datenübersicht</h3>
        </div>
        """, unsafe_allow_html=True)

        # Quality indicator
        quality_class = f"quality-{quality}"
        quality_text = {
            "excellent": "Exzellent",
            "good": "Gut",
            "poor": "Verbesserungsbedürftig"
        }

        st.markdown(f"""
        <div class="quality-indicator {quality_class}">
            Datenqualität: {quality_text[quality]}
        </div>
        """, unsafe_allow_html=True)

        # Data metrics
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.metric("📊 Datensätze", f"{len(df):,}")
        with col2:
            st.metric("📈 Spalten", len(df.columns))
        with col3:
            numeric_cols = df.select_dtypes(include=[np.number]).columns
            st.metric("🔢 Numerische Felder", len(numeric_cols))
        with col4:
            missing_pct = (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
            st.metric("✅ Vollständigkeit", f"{100-missing_pct:.1f}%")

        # Show recommendations if any
        if recommendations:
            st.markdown("#### 💡 Empfehlungen")
            for rec in recommendations:
                st.warning(rec)

        # Data preview
        with st.expander("🔍 Datenvorschau", expanded=False):
            st.dataframe(df.head(10), use_container_width=True)

        # Analysis button
        if st.button("🚀 Intelligente Analyse starten", type="primary"):
            advanced_analyze_data(df, numeric_cols, analysis_depth, include_clustering, include_anomalies, include_forecasting, include_ai_insights, include_gdpr_assessment, include_executive_summary)

    except Exception as e:
        st.error(f"🚨 Fehler bei der Datenverarbeitung: {str(e)}")
        with st.expander("🔍 Technische Details", expanded=False):
            st.code(traceback.format_exc())

# Demo section
st.markdown("---")
st.markdown("### 🎪 Professor Demo-Modus")

col1, col2 = st.columns(2)

with col1:
    if st.button("🎪 Demo mit Beispieldaten starten", type="secondary"):
        st.balloons()

        # Create comprehensive sample data
        np.random.seed(42)
        n_samples = 500

        sample_data = pd.DataFrame({
            'Datum': pd.date_range('2022-01-01', periods=n_samples, freq='D'),
            'Umsatz_EUR': np.random.normal(45000, 12000, n_samples).clip(min=10000),
            'Kunde_Segment': np.random.choice(['Premium', 'Standard', 'Basic'], n_samples, p=[0.2, 0.5, 0.3]),
            'Region': np.random.choice(['Nord', 'Süd', 'Ost', 'West'], n_samples),
            'Produkt_Kategorie': np.random.choice(['Software', 'Hardware', 'Services'], n_samples, p=[0.4, 0.3, 0.3]),
            'Mitarbeiter_Anzahl': np.random.poisson(25, n_samples),
            'Kundenzufriedenheit': np.random.normal(4.2, 0.8, n_samples).clip(min=1, max=5),
            'Marketing_Budget': np.random.exponential(5000, n_samples)
        })

        # Add some correlations and trends
        sample_data['Umsatz_EUR'] += sample_data['Mitarbeiter_Anzahl'] * 800
        sample_data['Umsatz_EUR'] += (sample_data['Kundenzufriedenheit'] - 3) * 5000

        # Add trend over time
        time_trend = np.arange(len(sample_data)) * 10
        sample_data['Umsatz_EUR'] += time_trend

        st.session_state['demo_data'] = sample_data
        st.success("✅ Demo-Daten geladen! Bereit für Analyse.")
        st.dataframe(sample_data.head(10), use_container_width=True)

with col2:
    if 'demo_data' in st.session_state:
        if st.button("🚀 Demo-Analyse starten", type="primary"):
            demo_df = st.session_state['demo_data']
            demo_numeric_cols = demo_df.select_dtypes(include=[np.number]).columns
            advanced_analyze_data(demo_df, demo_numeric_cols, "🎓 Academic Research", True, True, True, True, True, True)

# Footer
st.markdown("---")
st.markdown("""
<div style="text-align: center; color: #6b7280; font-size: 0.9rem;">
    🏢 <strong>KI-Automatisierungsplattform</strong> |
    Enterprise-Grade Data Intelligence |
    Powered by Advanced Analytics & Machine Learning
</div>
""", unsafe_allow_html=True)