import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from datetime import datetime, timedelta
from io import BytesIO
import warnings
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from docx import Document
import re

warnings.filterwarnings('ignore')

# Configure page with SAP-style branding
st.set_page_config(
    page_title="Universal Automation Platform - SAP Edition",
    page_icon="<�",
    layout="wide",
    initial_sidebar_state="expanded"
)

# SAP Corporate CSS Styling
st.markdown("""
<style>
    /* SAP Color Palette */
    :root {
        --sap-blue: #0070f2;
        --sap-dark-blue: #0854a0;
        --sap-light-blue: #d1efff;
        --sap-gray: #6a6d70;
        --sap-light-gray: #f7f7f7;
        --sap-green: #30914c;
        --sap-orange: #ff6600;
        --sap-red: #cc1919;
    }

    /* Main container styling */
    .main > div {
        padding-top: 1rem;
    }

    /* SAP Header styling */
    .sap-header {
        background: linear-gradient(135deg, var(--sap-blue) 0%, var(--sap-dark-blue) 100%);
        padding: 2rem;
        border-radius: 8px;
        margin-bottom: 2rem;
        box-shadow: 0 4px 12px rgba(0, 112, 242, 0.15);
        border-left: 4px solid var(--sap-orange);
    }

    .sap-title {
        color: white;
        font-size: 2.2rem;
        font-weight: 700;
        margin-bottom: 0.5rem;
        text-align: left;
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    }

    .sap-subtitle {
        color: var(--sap-light-blue);
        font-size: 1.1rem;
        text-align: left;
        margin-bottom: 0.5rem;
    }

    .sap-tagline {
        color: #e8f4ff;
        font-size: 0.95rem;
        text-align: left;
        font-style: italic;
    }

    /* SAP Card styling */
    .sap-card {
        background: white;
        padding: 1.5rem;
        border-radius: 8px;
        border: 1px solid #e5e5e5;
        margin: 1rem 0;
        box-shadow: 0 2px 8px rgba(0, 0, 0, 0.08);
        border-left: 4px solid var(--sap-blue);
    }

    /* SAP Button styling */
    .stButton > button {
        background: var(--sap-blue);
        color: white;
        border: none;
        border-radius: 4px;
        padding: 0.6rem 1.5rem;
        font-weight: 600;
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        transition: all 0.2s ease;
        box-shadow: 0 2px 4px rgba(0, 112, 242, 0.2);
    }

    .stButton > button:hover {
        background: var(--sap-dark-blue);
        transform: translateY(-1px);
        box-shadow: 0 4px 8px rgba(0, 112, 242, 0.3);
    }

    /* SAP Status indicators */
    .sap-status-success {
        background-color: #e8f5e8;
        color: var(--sap-green);
        padding: 0.5rem 1rem;
        border-radius: 4px;
        border-left: 4px solid var(--sap-green);
        margin: 0.5rem 0;
    }

    .sap-status-warning {
        background-color: #fff3e0;
        color: var(--sap-orange);
        padding: 0.5rem 1rem;
        border-radius: 4px;
        border-left: 4px solid var(--sap-orange);
        margin: 0.5rem 0;
    }

    .sap-status-error {
        background-color: #ffebee;
        color: var(--sap-red);
        padding: 0.5rem 1rem;
        border-radius: 4px;
        border-left: 4px solid var(--sap-red);
        margin: 0.5rem 0;
    }

    /* SAP Metrics */
    .sap-metric {
        background: var(--sap-light-gray);
        padding: 1rem;
        border-radius: 4px;
        text-align: center;
        border: 1px solid #e0e0e0;
    }

    .sap-metric-value {
        font-size: 1.8rem;
        font-weight: 700;
        color: var(--sap-blue);
        margin-bottom: 0.25rem;
    }

    .sap-metric-label {
        font-size: 0.9rem;
        color: var(--sap-gray);
        font-weight: 500;
    }

    /* SAP Sidebar styling */
    .css-1d391kg {
        background-color: var(--sap-light-gray);
    }

    /* SAP File uploader */
    .stFileUploader > div > div > div > div {
        border: 2px dashed var(--sap-blue);
        border-radius: 8px;
    }
</style>
""", unsafe_allow_html=True)

# SAP Corporate Header
st.markdown("""
<div class="sap-header">
    <div class="sap-title"><� Universal Automation Platform</div>
    <div class="sap-subtitle">SAP Edition - Enterprise Data Processing Suite</div>
    <div class="sap-tagline">Professional data processing and reporting for modern enterprises</div>
</div>
""", unsafe_allow_html=True)

class SAPTemplateAnalyzer:
    """SAP PowerPoint template analysis and learning system"""

    def __init__(self):
        self.template_data = None
        self.learned_styles = {}
        self.layout_patterns = []

    def analyze_template(self, pptx_file):
        """Analyze template PowerPoint to extract complete design elements"""
        try:
            prs = Presentation(BytesIO(pptx_file))

            # Store the original presentation for layout copying
            self.template_presentation = prs

            template_info = {
                'slide_count': len(prs.slides),
                'layouts': [],
                'colors': [],
                'fonts': [],
                'slide_patterns': [],
                'master_layouts': [],
                'backgrounds': [],
                'logos': [],
                'images': [],
                'content_patterns': {
                    'title_positions': [],
                    'content_areas': [],
                    'bullet_styles': [],
                    'shape_arrangements': []
                },
                'slide_masters': []
            }

            # Extract slide masters first
            for master in prs.slide_masters:
                master_info = {
                    'name': master.name if hasattr(master, 'name') else 'Master',
                    'layouts': [],
                    'background': self._extract_background(master),
                    'placeholders': []
                }

                # Extract master layouts
                for layout in master.slide_layouts:
                    layout_info = {
                        'name': layout.name if hasattr(layout, 'name') else 'Layout',
                        'placeholders': self._extract_placeholders(layout),
                        'background': self._extract_background(layout)
                    }
                    master_info['layouts'].append(layout_info)

                template_info['slide_masters'].append(master_info)

            # Extract slide layouts and patterns
            for i, slide in enumerate(prs.slides):
                slide_info = {
                    'slide_number': i + 1,
                    'layout_type': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown',
                    'layout_index': slide.slide_layout.element.get('idx') if hasattr(slide.slide_layout.element, 'get') else 0,
                    'layout_object': slide.slide_layout,  # Store reference to actual layout
                    'shapes': [],
                    'colors': [],
                    'fonts': [],
                    'title_shape': None,
                    'content_shapes': [],
                    'background_info': self._extract_background(slide),
                    'spacing_patterns': {},
                    'images': [],
                    'logos': []
                }

                # Analyze shapes and their properties
                for shape in slide.shapes:
                    shape_info = {
                        'type': shape.shape_type,
                        'left': shape.left.inches if hasattr(shape.left, 'inches') else None,
                        'top': shape.top.inches if hasattr(shape.top, 'inches') else None,
                        'width': shape.width.inches if hasattr(shape.width, 'inches') else None,
                        'height': shape.height.inches if hasattr(shape.height, 'inches') else None,
                        'is_title': False,
                        'is_placeholder': hasattr(shape, 'is_placeholder') and shape.is_placeholder,
                        'placeholder_format': None,
                        'text_alignment': None,
                        'indent_level': 0,
                        'is_image': False,
                        'is_logo': False,
                        'shape_object': shape  # Store reference to actual shape
                    }

                    # Detect images and potential logos
                    if shape.shape_type == 13:  # Picture shape type
                        shape_info['is_image'] = True

                        # Detect logos by position and size
                        if (shape_info['top'] and shape_info['top'] < 1.5 and
                            shape_info['width'] and shape_info['width'] < 3):
                            shape_info['is_logo'] = True
                            slide_info['logos'].append(shape_info.copy())
                            template_info['logos'].append({
                                'slide': i + 1,
                                'position': {
                                    'left': shape_info['left'],
                                    'top': shape_info['top'],
                                    'width': shape_info['width'],
                                    'height': shape_info['height']
                                },
                                'shape_object': shape
                            })
                        else:
                            slide_info['images'].append(shape_info.copy())
                            template_info['images'].append({
                                'slide': i + 1,
                                'position': {
                                    'left': shape_info['left'],
                                    'top': shape_info['top'],
                                    'width': shape_info['width'],
                                    'height': shape_info['height']
                                },
                                'shape_object': shape
                            })

                    # Identify title shapes
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        try:
                            if shape.placeholder_format.type == 1:  # Title placeholder
                                shape_info['is_title'] = True
                                slide_info['title_shape'] = shape_info.copy()
                                template_info['content_patterns']['title_positions'].append({
                                    'left': shape_info['left'],
                                    'top': shape_info['top'],
                                    'width': shape_info['width'],
                                    'height': shape_info['height']
                                })
                        except:
                            pass

                    # Check if it's a title shape by position (fallback)
                    if not shape_info['is_title'] and shape_info['top'] and shape_info['top'] < 2:
                        shape_info['is_title'] = True
                        slide_info['title_shape'] = shape_info.copy()

                    # Extract text formatting if available
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        try:
                            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                                # Capture paragraph-level formatting
                                bullet_info = {
                                    'level': paragraph.level,
                                    'alignment': paragraph.alignment,
                                    'space_before': getattr(paragraph.space_before, 'pt', None) if paragraph.space_before else None,
                                    'space_after': getattr(paragraph.space_after, 'pt', None) if paragraph.space_after else None,
                                    'line_spacing': getattr(paragraph.line_spacing, 'pt', None) if paragraph.line_spacing else None
                                }
                                template_info['content_patterns']['bullet_styles'].append(bullet_info)

                                for run in paragraph.runs:
                                    if run.font.name:
                                        slide_info['fonts'].append(run.font.name)
                                    if hasattr(run.font, 'color') and run.font.color.rgb:
                                        color_hex = f"#{run.font.color.rgb}"
                                        slide_info['colors'].append(color_hex)

                            # Store content area information
                            if not shape_info['is_title']:
                                content_area = {
                                    'left': shape_info['left'],
                                    'top': shape_info['top'],
                                    'width': shape_info['width'],
                                    'height': shape_info['height'],
                                    'paragraph_count': len(shape.text_frame.paragraphs)
                                }
                                template_info['content_patterns']['content_areas'].append(content_area)
                                slide_info['content_shapes'].append(shape_info.copy())

                        except:
                            pass

                    slide_info['shapes'].append(shape_info)

                template_info['slide_patterns'].append(slide_info)

            # Extract dominant colors and fonts
            all_fonts = []
            all_colors = []
            for slide in template_info['slide_patterns']:
                all_fonts.extend(slide['fonts'])
                all_colors.extend(slide['colors'])

            # Get most common fonts and colors
            from collections import Counter
            template_info['fonts'] = [font for font, count in Counter(all_fonts).most_common(5)]
            template_info['colors'] = [color for color, count in Counter(all_colors).most_common(5)]

            self.template_data = template_info
            self.learned_styles = {
                'primary_font': template_info['fonts'][0] if template_info['fonts'] else 'Calibri',
                'primary_colors': template_info['colors'][:3] if template_info['colors'] else ['#0070f2', '#ff6600', '#30914c'],
                'layout_patterns': template_info['slide_patterns'],
                'title_positioning': self._calculate_average_position(template_info['content_patterns']['title_positions']),
                'content_positioning': self._calculate_average_position(template_info['content_patterns']['content_areas']),
                'bullet_formatting': self._analyze_bullet_patterns(template_info['content_patterns']['bullet_styles']),
                'slide_sequencing': self._analyze_slide_sequence(template_info['slide_patterns'])
            }

            return template_info

        except Exception as e:
            st.error(f"Error analyzing template: {str(e)}")
            return None

    def _calculate_average_position(self, positions):
        """Calculate average position from a list of position dictionaries"""
        if not positions:
            return None

        avg_pos = {
            'left': sum(pos['left'] for pos in positions if pos['left']) / len([pos for pos in positions if pos['left']]),
            'top': sum(pos['top'] for pos in positions if pos['top']) / len([pos for pos in positions if pos['top']]),
            'width': sum(pos['width'] for pos in positions if pos['width']) / len([pos for pos in positions if pos['width']]),
            'height': sum(pos['height'] for pos in positions if pos['height']) / len([pos for pos in positions if pos['height']])
        }
        return avg_pos

    def _analyze_bullet_patterns(self, bullet_styles):
        """Analyze bullet formatting patterns"""
        if not bullet_styles:
            return {}

        from collections import Counter
        levels = [style['level'] for style in bullet_styles if style['level'] is not None]
        alignments = [style['alignment'] for style in bullet_styles if style['alignment'] is not None]

        return {
            'common_levels': Counter(levels).most_common(3),
            'common_alignments': Counter(alignments).most_common(2),
            'avg_space_before': sum(style['space_before'] for style in bullet_styles if style['space_before']) / max(1, len([s for s in bullet_styles if s['space_before']])),
            'avg_space_after': sum(style['space_after'] for style in bullet_styles if style['space_after']) / max(1, len([s for s in bullet_styles if s['space_after']]))
        }

    def _analyze_slide_sequence(self, slide_patterns):
        """Analyze slide sequencing patterns"""
        if not slide_patterns:
            return []

        sequence = []
        for slide in slide_patterns:
            slide_type = 'title' if slide['slide_number'] == 1 else 'content'
            if slide['title_shape'] and len(slide['content_shapes']) > 0:
                slide_type = 'title_and_content'
            elif slide['title_shape']:
                slide_type = 'title_only'
            elif len(slide['content_shapes']) > 0:
                slide_type = 'content_only'

            sequence.append({
                'slide_number': slide['slide_number'],
                'type': slide_type,
                'layout_type': slide['layout_type'],
                'shape_count': len(slide['shapes'])
            })

        return sequence

    def _extract_background(self, slide_or_layout):
        """Extract background information from slide or layout"""
        try:
            background_info = {
                'fill_type': None,
                'color': None,
                'image': None,
                'gradient': None
            }

            # Try to get background fill
            if hasattr(slide_or_layout, 'background'):
                bg = slide_or_layout.background
                if hasattr(bg, 'fill'):
                    fill = bg.fill
                    if hasattr(fill, 'type'):
                        background_info['fill_type'] = str(fill.type)

                        # Solid color background
                        if hasattr(fill, 'fore_color') and fill.fore_color:
                            try:
                                background_info['color'] = f"#{fill.fore_color.rgb}"
                            except:
                                pass

            return background_info
        except:
            return {'fill_type': None, 'color': None, 'image': None, 'gradient': None}

    def _extract_placeholders(self, layout):
        """Extract placeholder information from layout"""
        try:
            placeholders = []
            for placeholder in layout.placeholders:
                placeholder_info = {
                    'type': placeholder.placeholder_format.type if hasattr(placeholder, 'placeholder_format') else None,
                    'idx': placeholder.placeholder_format.idx if hasattr(placeholder, 'placeholder_format') else None,
                    'left': placeholder.left.inches if hasattr(placeholder.left, 'inches') else None,
                    'top': placeholder.top.inches if hasattr(placeholder.top, 'inches') else None,
                    'width': placeholder.width.inches if hasattr(placeholder.width, 'inches') else None,
                    'height': placeholder.height.inches if hasattr(placeholder.height, 'inches') else None,
                    'object': placeholder  # Store reference for copying
                }
                placeholders.append(placeholder_info)
            return placeholders
        except:
            return []

    def get_template_summary(self):
        """Get a summary of the learned template"""
        if not self.template_data:
            return None

        return {
            'total_slides': self.template_data['slide_count'],
            'unique_layouts': len(set([slide['layout_type'] for slide in self.template_data['slide_patterns']])),
            'dominant_fonts': self.template_data['fonts'][:3],
            'dominant_colors': self.template_data['colors'][:3],
            'avg_shapes_per_slide': sum(len(slide['shapes']) for slide in self.template_data['slide_patterns']) / max(1, len(self.template_data['slide_patterns'])),
            'content_areas_detected': len(self.template_data['content_patterns']['content_areas']),
            'bullet_styles_found': len(set([str(style) for style in self.template_data['content_patterns']['bullet_styles']])),
            'logos_found': len(self.template_data.get('logos', [])),
            'images_found': len(self.template_data.get('images', [])),
            'backgrounds_detected': len([s for s in self.template_data['slide_patterns'] if s.get('background_info', {}).get('color')]),
            'slide_masters': len(self.template_data.get('slide_masters', []))
        }

    def apply_template_structure(self, slide, slide_type='content'):
        """Copy template design completely to new slide"""
        if not hasattr(self, 'template_presentation') or not self.template_presentation:
            return slide

        try:
            # Find the best matching template slide
            template_slide = self._get_matching_template_slide(slide_type)
            if not template_slide:
                return slide

            # Copy background from template
            self._copy_background(template_slide, slide)

            # Copy all non-content shapes (logos, images, design elements)
            self._copy_design_shapes(template_slide, slide)

            # Apply template fonts and colors to existing content
            self._apply_template_formatting(slide)

            return slide

        except Exception as e:
            return slide

    def _get_matching_template_slide(self, slide_type):
        """Get the best matching slide from template based on slide type"""
        if not hasattr(self, 'template_presentation'):
            return None

        try:
            # For title slides, use the first slide
            if slide_type == 'title':
                return self.template_presentation.slides[0] if len(self.template_presentation.slides) > 0 else None

            # For content slides, find a slide with content layout
            for slide in self.template_presentation.slides[1:]:  # Skip title slide
                if len([s for s in slide.shapes if hasattr(s, 'text_frame') and s.text_frame]) > 0:
                    return slide

            # Fallback to second slide if available
            return self.template_presentation.slides[1] if len(self.template_presentation.slides) > 1 else self.template_presentation.slides[0]

        except:
            return None

    def _copy_background(self, template_slide, target_slide):
        """Copy background from template slide to target slide"""
        try:
            # Copy background fill
            if hasattr(template_slide, 'background') and hasattr(target_slide, 'background'):
                template_bg = template_slide.background
                target_bg = target_slide.background

                if hasattr(template_bg, 'fill') and hasattr(target_bg, 'fill'):
                    template_fill = template_bg.fill
                    target_fill = target_bg.fill

                    # Copy solid color backgrounds
                    if hasattr(template_fill, 'fore_color') and template_fill.fore_color:
                        try:
                            target_fill.solid()
                            target_fill.fore_color.rgb = template_fill.fore_color.rgb
                        except:
                            pass

        except:
            pass

    def _copy_design_shapes(self, template_slide, target_slide):
        """Copy non-content shapes (logos, images, design elements) from template"""
        try:
            for shape in template_slide.shapes:
                # Skip text placeholders - we'll handle content separately
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    continue

                # Copy images, logos, and design shapes
                if shape.shape_type == 13:  # Picture
                    self._copy_image_shape(shape, target_slide)
                elif hasattr(shape, 'shape_type') and shape.shape_type in [1, 5, 9]:  # AutoShape, Rectangle, etc.
                    self._copy_shape(shape, target_slide)

        except:
            pass

    def _copy_image_shape(self, template_shape, target_slide):
        """Copy an image shape from template to target slide"""
        try:
            # Get image data
            image_part = template_shape.image.image_part
            image_bytes = image_part.blob

            # Add image to target slide with same position and size
            left = template_shape.left
            top = template_shape.top
            width = template_shape.width
            height = template_shape.height

            # Create image in target slide
            pic = target_slide.shapes.add_picture(
                BytesIO(image_bytes), left, top, width, height
            )

        except:
            pass

    def _copy_shape(self, template_shape, target_slide):
        """Copy a shape from template to target slide"""
        try:
            # This is complex - for now, copy basic rectangle shapes
            if template_shape.shape_type == 1:  # AutoShape
                left = template_shape.left
                top = template_shape.top
                width = template_shape.width
                height = template_shape.height

                # Add rectangle shape
                shape = target_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top, width, height
                )

                # Copy fill
                if hasattr(template_shape, 'fill') and hasattr(shape, 'fill'):
                    template_fill = template_shape.fill
                    if hasattr(template_fill, 'fore_color'):
                        try:
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = template_fill.fore_color.rgb
                        except:
                            pass

        except:
            pass

    def _apply_template_formatting(self, slide):
        """Apply template fonts and colors to slide content"""
        try:
            primary_font = self.learned_styles.get('primary_font', 'Calibri')
            primary_colors = self.learned_styles.get('primary_colors', [])

            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Apply template font
                            try:
                                run.font.name = primary_font
                            except:
                                pass

                            # Apply template color
                            if primary_colors:
                                try:
                                    color_hex = primary_colors[0].replace('#', '')
                                    if len(color_hex) == 6:
                                        rgb_color = RGBColor(
                                            int(color_hex[0:2], 16),
                                            int(color_hex[2:4], 16),
                                            int(color_hex[4:6], 16)
                                        )
                                        run.font.color.rgb = rgb_color
                                except:
                                    pass

        except:
            pass

    def get_recommended_slide_order(self, content_sections):
        """Get recommended slide order based on template patterns"""
        if not self.learned_styles or not self.learned_styles.get('slide_sequencing'):
            return content_sections

        template_sequence = self.learned_styles['slide_sequencing']

        # Analyze template patterns
        title_slides = [s for s in template_sequence if s['type'] in ['title', 'title_only']]
        content_slides = [s for s in template_sequence if s['type'] in ['title_and_content', 'content']]

        # Recommend structure based on template
        recommended_order = []

        # Always start with title slide
        recommended_order.append({'type': 'title', 'content': 'title_slide'})

        # If template has early summary/overview slides
        if len(template_sequence) > 1 and template_sequence[1]['type'] in ['title_and_content']:
            recommended_order.append({'type': 'overview', 'content': 'executive_summary'})

        # Add content sections based on template pattern
        for section in content_sections:
            if section == 'data_quality' and any(s['shape_count'] > 5 for s in content_slides):
                recommended_order.append({'type': 'detailed_content', 'content': section})
            else:
                recommended_order.append({'type': 'standard_content', 'content': section})

        # If template ends with summary/recommendations
        if len(template_sequence) > 2 and template_sequence[-1]['shape_count'] < 4:
            recommended_order.append({'type': 'summary', 'content': 'recommendations'})

        return recommended_order

class SAPDataProcessor:
    """SAP-style data processing engine"""

    def __init__(self):
        self.supported_formats = ['csv', 'xlsx', 'xls', 'pptx', 'docx']

    def extract_word_document_data(self, docx_file):
        """Extract comprehensive content from Word documents including text, tables, and metadata"""
        try:
            doc = Document(BytesIO(docx_file))
            extracted_data = {
                'document_structure': {
                    'paragraphs': [],
                    'tables': [],
                    'headers': [],
                    'footers': [],
                    'metadata': {}
                },
                'processed_content': {
                    'text_content': [],
                    'data_tables': [],
                    'key_points': [],
                    'decisions': [],
                    'metrics': []
                },
                'document_stats': {
                    'total_paragraphs': 0,
                    'total_tables': 0,
                    'total_words': 0,
                    'bullet_points': 0,
                    'numbered_lists': 0
                }
            }

            # Extract document metadata
            try:
                core_props = doc.core_properties
                extracted_data['document_structure']['metadata'] = {
                    'title': core_props.title or '',
                    'author': core_props.author or '',
                    'subject': core_props.subject or '',
                    'created': core_props.created.strftime('%Y-%m-%d %H:%M:%S') if core_props.created else '',
                    'modified': core_props.modified.strftime('%Y-%m-%d %H:%M:%S') if core_props.modified else '',
                    'last_modified_by': core_props.last_modified_by or ''
                }
            except:
                extracted_data['document_structure']['metadata'] = {}

            # Extract headers and footers
            for section in doc.sections:
                # Headers
                if section.header:
                    for paragraph in section.header.paragraphs:
                        if paragraph.text.strip():
                            extracted_data['document_structure']['headers'].append({
                                'text': paragraph.text.strip(),
                                'style': paragraph.style.name if paragraph.style else 'Normal'
                            })

                # Footers
                if section.footer:
                    for paragraph in section.footer.paragraphs:
                        if paragraph.text.strip():
                            extracted_data['document_structure']['footers'].append({
                                'text': paragraph.text.strip(),
                                'style': paragraph.style.name if paragraph.style else 'Normal'
                            })

            # Extract paragraphs with detailed analysis
            word_count = 0
            bullet_count = 0
            numbered_count = 0

            for para_idx, paragraph in enumerate(doc.paragraphs):
                if not paragraph.text.strip():
                    continue

                para_text = self._clean_paragraph_text(paragraph.text)
                word_count += len(para_text.split())

                para_data = {
                    'index': para_idx,
                    'text': para_text,
                    'style': paragraph.style.name if paragraph.style else 'Normal',
                    'level': getattr(paragraph, 'level', 0),
                    'is_bullet': self._is_bullet_point(paragraph),
                    'is_numbered': self._is_numbered_list(paragraph),
                    'is_heading': self._is_heading(paragraph),
                    'formatting': self._extract_paragraph_formatting(paragraph),
                    'content_type': self._classify_paragraph_content(para_text)
                }

                # Count bullet points and numbered lists
                if para_data['is_bullet']:
                    bullet_count += 1
                if para_data['is_numbered']:
                    numbered_count += 1

                extracted_data['document_structure']['paragraphs'].append(para_data)

                # Classify content for processing
                if para_data['content_type'] in ['key_point', 'decision']:
                    extracted_data['processed_content'][para_data['content_type'] + 's'].append(para_text)
                elif para_data['content_type'] == 'metric':
                    extracted_data['processed_content']['metrics'].append(para_text)
                else:
                    extracted_data['processed_content']['text_content'].append(para_text)

            # Extract tables with sophisticated parsing
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                has_header = False

                for row_idx, row in enumerate(table.rows):
                    row_data = []
                    for cell_idx, cell in enumerate(row.cells):
                        cell_text = self._clean_cell_text(cell.text)
                        row_data.append(cell_text)

                    table_data.append(row_data)

                # Detect if first row is header
                if table_data and len(table_data) > 1:
                    first_row = table_data[0]
                    if self._likely_table_header(first_row):
                        has_header = True

                # Process table into DataFrame if it contains data
                df_table = None
                if table_data and len(table_data) > (1 if has_header else 0):
                    try:
                        if has_header and len(table_data) > 1:
                            df_table = pd.DataFrame(table_data[1:], columns=table_data[0])
                        else:
                            df_table = pd.DataFrame(table_data)

                        # Clean and process DataFrame
                        df_table = self._clean_dataframe(df_table)

                        # Only add if it has meaningful data
                        if not df_table.empty and df_table.shape[1] > 1:
                            extracted_data['processed_content']['data_tables'].append({
                                'table_index': table_idx,
                                'dataframe': df_table,
                                'has_header': has_header,
                                'shape': df_table.shape,
                                'numeric_columns': len(df_table.select_dtypes(include=[np.number]).columns)
                            })

                    except Exception as e:
                        # If DataFrame creation fails, store as raw table
                        pass

                table_info = {
                    'table_index': table_idx,
                    'raw_data': table_data,
                    'rows': len(table_data),
                    'columns': len(table_data[0]) if table_data else 0,
                    'has_header': has_header,
                    'has_dataframe': df_table is not None and not df_table.empty
                }

                extracted_data['document_structure']['tables'].append(table_info)

            # Update document statistics
            extracted_data['document_stats'].update({
                'total_paragraphs': len(extracted_data['document_structure']['paragraphs']),
                'total_tables': len(extracted_data['document_structure']['tables']),
                'total_words': word_count,
                'bullet_points': bullet_count,
                'numbered_lists': numbered_count,
                'data_tables_found': len(extracted_data['processed_content']['data_tables']),
                'key_points_found': len(extracted_data['processed_content']['key_points']),
                'decisions_found': len(extracted_data['processed_content']['decisions']),
                'metrics_found': len(extracted_data['processed_content']['metrics'])
            })

            return extracted_data

        except Exception as e:
            st.error(f"Error extracting Word document data: {str(e)}")
            return None

    def _clean_paragraph_text(self, text):
        """Clean paragraph text by removing track changes and comments"""
        # Remove common track changes markers
        text = re.sub(r'\[.*?\]', '', text)  # Remove tracked deletions
        text = re.sub(r'\{.*?\}', '', text)  # Remove comments
        text = re.sub(r'<.*?>', '', text)    # Remove any XML-like tags

        # Clean up whitespace
        text = ' '.join(text.split())
        return text.strip()

    def _clean_cell_text(self, text):
        """Clean table cell text"""
        text = self._clean_paragraph_text(text)
        # Remove table-specific artifacts
        text = text.replace('\a', ' ')  # Remove vertical tab characters
        return text.strip()

    def _is_bullet_point(self, paragraph):
        """Detect if paragraph is a bullet point"""
        text = paragraph.text.strip()
        if not text:
            return False

        # Check for common bullet markers
        bullet_markers = ['•', '●', '○', '▪', '▫', '■', '□', '-', '*']
        if any(text.startswith(marker) for marker in bullet_markers):
            return True

        # Check paragraph style
        style_name = paragraph.style.name.lower() if paragraph.style else ''
        if 'bullet' in style_name or 'list' in style_name:
            return True

        return False

    def _is_numbered_list(self, paragraph):
        """Detect if paragraph is a numbered list item"""
        text = paragraph.text.strip()
        if not text:
            return False

        # Check for numbered patterns
        numbered_patterns = [
            r'^\d+\.',     # 1. 2. 3.
            r'^\d+\)',     # 1) 2) 3)
            r'^\(\d+\)',   # (1) (2) (3)
            r'^[a-z]\.',   # a. b. c.
            r'^[A-Z]\.',   # A. B. C.
            r'^[ivx]+\.',  # i. ii. iii.
            r'^[IVX]+\.'   # I. II. III.
        ]

        for pattern in numbered_patterns:
            if re.match(pattern, text):
                return True

        return False

    def _is_heading(self, paragraph):
        """Detect if paragraph is a heading"""
        if paragraph.style:
            style_name = paragraph.style.name.lower()
            if 'heading' in style_name or 'title' in style_name:
                return True
        return False

    def _extract_paragraph_formatting(self, paragraph):
        """Extract formatting information from paragraph"""
        formatting = {
            'is_bold': False,
            'is_italic': False,
            'font_size': None,
            'font_name': None
        }

        try:
            if paragraph.runs:
                first_run = paragraph.runs[0]
                formatting['is_bold'] = first_run.bold or False
                formatting['is_italic'] = first_run.italic or False
                if first_run.font.size:
                    formatting['font_size'] = first_run.font.size.pt
                formatting['font_name'] = first_run.font.name
        except:
            pass

        return formatting

    def _classify_paragraph_content(self, text):
        """Classify paragraph content type for better processing"""
        text_lower = text.lower()

        # Key point indicators
        key_indicators = ['key point', 'important', 'note:', 'remember', 'action item', 'takeaway']
        if any(indicator in text_lower for indicator in key_indicators):
            return 'key_point'

        # Decision indicators
        decision_indicators = ['decision', 'agreed', 'decided', 'resolved', 'concluded', 'approved']
        if any(indicator in text_lower for indicator in decision_indicators):
            return 'decision'

        # Metric indicators
        metric_patterns = [
            r'\d+%',           # Percentages
            r'\$[\d,]+',       # Dollar amounts
            r'\d+\.\d+',       # Decimal numbers
            r'\d{1,3}(,\d{3})*' # Large numbers with commas
        ]
        if any(re.search(pattern, text) for pattern in metric_patterns):
            return 'metric'

        return 'general_text'

    def _likely_table_header(self, row):
        """Determine if a table row is likely a header"""
        if not row:
            return False

        # Check if all cells have content
        if all(cell.strip() for cell in row):
            # Check for common header patterns
            header_indicators = ['name', 'date', 'amount', 'type', 'status', 'description', 'value', 'total']
            if any(any(indicator in cell.lower() for indicator in header_indicators) for cell in row):
                return True

        return False

    def _clean_dataframe(self, df):
        """Clean and optimize DataFrame from Word table"""
        if df.empty:
            return df

        # Remove completely empty rows and columns
        df = df.dropna(how='all').dropna(axis=1, how='all')

        # Try to convert numeric columns
        for col in df.columns:
            try:
                # Remove common non-numeric characters
                df[col] = df[col].astype(str).str.replace('$', '').str.replace(',', '').str.replace('%', '')
                # Try to convert to numeric
                df[col] = pd.to_numeric(df[col], errors='ignore')
            except:
                pass

        return df

    def process_multiple_files(self, uploaded_files):
        """Process multiple files and organize by type"""
        try:
            processed_data = {
                'csv_excel_files': [],
                'word_documents': [],
                'powerpoint_files': [],
                'combined_dataframes': [],
                'all_word_data': [],
                'all_pptx_data': [],
                'file_summary': {
                    'total_files': len(uploaded_files),
                    'csv_excel_count': 0,
                    'word_count': 0,
                    'pptx_count': 0,
                    'processing_errors': []
                }
            }

            for file_obj in uploaded_files:
                file_content = file_obj.getvalue()
                file_name = file_obj.name
                file_extension = file_name.split('.')[-1].lower()

                try:
                    if file_extension in ['csv', 'xlsx', 'xls']:
                        # Process CSV/Excel files
                        df = self.process_csv_excel(file_content, file_name)
                        if df is not None:
                            processed_data['csv_excel_files'].append({
                                'filename': file_name,
                                'dataframe': df,
                                'shape': df.shape,
                                'columns': list(df.columns)
                            })
                            processed_data['combined_dataframes'].append(df)
                            processed_data['file_summary']['csv_excel_count'] += 1

                    elif file_extension == 'docx':
                        # Process Word documents
                        docx_data = self.extract_word_document_data(file_content)
                        if docx_data:
                            docx_data['filename'] = file_name
                            processed_data['word_documents'].append(docx_data)
                            processed_data['all_word_data'].append(docx_data)
                            processed_data['file_summary']['word_count'] += 1

                    elif file_extension == 'pptx':
                        # Process PowerPoint files
                        pptx_data = self.extract_powerpoint_data(file_content)
                        if pptx_data:
                            pptx_data['filename'] = file_name
                            processed_data['powerpoint_files'].append(pptx_data)
                            processed_data['all_pptx_data'].append(pptx_data)
                            processed_data['file_summary']['pptx_count'] += 1

                except Exception as e:
                    processed_data['file_summary']['processing_errors'].append({
                        'filename': file_name,
                        'error': str(e)
                    })

            return processed_data

        except Exception as e:
            st.error(f"Error processing multiple files: {str(e)}")
            return None

    def combine_csv_excel_data(self, dataframes_list):
        """Combine multiple CSV/Excel DataFrames intelligently"""
        if not dataframes_list:
            return None

        try:
            # If only one DataFrame, return it
            if len(dataframes_list) == 1:
                return dataframes_list[0]

            # Try to concatenate if columns are similar
            combined_df = None

            # Check if all DataFrames have similar column structures
            first_cols = set(dataframes_list[0].columns)
            similar_structure = all(
                len(set(df.columns).intersection(first_cols)) / len(set(df.columns).union(first_cols)) > 0.5
                for df in dataframes_list[1:]
            )

            if similar_structure:
                # Concatenate DataFrames with similar structures
                combined_df = pd.concat(dataframes_list, ignore_index=True, sort=False)
                combined_df = combined_df.loc[:, ~combined_df.columns.duplicated()]
            else:
                # Merge DataFrames with different structures side by side
                combined_df = dataframes_list[0].copy()
                for i, df in enumerate(dataframes_list[1:], 1):
                    # Add suffix to avoid column name conflicts
                    df_suffixed = df.add_suffix(f'_file{i+1}')
                    # Concatenate horizontally
                    combined_df = pd.concat([combined_df, df_suffixed], axis=1)

            return combined_df

        except Exception as e:
            st.error(f"Error combining CSV/Excel data: {str(e)}")
            return dataframes_list[0] if dataframes_list else None

    def merge_word_document_insights(self, word_data_list):
        """Merge insights from multiple Word documents"""
        if not word_data_list:
            return None

        try:
            merged_insights = {
                'combined_stats': {
                    'total_documents': len(word_data_list),
                    'total_words': 0,
                    'total_tables': 0,
                    'total_paragraphs': 0,
                    'total_key_points': 0,
                    'total_decisions': 0,
                    'total_metrics': 0,
                    'data_tables_found': 0
                },
                'all_key_points': [],
                'all_decisions': [],
                'all_metrics': [],
                'all_text_content': [],
                'combined_data_tables': [],
                'document_summaries': [],
                'authors': set(),
                'creation_dates': []
            }

            for doc_data in word_data_list:
                filename = doc_data.get('filename', 'Unknown')
                doc_stats = doc_data['document_stats']

                # Aggregate statistics
                merged_insights['combined_stats']['total_words'] += doc_stats['total_words']
                merged_insights['combined_stats']['total_tables'] += doc_stats['total_tables']
                merged_insights['combined_stats']['total_paragraphs'] += doc_stats['total_paragraphs']
                merged_insights['combined_stats']['total_key_points'] += doc_stats['key_points_found']
                merged_insights['combined_stats']['total_decisions'] += doc_stats['decisions_found']
                merged_insights['combined_stats']['total_metrics'] += doc_stats['metrics_found']
                merged_insights['combined_stats']['data_tables_found'] += doc_stats['data_tables_found']

                # Collect content with source information
                processed_content = doc_data['processed_content']

                for point in processed_content['key_points']:
                    merged_insights['all_key_points'].append({
                        'content': point,
                        'source': filename
                    })

                for decision in processed_content['decisions']:
                    merged_insights['all_decisions'].append({
                        'content': decision,
                        'source': filename
                    })

                for metric in processed_content['metrics']:
                    merged_insights['all_metrics'].append({
                        'content': metric,
                        'source': filename
                    })

                for text in processed_content['text_content'][:10]:  # Limit text content
                    merged_insights['all_text_content'].append({
                        'content': text,
                        'source': filename
                    })

                # Collect data tables
                for table_data in processed_content['data_tables']:
                    table_info = table_data.copy()
                    table_info['source'] = filename
                    merged_insights['combined_data_tables'].append(table_info)

                # Document metadata
                metadata = doc_data['document_structure']['metadata']
                if metadata.get('author'):
                    merged_insights['authors'].add(metadata['author'])
                if metadata.get('created'):
                    merged_insights['creation_dates'].append(metadata['created'])

                # Create document summary
                merged_insights['document_summaries'].append({
                    'filename': filename,
                    'words': doc_stats['total_words'],
                    'tables': doc_stats['total_tables'],
                    'key_points': doc_stats['key_points_found'],
                    'decisions': doc_stats['decisions_found'],
                    'author': metadata.get('author', 'Unknown')
                })

            # Convert sets to lists for JSON serialization
            merged_insights['authors'] = list(merged_insights['authors'])

            # Create PowerPoint-compatible structure
            powerpoint_compatible = {
                'document_stats': {
                    'total_words': merged_insights['combined_stats']['total_words'],
                    'total_tables': merged_insights['combined_stats']['total_tables'],
                    'total_paragraphs': merged_insights['combined_stats']['total_paragraphs'],
                    'key_points_found': merged_insights['combined_stats']['total_key_points'],
                    'decisions_found': merged_insights['combined_stats']['total_decisions'],
                    'metrics_found': merged_insights['combined_stats']['total_metrics'],
                    'data_tables_found': merged_insights['combined_stats']['data_tables_found'],
                    'bullet_points': merged_insights['combined_stats'].get('bullet_points', 0),
                    'numbered_lists': merged_insights['combined_stats'].get('numbered_lists', 0)
                },
                'processed_content': {
                    'key_points': [item['content'] for item in merged_insights['all_key_points']],
                    'decisions': [item['content'] for item in merged_insights['all_decisions']],
                    'metrics': [item['content'] for item in merged_insights['all_metrics']],
                    'text_content': [item['content'] for item in merged_insights['all_text_content']],
                    'data_tables': merged_insights['combined_data_tables']
                },
                'document_structure': {
                    'metadata': {
                        'authors': ', '.join(merged_insights['authors']) if merged_insights['authors'] else '',
                        'creation_dates': merged_insights['creation_dates'],
                        'total_documents': merged_insights['combined_stats']['total_documents']
                    },
                    'headers': [],
                    'footers': []
                },
                # Keep the original merged structure for backwards compatibility
                'merged_insights': merged_insights
            }

            return powerpoint_compatible

        except Exception as e:
            st.error(f"Error merging Word document insights: {str(e)}")
            return None

    def merge_powerpoint_data(self, pptx_data_list):
        """Merge data from multiple PowerPoint presentations"""
        if not pptx_data_list:
            return None

        try:
            merged_pptx = {
                'combined_stats': {
                    'total_presentations': len(pptx_data_list),
                    'total_slides': 0,
                    'total_text_blocks': 0,
                    'total_tables': 0,
                    'total_words': 0
                },
                'all_slides': [],
                'all_text_content': [],
                'all_tables': [],
                'presentation_summaries': []
            }

            for pptx_data in pptx_data_list:
                filename = pptx_data.get('filename', 'Unknown')

                # Aggregate statistics
                merged_pptx['combined_stats']['total_slides'] += pptx_data['slide_count']
                merged_pptx['combined_stats']['total_text_blocks'] += len(pptx_data['text_content'])
                merged_pptx['combined_stats']['total_tables'] += len(pptx_data['tables'])

                # Calculate total words
                total_words = sum(len(text.split()) for text in pptx_data['text_content'])
                merged_pptx['combined_stats']['total_words'] += total_words

                # Collect content with source information
                for slide in pptx_data['slides']:
                    slide_copy = slide.copy()
                    slide_copy['source'] = filename
                    merged_pptx['all_slides'].append(slide_copy)

                for text in pptx_data['text_content']:
                    merged_pptx['all_text_content'].append({
                        'content': text,
                        'source': filename
                    })

                for table in pptx_data['tables']:
                    table_copy = table.copy()
                    table_copy['source'] = filename
                    merged_pptx['all_tables'].append(table_copy)

                # Create presentation summary
                merged_pptx['presentation_summaries'].append({
                    'filename': filename,
                    'slides': pptx_data['slide_count'],
                    'text_blocks': len(pptx_data['text_content']),
                    'tables': len(pptx_data['tables']),
                    'words': total_words
                })

            return merged_pptx

        except Exception as e:
            st.error(f"Error merging PowerPoint data: {str(e)}")
            return None

    def extract_powerpoint_data(self, pptx_file):
        """Extract text and table data from PowerPoint presentations"""
        try:
            prs = Presentation(BytesIO(pptx_file))
            extracted_data = {
                'slides': [],
                'text_content': [],
                'tables': [],
                'slide_count': len(prs.slides)
            }

            for i, slide in enumerate(prs.slides):
                slide_data = {
                    'slide_number': i + 1,
                    'title': '',
                    'text_content': [],
                    'tables': []
                }

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape == slide.shapes.title:
                            slide_data['title'] = shape.text.strip()
                        else:
                            slide_data['text_content'].append(shape.text.strip())

                    # Extract tables
                    if shape.shape_type == 19:  # Table shape
                        try:
                            table_data = []
                            table = shape.table
                            for row_idx, row in enumerate(table.rows):
                                row_data = []
                                for cell in row.cells:
                                    row_data.append(cell.text.strip())
                                table_data.append(row_data)

                            if table_data:
                                slide_data['tables'].append(table_data)
                                extracted_data['tables'].append({
                                    'slide': i + 1,
                                    'data': table_data
                                })
                        except:
                            pass

                extracted_data['slides'].append(slide_data)
                extracted_data['text_content'].extend(slide_data['text_content'])

            return extracted_data

        except Exception as e:
            st.error(f"Error extracting PowerPoint data: {str(e)}")
            return None

    def process_csv_excel(self, file_content, file_name):
        """Process CSV and Excel files"""
        try:
            if file_name.lower().endswith('.csv'):
                # Try different encodings for CSV
                encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
                for encoding in encodings:
                    try:
                        df = pd.read_csv(BytesIO(file_content), encoding=encoding)
                        return df
                    except UnicodeDecodeError:
                        continue
                raise ValueError("Could not decode CSV file with any supported encoding")
            else:
                # Excel files
                df = pd.read_excel(BytesIO(file_content))
                return df

        except Exception as e:
            st.error(f"Error processing file: {str(e)}")
            return None

    def analyze_data_quality(self, df):
        """Analyze data quality with SAP-style metrics"""
        if df is None or df.empty:
            return None

        total_cells = len(df) * len(df.columns)
        missing_cells = df.isnull().sum().sum()
        duplicate_rows = df.duplicated().sum()

        quality_metrics = {
            'total_records': len(df),
            'total_columns': len(df.columns),
            'completeness_pct': ((total_cells - missing_cells) / total_cells * 100) if total_cells > 0 else 0,
            'duplicate_pct': (duplicate_rows / len(df) * 100) if len(df) > 0 else 0,
            'numeric_columns': len(df.select_dtypes(include=[np.number]).columns),
            'text_columns': len(df.select_dtypes(include=['object', 'string']).columns),
            'date_columns': len(df.select_dtypes(include=['datetime64']).columns)
        }

        # Quality assessment
        if quality_metrics['completeness_pct'] >= 95 and quality_metrics['duplicate_pct'] <= 1:
            quality_metrics['quality_level'] = 'Excellent'
            quality_metrics['quality_color'] = 'success'
        elif quality_metrics['completeness_pct'] >= 85 and quality_metrics['duplicate_pct'] <= 5:
            quality_metrics['quality_level'] = 'Good'
            quality_metrics['quality_color'] = 'warning'
        else:
            quality_metrics['quality_level'] = 'Needs Improvement'
            quality_metrics['quality_color'] = 'error'

        return quality_metrics

    def generate_basic_insights(self, df):
        """Generate basic data insights"""
        if df is None or df.empty:
            return None

        insights = {
            'summary_stats': {},
            'column_insights': [],
            'data_types': df.dtypes.to_dict()
        }

        # Numeric column analysis
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        if len(numeric_cols) > 0:
            insights['summary_stats'] = df[numeric_cols].describe().round(2)

        # Column-level insights
        for col in df.columns:
            col_insight = {
                'column': col,
                'type': str(df[col].dtype),
                'missing_count': df[col].isnull().sum(),
                'missing_pct': (df[col].isnull().sum() / len(df) * 100),
                'unique_values': df[col].nunique()
            }

            if df[col].dtype in ['object', 'string']:
                col_insight['top_values'] = df[col].value_counts().head(3).to_dict()
            elif np.issubdtype(df[col].dtype, np.number):
                col_insight['min_value'] = df[col].min()
                col_insight['max_value'] = df[col].max()
                col_insight['mean_value'] = df[col].mean()

            insights['column_insights'].append(col_insight)

        return insights

# Initialize SAP processor and template analyzer
sap_processor = SAPDataProcessor()

# Initialize template analyzer in session state to persist across reruns
if 'template_analyzer' not in st.session_state:
    st.session_state['template_analyzer'] = SAPTemplateAnalyzer()
template_analyzer = st.session_state['template_analyzer']

# Sidebar configuration
with st.sidebar:
    st.markdown("### � SAP Configuration Panel")

    # File processing options
    st.markdown("#### File Processing")
    auto_detect_types = st.checkbox("Auto-detect data types", value=True)
    remove_duplicates = st.checkbox("Remove duplicate records", value=False)
    handle_missing = st.selectbox(
        "Missing data handling",
        ["Keep as-is", "Remove rows", "Fill with mean", "Fill with median"],
        index=0
    )

    # Visualization options
    st.markdown("#### Visualization")
    chart_theme = st.selectbox("Chart theme", ["SAP Corporate", "Modern", "Classic"], index=0)
    include_charts = st.checkbox("Include charts in report", value=True)
    max_categories = st.slider("Maximum categories in charts", 5, 20, 10)

    # Template learning section
    st.markdown("#### 🎨 Template Learning")
    use_template = st.checkbox("Enable template matching", value=False, help="Use a custom PowerPoint template for report generation")

    template_file = None
    if use_template:
        template_file = st.file_uploader(
            "Upload template PowerPoint",
            type=['pptx'],
            help="Upload a sample PowerPoint with your preferred layout and styling",
            key="template_uploader"
        )

        if template_file and st.button("🔍 Learn from Template", key="learn_template"):
            with st.spinner("Analyzing template..."):
                template_data = template_analyzer.analyze_template(template_file.getvalue())
                if template_data:
                    st.success(f"✅ Template learned! Found {template_data['slide_count']} slides")
                    st.session_state['template_learned'] = True
                    st.session_state['template_data'] = template_data

    # Report options
    st.markdown("#### Report Generation")
    report_title = st.text_input("Report title", value="SAP Data Analysis Report")
    include_summary = st.checkbox("Include executive summary", value=True)
    include_quality_assessment = st.checkbox("Include quality assessment", value=True)

    st.markdown("---")
    st.markdown("### =� SAP System Status")
    st.markdown('<div class="sap-status-success"> All systems operational</div>', unsafe_allow_html=True)
    st.markdown(f"=P **Last update:** {datetime.now().strftime('%H:%M:%S')}")

# Main content area
st.markdown("### =� File Upload & Processing")

# File upload section
col1, col2 = st.columns([3, 1])

with col1:
    uploaded_files = st.file_uploader(
        "Upload your business data files",
        type=['csv', 'xlsx', 'xls', 'pptx', 'docx'],
        help="Supported formats: CSV, Excel (.xlsx, .xls), PowerPoint (.pptx), Word (.docx). Upload multiple files for combined analysis.",
        accept_multiple_files=True
    )

with col2:
    if uploaded_files:
        total_size = sum(len(f.getvalue()) for f in uploaded_files) / (1024 * 1024)  # MB
        file_types = list(set(f.name.split('.')[-1].upper() for f in uploaded_files))

        st.markdown('<div class="sap-card">', unsafe_allow_html=True)
        st.markdown(f'<div class="sap-metric-value">{len(uploaded_files)}</div>', unsafe_allow_html=True)
        st.markdown(f'<div class="sap-metric-label">Files</div>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

        st.markdown('<div class="sap-card">', unsafe_allow_html=True)
        st.markdown(f'<div class="sap-metric-value">{total_size:.1f} MB</div>', unsafe_allow_html=True)
        st.markdown(f'<div class="sap-metric-label">Total Size</div>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

        st.markdown('<div class="sap-card">', unsafe_allow_html=True)
        formats_text = ", ".join(file_types) if len(file_types) <= 2 else f"{len(file_types)} types"
        st.markdown(f'<div class="sap-metric-value">{formats_text}</div>', unsafe_allow_html=True)
        st.markdown(f'<div class="sap-metric-label">Formats</div>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

# Template preview section
if 'template_learned' in st.session_state and st.session_state['template_learned']:
    st.markdown("### 🎨 Template Preview")

    # Get template summary
    template_summary = template_analyzer.get_template_summary()

    if template_summary:
        st.markdown('<div class="sap-card">', unsafe_allow_html=True)

        col1, col2, col3, col4 = st.columns(4)

        with col1:
            st.markdown('<div class="sap-metric">', unsafe_allow_html=True)
            st.markdown(f'<div class="sap-metric-value">{template_summary["total_slides"]}</div>', unsafe_allow_html=True)
            st.markdown('<div class="sap-metric-label">Template Slides</div>', unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)

        with col2:
            st.markdown('<div class="sap-metric">', unsafe_allow_html=True)
            st.markdown(f'<div class="sap-metric-value">{template_summary["unique_layouts"]}</div>', unsafe_allow_html=True)
            st.markdown('<div class="sap-metric-label">Layout Types</div>', unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)

        with col3:
            st.markdown('<div class="sap-metric">', unsafe_allow_html=True)
            st.markdown(f'<div class="sap-metric-value">{len(template_summary["dominant_fonts"])}</div>', unsafe_allow_html=True)
            st.markdown('<div class="sap-metric-label">Font Styles</div>', unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)

        with col4:
            st.markdown('<div class="sap-metric">', unsafe_allow_html=True)
            st.markdown(f'<div class="sap-metric-value">{template_summary["avg_shapes_per_slide"]:.1f}</div>', unsafe_allow_html=True)
            st.markdown('<div class="sap-metric-label">Avg Shapes/Slide</div>', unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)

        st.markdown('</div>', unsafe_allow_html=True)

        # Show template details
        with st.expander("🔍 Template Analysis Details", expanded=False):
            col1, col2 = st.columns(2)

            with col1:
                st.markdown("**Dominant Fonts:**")
                for font in template_summary['dominant_fonts']:
                    st.write(f"• {font}")

            with col2:
                st.markdown("**Color Scheme:**")
                for color in template_summary['dominant_colors']:
                    st.markdown(f'<span style="color: {color}; font-weight: bold;">● {color}</span>', unsafe_allow_html=True)

        # Show additional template analysis
        if template_summary.get('content_areas_detected', 0) > 0:
            st.markdown("**Template Structure Analysis:**")
            col1, col2 = st.columns(2)
            with col1:
                st.write(f"📍 Content Areas Detected: {template_summary['content_areas_detected']}")
            with col2:
                st.write(f"🎯 Bullet Styles Found: {template_summary['bullet_styles_found']}")

            # Show slide sequencing if available
            if hasattr(template_analyzer, 'learned_styles') and template_analyzer.learned_styles.get('slide_sequencing'):
                sequence = template_analyzer.learned_styles['slide_sequencing']
                st.markdown("**Template Slide Sequence:**")
                sequence_text = " → ".join([f"{s['type']}" for s in sequence[:5]])
                if len(sequence) > 5:
                    sequence_text += f" ... ({len(sequence)} total slides)"
                st.write(sequence_text)

        st.info("🎯 **Template Active:** Generated reports will match this template's styling and layout patterns.")

def create_sap_visualization(df, chart_type="overview"):
    """Create SAP-style visualizations"""
    if df is None or df.empty:
        return None

    numeric_cols = df.select_dtypes(include=[np.number]).columns

    if chart_type == "overview" and len(numeric_cols) > 0:
        # Create overview dashboard
        fig = make_subplots(
            rows=2, cols=2,
            subplot_titles=('Data Distribution', 'Correlation Matrix', 'Missing Values', 'Data Types'),
            specs=[[{"type": "histogram"}, {"type": "heatmap"}],
                   [{"type": "bar"}, {"type": "pie"}]]
        )

        # Data distribution (first numeric column)
        if len(numeric_cols) > 0:
            fig.add_trace(
                go.Histogram(x=df[numeric_cols[0]], name=numeric_cols[0], marker_color='#0070f2'),
                row=1, col=1
            )

        # Correlation matrix
        if len(numeric_cols) > 1:
            corr_matrix = df[numeric_cols].corr()
            fig.add_trace(
                go.Heatmap(
                    z=corr_matrix.values,
                    x=corr_matrix.columns,
                    y=corr_matrix.columns,
                    colorscale='RdBu',
                    zmid=0
                ),
                row=1, col=2
            )

        # Missing values
        missing_data = df.isnull().sum()
        missing_data = missing_data[missing_data > 0]
        if len(missing_data) > 0:
            fig.add_trace(
                go.Bar(x=missing_data.values, y=missing_data.index,
                       orientation='h', marker_color='#ff6600'),
                row=2, col=1
            )

        # Data types
        type_counts = df.dtypes.value_counts()
        fig.add_trace(
            go.Pie(labels=type_counts.index.astype(str), values=type_counts.values,
                   marker_colors=['#0070f2', '#30914c', '#ff6600', '#cc1919'][:len(type_counts)]),
            row=2, col=2
        )

        fig.update_layout(
            height=600,
            showlegend=False,
            title_text="SAP Data Overview Dashboard",
            title_font_size=16
        )

        return fig

    return None

def aggregate_multi_source_data(df, insights, pptx_data=None, docx_data=None):
    """Aggregate and correlate data from all input sources for unified executive insights"""
    try:
        aggregated = {
            'executive_summary': {
                'total_data_sources': 0,
                'total_data_points': 0,
                'total_content_items': 0,
                'data_quality_score': 0,
                'key_metrics': [],
                'main_insights': [],
                'recommendations': []
            },
            'data_overview': {
                'csv_excel_summary': None,
                'document_summary': None,
                'presentation_summary': None
            },
            'cross_source_insights': {
                'common_themes': [],
                'data_correlations': [],
                'content_patterns': [],
                'decision_alignment': []
            },
            'unified_metrics': {
                'total_records': 0,
                'total_words': 0,
                'total_slides': 0,
                'total_tables': 0,
                'total_decisions': 0,
                'total_key_points': 0
            },
            # Add PowerPoint-compatible structures
            'document_stats': {
                'total_words': 0,
                'total_tables': 0,
                'total_paragraphs': 0,
                'key_points_found': 0,
                'decisions_found': 0,
                'metrics_found': 0,
                'data_tables_found': 0,
                'bullet_points': 0,
                'numbered_lists': 0
            },
            'processed_content': {
                'key_points': [],
                'decisions': [],
                'metrics': [],
                'text_content': [],
                'data_tables': []
            },
            'document_structure': {
                'metadata': {},
                'headers': [],
                'footers': []
            }
        }

        source_count = 0

        # Process CSV/Excel data
        if df is not None and insights:
            source_count += 1
            aggregated['unified_metrics']['total_records'] = len(df)
            aggregated['unified_metrics']['total_tables'] += 1

            if insights.get('summary_stats') is not None and not insights['summary_stats'].empty:
                aggregated['data_overview']['csv_excel_summary'] = {
                    'shape': df.shape,
                    'numeric_columns': len(df.select_dtypes(include=[np.number]).columns),
                    'quality_metrics': insights.get('quality_metrics', {}),
                    'key_columns': list(df.columns)[:5]
                }

                # Extract key metrics from data
                numeric_cols = df.select_dtypes(include=[np.number]).columns
                if len(numeric_cols) > 0:
                    for col in numeric_cols[:3]:  # Top 3 numeric columns
                        if df[col].notna().sum() > 0:
                            aggregated['executive_summary']['key_metrics'].append({
                                'metric': col,
                                'value': f"{df[col].mean():.2f}" if df[col].dtype in ['float64', 'int64'] else str(df[col].iloc[0]),
                                'source': 'Data Analysis',
                                'type': 'numeric'
                            })

        # Process Word document data
        if docx_data:
            if isinstance(docx_data, dict) and 'combined_stats' in docx_data:
                # Multiple documents merged
                source_count += docx_data['combined_stats']['total_documents']
                aggregated['unified_metrics']['total_words'] = docx_data['combined_stats']['total_words']
                aggregated['unified_metrics']['total_decisions'] = docx_data['combined_stats']['total_decisions']
                aggregated['unified_metrics']['total_key_points'] = docx_data['combined_stats']['total_key_points']
                aggregated['unified_metrics']['total_tables'] += docx_data['combined_stats']['data_tables_found']

                # Update document_stats for PowerPoint compatibility
                aggregated['document_stats']['total_words'] = docx_data['combined_stats']['total_words']
                aggregated['document_stats']['total_tables'] = docx_data['combined_stats']['data_tables_found']
                aggregated['document_stats']['total_paragraphs'] = docx_data['combined_stats']['total_paragraphs']
                aggregated['document_stats']['key_points_found'] = docx_data['combined_stats']['total_key_points']
                aggregated['document_stats']['decisions_found'] = docx_data['combined_stats']['total_decisions']
                aggregated['document_stats']['metrics_found'] = docx_data['combined_stats']['total_metrics']
                aggregated['document_stats']['data_tables_found'] = docx_data['combined_stats']['data_tables_found']
                aggregated['document_stats']['bullet_points'] = docx_data['combined_stats'].get('bullet_points', 0)
                aggregated['document_stats']['numbered_lists'] = docx_data['combined_stats'].get('numbered_lists', 0)

                # Update processed_content for PowerPoint compatibility
                aggregated['processed_content']['key_points'] = [item['content'] for item in docx_data.get('all_key_points', [])]
                aggregated['processed_content']['decisions'] = [item['content'] for item in docx_data.get('all_decisions', [])]
                aggregated['processed_content']['metrics'] = [item['content'] for item in docx_data.get('all_metrics', [])]
                aggregated['processed_content']['text_content'] = [item['content'] for item in docx_data.get('all_text_content', [])]
                aggregated['processed_content']['data_tables'] = docx_data.get('combined_data_tables', [])

                # Update document_structure for PowerPoint compatibility
                if docx_data.get('authors'):
                    aggregated['document_structure']['metadata']['authors'] = ', '.join(docx_data['authors'])
                if docx_data.get('creation_dates'):
                    aggregated['document_structure']['metadata']['creation_dates'] = docx_data['creation_dates']

                aggregated['data_overview']['document_summary'] = {
                    'document_count': docx_data['combined_stats']['total_documents'],
                    'total_words': docx_data['combined_stats']['total_words'],
                    'key_points': len(docx_data.get('all_key_points', [])),
                    'decisions': len(docx_data.get('all_decisions', [])),
                    'authors': docx_data.get('authors', [])
                }

                # Extract top insights from documents
                for point_data in docx_data.get('all_key_points', [])[:5]:
                    aggregated['executive_summary']['main_insights'].append({
                        'content': point_data['content'],
                        'source': f"Document: {point_data['source']}",
                        'type': 'key_point'
                    })

                for decision_data in docx_data.get('all_decisions', [])[:3]:
                    aggregated['executive_summary']['main_insights'].append({
                        'content': decision_data['content'],
                        'source': f"Document: {decision_data['source']}",
                        'type': 'decision'
                    })

            elif isinstance(docx_data, dict) and 'document_stats' in docx_data:
                # Single document
                source_count += 1
                doc_stats = docx_data['document_stats']
                aggregated['unified_metrics']['total_words'] = doc_stats['total_words']
                aggregated['unified_metrics']['total_decisions'] = doc_stats['decisions_found']
                aggregated['unified_metrics']['total_key_points'] = doc_stats['key_points_found']
                aggregated['unified_metrics']['total_tables'] += doc_stats['data_tables_found']

                # Copy document_stats directly for single document
                aggregated['document_stats'] = doc_stats.copy()

                # Copy processed_content for single document
                if 'processed_content' in docx_data:
                    aggregated['processed_content'] = docx_data['processed_content'].copy()

                # Copy document_structure for single document
                if 'document_structure' in docx_data:
                    aggregated['document_structure'] = docx_data['document_structure'].copy()

        # Process PowerPoint data
        if pptx_data:
            if isinstance(pptx_data, dict) and 'combined_stats' in pptx_data:
                # Multiple presentations merged
                source_count += pptx_data['combined_stats']['total_presentations']
                aggregated['unified_metrics']['total_slides'] = pptx_data['combined_stats']['total_slides']
                aggregated['unified_metrics']['total_words'] += pptx_data['combined_stats']['total_words']
                aggregated['unified_metrics']['total_tables'] += pptx_data['combined_stats']['total_tables']

                aggregated['data_overview']['presentation_summary'] = {
                    'presentation_count': pptx_data['combined_stats']['total_presentations'],
                    'total_slides': pptx_data['combined_stats']['total_slides'],
                    'total_words': pptx_data['combined_stats']['total_words'],
                    'total_tables': pptx_data['combined_stats']['total_tables']
                }

            elif isinstance(pptx_data, dict) and 'slide_count' in pptx_data:
                # Single presentation
                source_count += 1
                aggregated['unified_metrics']['total_slides'] = pptx_data['slide_count']
                total_words = sum(len(text.split()) for text in pptx_data.get('text_content', []))
                aggregated['unified_metrics']['total_words'] += total_words
                aggregated['unified_metrics']['total_tables'] += len(pptx_data.get('tables', []))

        # Calculate overall metrics
        aggregated['executive_summary']['total_data_sources'] = source_count
        aggregated['executive_summary']['total_data_points'] = (
            aggregated['unified_metrics']['total_records'] +
            aggregated['unified_metrics']['total_words'] +
            aggregated['unified_metrics']['total_slides']
        )
        aggregated['executive_summary']['total_content_items'] = (
            aggregated['unified_metrics']['total_key_points'] +
            aggregated['unified_metrics']['total_decisions'] +
            aggregated['unified_metrics']['total_tables']
        )

        # Calculate data quality score
        quality_factors = []
        if df is not None and insights and insights.get('quality_metrics'):
            quality_factors.append(insights['quality_metrics'].get('completeness_pct', 0))
        if aggregated['unified_metrics']['total_words'] > 0:
            quality_factors.append(min(100, aggregated['unified_metrics']['total_words'] / 100))  # Word richness
        total_content_items = (
            aggregated['unified_metrics'].get('total_key_points', 0) +
            aggregated['unified_metrics'].get('total_decisions', 0) +
            aggregated['unified_metrics'].get('total_tables', 0)
        )
        if total_content_items > 0:
            quality_factors.append(min(100, total_content_items * 10))  # Content richness

        aggregated['executive_summary']['data_quality_score'] = sum(quality_factors) / len(quality_factors) if quality_factors else 0

        # Generate unified recommendations
        recommendations = []
        if aggregated['unified_metrics']['total_records'] > 1000:
            recommendations.append("Consider implementing automated data processing workflows")
        if aggregated['unified_metrics']['total_decisions'] > 5:
            recommendations.append("Establish decision tracking and follow-up processes")
        if aggregated['unified_metrics']['total_key_points'] > 10:
            recommendations.append("Create centralized knowledge management system")
        if source_count > 3:
            recommendations.append("Implement unified data governance across all sources")

        aggregated['executive_summary']['recommendations'] = recommendations

        # Identify cross-source patterns
        if source_count > 1:
            themes = []
            if aggregated['unified_metrics']['total_decisions'] > 0 and aggregated['unified_metrics']['total_records'] > 0:
                themes.append("Data-driven decision making is evident across sources")
            if aggregated['unified_metrics']['total_tables'] > 2:
                themes.append("Structured data analysis is a common theme")
            if aggregated['unified_metrics']['total_words'] > 5000:
                themes.append("Rich content documentation supports business processes")

            aggregated['cross_source_insights']['common_themes'] = themes

        return aggregated

    except Exception as e:
        st.error(f"Error aggregating multi-source data: {str(e)}")
        return None

def generate_sap_powerpoint_report(df, insights, pptx_data=None, docx_data=None):
    """Generate unified executive SAP-style PowerPoint report aggregating all input sources"""
    try:
        # Check if template should be used
        use_template_styling = ('template_learned' in st.session_state and
                              st.session_state['template_learned'] and
                              hasattr(template_analyzer, 'template_presentation') and
                              template_analyzer.template_presentation)

        # Debug messaging
        if 'template_learned' in st.session_state and st.session_state['template_learned']:
            if hasattr(template_analyzer, 'template_presentation') and template_analyzer.template_presentation:
                st.info(f"✅ Using template with {len(template_analyzer.template_presentation.slides)} slides")
            else:
                st.warning("⚠️ Template learned but presentation not found - using default styling")
        else:
            st.info("ℹ️ No template loaded - using default styling")

        if use_template_styling:
            # CLONE the template presentation instead of creating new
            try:
                # Save template to temporary buffer and reload it (creates a copy)
                template_buffer = BytesIO()
                template_analyzer.template_presentation.save(template_buffer)
                template_buffer.seek(0)

                # Load the template as our base presentation
                prs = Presentation(template_buffer)
                st.success("Successfully cloned template presentation")

                # Store reference to original template slides for duplication
                template_slides = list(prs.slides)

                # Function to duplicate a slide from template
                def duplicate_slide(source_slide_index=1):
                    """Duplicate a slide from template (default to second slide for content)"""
                    try:
                        if source_slide_index < len(template_slides):
                            source_slide = template_slides[source_slide_index]

                            # Get the slide layout
                            slide_layout = source_slide.slide_layout

                            # Add new slide with same layout
                            new_slide = prs.slides.add_slide(slide_layout)

                            # Copy all shapes except placeholders (which will be filled with new content)
                            for shape in source_slide.shapes:
                                if not (hasattr(shape, 'is_placeholder') and shape.is_placeholder):
                                    # Copy non-placeholder shapes (logos, images, design elements)
                                    if shape.shape_type == 13:  # Picture
                                        # Copy image
                                        try:
                                            image_part = shape.image.image_part
                                            new_slide.shapes.add_picture(
                                                BytesIO(image_part.blob),
                                                shape.left, shape.top, shape.width, shape.height
                                            )
                                        except:
                                            pass
                                    elif hasattr(shape, 'shape_type'):
                                        # Copy other shapes (rectangles, etc.)
                                        try:
                                            if shape.shape_type == 1:  # AutoShape
                                                new_shape = new_slide.shapes.add_shape(
                                                    MSO_SHAPE.RECTANGLE,
                                                    shape.left, shape.top, shape.width, shape.height
                                                )
                                                # Copy fill if available
                                                if hasattr(shape, 'fill') and hasattr(new_shape, 'fill'):
                                                    if hasattr(shape.fill, 'fore_color'):
                                                        new_shape.fill.solid()
                                                        new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
                                        except:
                                            pass

                            return new_slide
                        else:
                            # Fallback to basic slide if template doesn't have enough slides
                            return prs.slides.add_slide(prs.slide_layouts[1])
                    except:
                        return prs.slides.add_slide(prs.slide_layouts[1])

                # Function to replace text in slide placeholders while preserving formatting
                def replace_slide_content(slide, title_text, content_lines):
                    """Replace content in slide placeholders while keeping all formatting"""
                    try:
                        # Find and update title
                        if hasattr(slide, 'shapes') and slide.shapes.title:
                            title_shape = slide.shapes.title
                            if hasattr(title_shape, 'text_frame'):
                                # Clear existing text but keep formatting
                                title_shape.text_frame.clear()
                                p = title_shape.text_frame.paragraphs[0]
                                p.text = title_text

                        # Find content placeholder and update
                        for shape in slide.shapes:
                            if (hasattr(shape, 'is_placeholder') and shape.is_placeholder and
                                hasattr(shape, 'placeholder_format') and
                                shape.placeholder_format.type == 2):  # Content placeholder

                                if hasattr(shape, 'text_frame'):
                                    # Clear existing content but preserve formatting
                                    shape.text_frame.clear()

                                    # Add new content as bullets
                                    for i, line in enumerate(content_lines):
                                        if i == 0:
                                            p = shape.text_frame.paragraphs[0]
                                        else:
                                            p = shape.text_frame.add_paragraph()

                                        p.text = line
                                        p.level = 0  # Top level bullet

                                break  # Only update first content placeholder
                    except Exception as e:
                        pass

            except Exception as e:
                st.error(f"Error cloning template: {str(e)}")
                prs = Presentation()
                use_template_styling = False
        else:
            prs = Presentation()

        # Define colors (use template colors if available, otherwise SAP defaults)
        if use_template_styling and template_analyzer.learned_styles['primary_colors']:
            try:
                # Convert hex colors to RGB
                def hex_to_rgb(hex_color):
                    hex_color = hex_color.lstrip('#')
                    return RGBColor(*tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4)))

                primary_colors = template_analyzer.learned_styles['primary_colors']
                primary_color = hex_to_rgb(primary_colors[0]) if primary_colors else RGBColor(0, 112, 242)
                secondary_color = hex_to_rgb(primary_colors[1]) if len(primary_colors) > 1 else RGBColor(255, 102, 0)
                accent_color = hex_to_rgb(primary_colors[2]) if len(primary_colors) > 2 else RGBColor(106, 109, 112)
            except:
                # Fallback to SAP colors if template colors can't be parsed
                primary_color = RGBColor(0, 112, 242)
                secondary_color = RGBColor(255, 102, 0)
                accent_color = RGBColor(106, 109, 112)
        else:
            # Use SAP default colors
            primary_color = RGBColor(0, 112, 242)
            secondary_color = RGBColor(255, 102, 0)
            accent_color = RGBColor(106, 109, 112)

        # Define primary font (use template font if available)
        primary_font = (template_analyzer.learned_styles['primary_font']
                       if use_template_styling and template_analyzer.learned_styles['primary_font']
                       else 'Calibri')

        # Aggregate data from all sources for unified analysis
        aggregated_insights = aggregate_multi_source_data(df, insights, pptx_data, docx_data)

        # Configure report settings based on aggregated data
        if aggregated_insights:
            source_count = aggregated_insights['executive_summary']['total_data_sources']
            total_points = aggregated_insights['executive_summary']['total_data_points']
            report_title = f"Executive Summary Report - {source_count} Data Source{'s' if source_count != 1 else ''}"
        else:
            report_title = "SAP Enterprise Analytics Report"
            total_points = 0

        include_summary = True
        include_data_overview = bool(df is not None)
        include_document_analysis = bool(docx_data)
        include_presentation_analysis = bool(pptx_data)

        # Title slide - use first slide from template if available
        if use_template_styling and len(prs.slides) > 0:
            # Modify existing template title slide
            slide = prs.slides[0]
            subtitle_text = f"Unified Business Intelligence Analysis\n{total_points:,} Data Points | {datetime.now().strftime('%B %d, %Y')}"

            if 'replace_slide_content' in locals():
                # Use template content replacement
                replace_slide_content(slide, report_title, [subtitle_text])
            else:
                # Manual replacement as fallback
                if hasattr(slide, 'shapes') and slide.shapes.title:
                    slide.shapes.title.text = report_title
                # Try to find subtitle placeholder
                for shape in slide.shapes:
                    if (hasattr(shape, 'is_placeholder') and shape.is_placeholder and
                        hasattr(shape, 'placeholder_format') and
                        shape.placeholder_format.type in [3, 4]):  # Subtitle placeholders
                        if hasattr(shape, 'text_frame'):
                            shape.text_frame.clear()
                            shape.text_frame.paragraphs[0].text = subtitle_text
                        break
        else:
            # Fallback to creating new title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)

            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = report_title
            subtitle.text = f"Unified Business Intelligence Analysis\n{total_points:,} Data Points | {datetime.now().strftime('%B %d, %Y')}"

            # Format title with manual styling
            title_paragraph = title.text_frame.paragraphs[0]
            title_paragraph.font.color.rgb = primary_color
            title_paragraph.font.size = Pt(44)
            title_paragraph.font.bold = True
            title_paragraph.font.name = primary_font

        # Unified Executive Summary slide
        if include_summary and aggregated_insights:
            if use_template_styling and 'duplicate_slide' in locals():
                # Use template cloning approach
                slide = duplicate_slide(1)  # Duplicate second slide from template for content
            else:
                # Fallback to old approach
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)

            # Prepare content for template replacement
            exec_summary = aggregated_insights.get('executive_summary', {})
            unified_metrics = aggregated_insights.get('unified_metrics', {})
            data_sources = exec_summary.get('total_data_sources', 0)

            title_text = "Executive Summary"
            content_lines = [
                f"Unified Analysis of {data_sources} Data Source{'s' if data_sources != 1 else ''}",
                f"Total Data Points Analyzed: {exec_summary.get('total_data_points', 0):,}",
                f"Content Items Extracted: {exec_summary.get('total_content_items', 0)}",
                f"Overall Quality Score: {exec_summary.get('data_quality_score', 0):.1f}%"
            ]

            if use_template_styling and 'replace_slide_content' in locals():
                # Use template content replacement (preserves all formatting)
                replace_slide_content(slide, title_text, content_lines)
            else:
                # Fallback to manual formatting
                title = slide.shapes.title
                title.text = title_text
                title_paragraph = title.text_frame.paragraphs[0]
                title_paragraph.font.color.rgb = primary_color
                title_paragraph.font.name = primary_font

                body = slide.placeholders[1]
                tf = body.text_frame
                if "tf" in locals() and tf: tf.text = content_lines[0]

                for line in content_lines[1:]:
                    if "tf" in locals() and tf: p = tf.add_paragraph()
                    p.text = f"• {line}"
                    p.level = 1

            # Detailed breakdown
            total_records = unified_metrics.get('total_records', 0)
            if total_records > 0:
                if "tf" in locals() and tf: p = tf.add_paragraph()
                p.text = f"• Structured Data Records: {total_records:,}"
                p.level = 1

            total_words = unified_metrics.get('total_words', 0)
            if total_words > 0:
                if "tf" in locals() and tf: p = tf.add_paragraph()
                p.text = f"• Document Content: {total_words:,} words"
                p.level = 1

            total_slides = unified_metrics.get('total_slides', 0)
            if total_slides > 0:
                if "tf" in locals() and tf: p = tf.add_paragraph()
                p.text = f"• Presentation Content: {total_slides} slides"
                p.level = 1

            total_decisions = unified_metrics.get('total_decisions', 0)
            if total_decisions > 0:
                if "tf" in locals() and tf: p = tf.add_paragraph()
                p.text = f"• Business Decisions Identified: {total_decisions}"
                p.level = 1

            total_key_points = unified_metrics.get('total_key_points', 0)
            if total_key_points > 0:
                if "tf" in locals() and tf: p = tf.add_paragraph()
                p.text = f"• Key Insights Extracted: {total_key_points}"
                p.level = 1

        # Key Insights and Recommendations slide
        if (aggregated_insights and
            (aggregated_insights.get('executive_summary', {}).get('main_insights', []) or
             aggregated_insights.get('executive_summary', {}).get('recommendations', []))):

            if use_template_styling and 'duplicate_slide' in locals():
                # Use template cloning approach
                slide = duplicate_slide(1)  # Duplicate second slide from template for content
            else:
                # Fallback to old approach
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)

            # Prepare content for template replacement
            title_text = "Key Insights & Recommendations"
            content_lines = []

            # Main insights
            if aggregated_insights['executive_summary']['main_insights']:
                content_lines.append("Strategic Insights from Analysis")
                for insight in aggregated_insights['executive_summary']['main_insights'][:5]:
                    content_lines.append(insight['content'][:120] + ('...' if len(insight['content']) > 120 else ''))

            # Recommendations
            if aggregated_insights['executive_summary']['recommendations']:
                content_lines.append("Recommended Actions")
                for rec in aggregated_insights['executive_summary']['recommendations']:
                    content_lines.append(rec)

            if use_template_styling and 'replace_slide_content' in locals():
                # Use template content replacement (preserves all formatting)
                replace_slide_content(slide, title_text, content_lines)
            else:
                # Fallback to manual formatting
                title = slide.shapes.title
                title.text = title_text
                title_paragraph = title.text_frame.paragraphs[0]
                title_paragraph.font.color.rgb = primary_color
                title_paragraph.font.name = primary_font

                body = slide.placeholders[1]
                tf = body.text_frame

                if content_lines:
                    if "tf" in locals() and tf: tf.text = content_lines[0]

                    for line in content_lines[1:]:
                        if "tf" in locals() and tf: p = tf.add_paragraph()
                        p.text = f"• {line}"
                        p.level = 1

        # Cross-Source Analysis slide
        if (aggregated_insights and
            aggregated_insights.get('cross_source_insights', {}).get('common_themes', [])):

            if use_template_styling and 'duplicate_slide' in locals():
                # Use template cloning approach
                slide = duplicate_slide(1)  # Duplicate second slide from template for content
            else:
                # Fallback to old approach
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)

            # Prepare content for template replacement
            title_text = "Cross-Source Analysis"
            content_lines = ["Common Themes Across Data Sources"]

            for theme in aggregated_insights['cross_source_insights']['common_themes']:
                content_lines.append(theme)

            # Add data source overview to content
            if aggregated_insights['data_overview']:
                content_lines.append("Data Source Breakdown")

                if aggregated_insights['data_overview']['csv_excel_summary']:
                    csv_summary = aggregated_insights['data_overview']['csv_excel_summary']
                    content_lines.append(f"Structured Data: {csv_summary['shape'][0]:,} records, {csv_summary['shape'][1]} columns")

                if aggregated_insights['data_overview']['document_summary']:
                    doc_summary = aggregated_insights['data_overview']['document_summary']
                    content_lines.append(f"Documents: {doc_summary['document_count']} files, {doc_summary['total_words']:,} words")

                if aggregated_insights['data_overview']['presentation_summary']:
                    pres_summary = aggregated_insights['data_overview']['presentation_summary']
                    content_lines.append(f"Presentations: {pres_summary['presentation_count']} files, {pres_summary['total_slides']} slides")

            if use_template_styling and 'replace_slide_content' in locals():
                # Use template content replacement (preserves all formatting)
                replace_slide_content(slide, title_text, content_lines)
            else:
                # Fallback to manual formatting
                title = slide.shapes.title
                title.text = title_text
                title_paragraph = title.text_frame.paragraphs[0]
                title_paragraph.font.color.rgb = primary_color
                title_paragraph.font.name = primary_font

                body = slide.placeholders[1]
                tf = body.text_frame
                if "tf" in locals() and tf: tf.text = content_lines[0]

                for line in content_lines[1:]:
                    if "tf" in locals() and tf: p = tf.add_paragraph()
                    p.text = f"• {line}"
                    p.level = 1

        # Configure additional slide flags
        include_quality_assessment = bool(df is not None)

        # Data Quality slide
        if include_quality_assessment and df is not None:
            quality_metrics = sap_processor.analyze_data_quality(df)
            if quality_metrics:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)

                # Apply template structure
                if use_template_styling:
                    slide = template_analyzer.apply_template_structure(slide, 'title_and_content')

                title = slide.shapes.title
                body = slide.placeholders[1]

                title.text = "Data Quality Assessment"
                title_paragraph = title.text_frame.paragraphs[0]
                title_paragraph.font.color.rgb = primary_color
                title_paragraph.font.name = primary_font

                tf = body.text_frame
                if "tf" in locals() and tf: tf.text = f"Overall Quality: {quality_metrics['quality_level']}"

                metrics_text = [
                    f"Data Completeness: {quality_metrics['completeness_pct']:.1f}%",
                    f"Duplicate Records: {quality_metrics['duplicate_pct']:.1f}%",
                    f"Numeric Columns: {quality_metrics['numeric_columns']}",
                    f"Text Columns: {quality_metrics['text_columns']}",
                    f"Total Records: {quality_metrics['total_records']:,}"
                ]

                for metric in metrics_text:
                    if "tf" in locals() and tf: p = tf.add_paragraph()
                    p.text = f"• {metric}"
                    p.level = 1

        # PowerPoint Data slide (if PowerPoint was uploaded)
        if pptx_data:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)

            # Apply template structure
            if use_template_styling:
                slide = template_analyzer.apply_template_structure(slide, 'title_and_content')

            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = "PowerPoint Analysis Results"
            title_paragraph = title.text_frame.paragraphs[0]
            title_paragraph.font.color.rgb = primary_color
            title_paragraph.font.name = primary_font

            tf = body.text_frame
            if "tf" in locals() and tf: tf.text = f"Processed {pptx_data['slide_count']} slides"

            if "tf" in locals() and tf: p = tf.add_paragraph()
            p.text = f"• Text blocks extracted: {len(pptx_data['text_content'])}"
            p.level = 1

            if "tf" in locals() and tf: p = tf.add_paragraph()
            p.text = f"• Tables found: {len(pptx_data['tables'])}"
            p.level = 1

            # Show sample text content
            if pptx_data['text_content']:
                if "tf" in locals() and tf: p = tf.add_paragraph()
                p.text = "• Sample content:"
                p.level = 1

                for i, text in enumerate(pptx_data['text_content'][:3]):
                    if text.strip():
                        if "tf" in locals() and tf: p = tf.add_paragraph()
                        p.text = f"  - {text[:100]}{'...' if len(text) > 100 else ''}"
                        p.level = 2

        # Word Document Data slide (if Word document was uploaded)
        if docx_data:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)

            # Apply template structure
            if use_template_styling:
                slide = template_analyzer.apply_template_structure(slide, 'title_and_content')

            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = "Word Document Analysis Results"
            title_paragraph = title.text_frame.paragraphs[0]
            title_paragraph.font.color.rgb = primary_color
            title_paragraph.font.name = primary_font

            tf = body.text_frame
            # Use aggregated insights if available, otherwise fall back to docx_data
            if aggregated_insights and 'document_stats' in aggregated_insights:
                doc_stats = aggregated_insights['document_stats']
            else:
                doc_stats = docx_data.get('document_stats', {})

            total_words = doc_stats.get('total_words', 0)
            if "tf" in locals() and tf: tf.text = f"Processed document with {total_words:,} words"

            if "tf" in locals() and tf: p = tf.add_paragraph()
            p.text = f"• Tables extracted: {doc_stats.get('total_tables', 0)}"
            p.level = 1

            if "tf" in locals() and tf: p = tf.add_paragraph()
            p.text = f"• Data tables found: {doc_stats.get('data_tables_found', 0)}"
            p.level = 1

            if "tf" in locals() and tf: p = tf.add_paragraph()
            p.text = f"• Key points identified: {doc_stats.get('key_points_found', 0)}"
            p.level = 1

            if "tf" in locals() and tf: p = tf.add_paragraph()
            p.text = f"• Decisions captured: {doc_stats.get('decisions_found', 0)}"
            p.level = 1

            if "tf" in locals() and tf: p = tf.add_paragraph()
            p.text = f"• Metrics found: {doc_stats.get('metrics_found', 0)}"
            p.level = 1

            # Document metadata slide
            # Check for document structure in both original data and aggregated data
            document_metadata = {}
            if docx_data and docx_data.get('document_structure', {}).get('metadata', {}):
                document_metadata = docx_data.get('document_structure', {}).get('metadata', {})
            elif aggregated_insights and aggregated_insights.get('document_structure', {}).get('metadata', {}):
                document_metadata = aggregated_insights.get('document_structure', {}).get('metadata', {})

            if document_metadata:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)

                # Apply template structure
                if use_template_styling:
                    slide = template_analyzer.apply_template_structure(slide, 'title_and_content')

                title = slide.shapes.title
                body = slide.placeholders[1]

                title.text = "Document Information"
                title_paragraph = title.text_frame.paragraphs[0]
                title_paragraph.font.color.rgb = primary_color
                title_paragraph.font.name = primary_font

                tf = body.text_frame
                metadata = document_metadata

                if "tf" in locals() and tf: tf.text = "Document Properties"

                if metadata.get('title'):
                    if "tf" in locals() and tf: p = tf.add_paragraph()
                    p.text = f"• Title: {metadata['title']}"
                    p.level = 1

                if metadata.get('author') or metadata.get('authors'):
                    author_text = metadata.get('author') or metadata.get('authors', '')
                    if "tf" in locals() and tf: p = tf.add_paragraph()
                    p.text = f"• Author(s): {author_text}"
                    p.level = 1

                if metadata.get('created'):
                    if "tf" in locals() and tf: p = tf.add_paragraph()
                    p.text = f"• Created: {metadata['created']}"
                    p.level = 1

                if metadata.get('subject'):
                    if "tf" in locals() and tf: p = tf.add_paragraph()
                    p.text = f"• Subject: {metadata['subject']}"
                    p.level = 1

                if metadata.get('total_documents'):
                    if "tf" in locals() and tf: p = tf.add_paragraph()
                    p.text = f"• Documents Processed: {metadata['total_documents']}"
                    p.level = 1

            # Key findings slide from Word document
            if (docx_data and
                (docx_data.get('processed_content', {}).get('key_points', []) or
                 docx_data.get('processed_content', {}).get('decisions', []))):
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)

                # Apply template structure
                if use_template_styling:
                    slide = template_analyzer.apply_template_structure(slide, 'title_and_content')

                title = slide.shapes.title
                body = slide.placeholders[1]

                title.text = "Key Findings from Document"
                title_paragraph = title.text_frame.paragraphs[0]
                title_paragraph.font.color.rgb = primary_color
                title_paragraph.font.name = primary_font

                tf = body.text_frame
                if "tf" in locals() and tf: tf.text = "Important Points and Decisions"

                # Add key points
                key_points = docx_data.get('processed_content', {}).get('key_points', [])[:5]
                if key_points:
                    if "tf" in locals() and tf: p = tf.add_paragraph()
                    p.text = "Key Points:"
                    p.level = 1

                    for point in key_points:
                        if "tf" in locals() and tf: p = tf.add_paragraph()
                        p.text = f"• {point[:150]}{'...' if len(point) > 150 else ''}"
                        p.level = 2

                # Add decisions
                decisions = docx_data.get('processed_content', {}).get('decisions', [])[:3]
                if decisions:
                    if "tf" in locals() and tf: p = tf.add_paragraph()
                    p.text = "Decisions Made:"
                    p.level = 1

                    for decision in decisions:
                        if "tf" in locals() and tf: p = tf.add_paragraph()
                        p.text = f"• {decision[:150]}{'...' if len(decision) > 150 else ''}"
                        p.level = 2

            # Data tables slide from Word document
            if (docx_data and
                docx_data.get('processed_content', {}).get('data_tables', [])):
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)

                # Apply template structure
                if use_template_styling:
                    slide = template_analyzer.apply_template_structure(slide, 'title_and_content')

                title = slide.shapes.title
                body = slide.placeholders[1]

                title.text = "Extracted Data Tables"
                title_paragraph = title.text_frame.paragraphs[0]
                title_paragraph.font.color.rgb = primary_color
                title_paragraph.font.name = primary_font

                tf = body.text_frame
                data_tables = docx_data.get('processed_content', {}).get('data_tables', [])
                if "tf" in locals() and tf: tf.text = f"Found {len(data_tables)} analyzable data tables"

                for i, table_info in enumerate(data_tables[:3]):  # Show first 3 tables
                    if "tf" in locals() and tf: p = tf.add_paragraph()
                    shape = table_info['shape']
                    numeric_cols = table_info['numeric_columns']
                    p.text = f"• Table {i+1}: {shape[0]} rows × {shape[1]} columns ({numeric_cols} numeric)"
                    p.level = 1

                    # Show sample data from largest table
                    if i == 0 and table_info['dataframe'] is not None:
                        df_sample = table_info['dataframe']
                        if not df_sample.empty:
                            if "tf" in locals() and tf: p = tf.add_paragraph()
                            p.text = "Sample columns:"
                            p.level = 1

                            for col in list(df_sample.columns)[:4]:  # Show first 4 columns
                                if "tf" in locals() and tf: p = tf.add_paragraph()
                                p.text = f"  - {col}"
                                p.level = 2

        # Data Overview slide
        if df is not None and insights:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)

            # Apply template structure
            if use_template_styling:
                slide = template_analyzer.apply_template_structure(slide, 'title_and_content')

            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = "Data Structure Overview"
            title_paragraph = title.text_frame.paragraphs[0]
            title_paragraph.font.color.rgb = primary_color
            title_paragraph.font.name = primary_font

            tf = body.text_frame
            if "tf" in locals() and tf: tf.text = "Column Analysis"

            # Show top columns with insights
            for col_insight in insights['column_insights'][:8]:
                if "tf" in locals() and tf: p = tf.add_paragraph()
                p.text = f"• {col_insight['column']} ({col_insight['type']})"
                p.level = 1

                if col_insight['missing_pct'] > 0:
                    if "tf" in locals() and tf: p = tf.add_paragraph()
                    p.text = f"  Missing: {col_insight['missing_pct']:.1f}%"
                    p.level = 2

        # Recommendations slide
        recommendations_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(recommendations_layout)

        # Apply template structure
        if use_template_styling:
            slide = template_analyzer.apply_template_structure(slide, 'title_and_content')

        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "SAP Recommendations"
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.color.rgb = primary_color
        title_paragraph.font.name = primary_font

        tf = body.text_frame
        if "tf" in locals() and tf: tf.text = "Next Steps for Data Processing"

        recommendations = [
            "Implement data quality monitoring",
            "Establish automated data validation",
            "Create standardized reporting templates",
            "Set up regular data governance reviews",
            "Consider SAP integration opportunities"
        ]

        for rec in recommendations:
            if "tf" in locals() and tf: p = tf.add_paragraph()
            p.text = f"• {rec}"
            p.level = 1

        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)

        return ppt_buffer

    except Exception as e:
        st.error(f"Error generating PowerPoint report: {str(e)}")
        return None

# Main processing logic
if uploaded_files:
    with st.spinner("Processing files..."):
        # Process multiple files
        processed_data = sap_processor.process_multiple_files(uploaded_files)

        if processed_data:
            # Display file processing summary
            st.markdown("### 📊 Multi-File Analysis Dashboard")

            # File summary metrics
            file_summary = processed_data['file_summary']
            col1, col2, col3, col4 = st.columns(4)

            with col1:
                st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                st.markdown(f'<div class="sap-metric-value">{file_summary["total_files"]}</div>', unsafe_allow_html=True)
                st.markdown('<div class="sap-metric-label">Total Files</div>', unsafe_allow_html=True)
                st.markdown('</div>', unsafe_allow_html=True)

            with col2:
                st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                st.markdown(f'<div class="sap-metric-value">{file_summary["csv_excel_count"]}</div>', unsafe_allow_html=True)
                st.markdown('<div class="sap-metric-label">Data Files</div>', unsafe_allow_html=True)
                st.markdown('</div>', unsafe_allow_html=True)

            with col3:
                st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                st.markdown(f'<div class="sap-metric-value">{file_summary["word_count"]}</div>', unsafe_allow_html=True)
                st.markdown('<div class="sap-metric-label">Word Docs</div>', unsafe_allow_html=True)
                st.markdown('</div>', unsafe_allow_html=True)

            with col4:
                st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                st.markdown(f'<div class="sap-metric-value">{file_summary["pptx_count"]}</div>', unsafe_allow_html=True)
                st.markdown('<div class="sap-metric-label">Presentations</div>', unsafe_allow_html=True)
                st.markdown('</div>', unsafe_allow_html=True)

            # Show processing errors if any
            if file_summary['processing_errors']:
                st.warning(f"⚠️ {len(file_summary['processing_errors'])} files had processing errors")
                with st.expander("View processing errors", expanded=False):
                    for error in file_summary['processing_errors']:
                        st.error(f"**{error['filename']}**: {error['error']}")

            # Process and combine CSV/Excel data
            combined_df = None
            if processed_data['combined_dataframes']:
                combined_df = sap_processor.combine_csv_excel_data(processed_data['combined_dataframes'])

                if combined_df is not None:
                    st.markdown("#### 📈 Combined Data Analysis")

                    # Data processing options for combined data
                    if handle_missing == "Remove rows":
                        combined_df = combined_df.dropna()
                    elif handle_missing == "Fill with mean":
                        numeric_cols = combined_df.select_dtypes(include=[np.number]).columns
                        combined_df[numeric_cols] = combined_df[numeric_cols].fillna(combined_df[numeric_cols].mean())
                    elif handle_missing == "Fill with median":
                        numeric_cols = combined_df.select_dtypes(include=[np.number]).columns
                        combined_df[numeric_cols] = combined_df[numeric_cols].fillna(combined_df[numeric_cols].median())

                    if remove_duplicates:
                        combined_df = combined_df.drop_duplicates()

                    # Analyze combined data
                    quality_metrics = sap_processor.analyze_data_quality(combined_df)
                    insights = sap_processor.generate_basic_insights(combined_df)

                    # Display combined data quality metrics
                    if quality_metrics:
                        col1, col2, col3, col4 = st.columns(4)

                        with col1:
                            st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                            st.markdown(f'<div class="sap-metric-value">{quality_metrics["total_records"]:,}</div>', unsafe_allow_html=True)
                            st.markdown('<div class="sap-metric-label">Combined Records</div>', unsafe_allow_html=True)
                            st.markdown('</div>', unsafe_allow_html=True)

                        with col2:
                            st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                            st.markdown(f'<div class="sap-metric-value">{quality_metrics["completeness_pct"]:.1f}%</div>', unsafe_allow_html=True)
                            st.markdown('<div class="sap-metric-label">Complete</div>', unsafe_allow_html=True)
                            st.markdown('</div>', unsafe_allow_html=True)

                        with col3:
                            st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                            st.markdown(f'<div class="sap-metric-value">{quality_metrics["numeric_columns"]}</div>', unsafe_allow_html=True)
                            st.markdown('<div class="sap-metric-label">Numeric Cols</div>', unsafe_allow_html=True)
                            st.markdown('</div>', unsafe_allow_html=True)

                        with col4:
                            quality_status = f'sap-status-{quality_metrics["quality_color"]}'
                            st.markdown(f'<div class="{quality_status}">{quality_metrics["quality_level"]}</div>', unsafe_allow_html=True)

                    # Combined data preview
                    st.markdown("##### 📊 Combined Data Preview")
                    st.dataframe(combined_df.head(10), use_container_width=True)

                    # Individual file details
                    with st.expander("📋 Individual File Details", expanded=False):
                        for file_data in processed_data['csv_excel_files']:
                            st.markdown(f"**{file_data['filename']}** - Shape: {file_data['shape']}")
                            st.write(f"Columns: {', '.join(file_data['columns'][:5])}{'...' if len(file_data['columns']) > 5 else ''}")

                    # Summary statistics for combined data
                    if insights and 'summary_stats' in insights and not insights['summary_stats'].empty:
                        st.markdown("##### 📊 Combined Summary Statistics")
                        st.dataframe(insights['summary_stats'], use_container_width=True)

                    # Visualizations for combined data
                    if include_charts:
                        st.markdown("##### 📊 Combined Data Visualizations")
                        fig = create_sap_visualization(combined_df)
                        if fig:
                            st.plotly_chart(fig, use_container_width=True)

            # Process and merge Word document insights
            merged_word_insights = None
            if processed_data['all_word_data']:
                merged_word_insights = sap_processor.merge_word_document_insights(processed_data['all_word_data'])

                if merged_word_insights:
                    st.markdown("#### 📄 Combined Word Document Analysis")

                    # Word document summary metrics
                    combined_stats = merged_word_insights.get('merged_insights', {}).get('combined_stats', {})
                    col1, col2, col3, col4 = st.columns(4)

                    with col1:
                        st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                        st.markdown(f'<div class="sap-metric-value">{combined_stats["total_documents"]}</div>', unsafe_allow_html=True)
                        st.markdown('<div class="sap-metric-label">Documents</div>', unsafe_allow_html=True)
                        st.markdown('</div>', unsafe_allow_html=True)

                    with col2:
                        st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                        st.markdown(f'<div class="sap-metric-value">{combined_stats["total_words"]:,}</div>', unsafe_allow_html=True)
                        st.markdown('<div class="sap-metric-label">Total Words</div>', unsafe_allow_html=True)
                        st.markdown('</div>', unsafe_allow_html=True)

                    with col3:
                        st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                        st.markdown(f'<div class="sap-metric-value">{combined_stats["total_key_points"]}</div>', unsafe_allow_html=True)
                        st.markdown('<div class="sap-metric-label">Key Points</div>', unsafe_allow_html=True)
                        st.markdown('</div>', unsafe_allow_html=True)

                    with col4:
                        st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                        st.markdown(f'<div class="sap-metric-value">{combined_stats["total_decisions"]}</div>', unsafe_allow_html=True)
                        st.markdown('<div class="sap-metric-label">Decisions</div>', unsafe_allow_html=True)
                        st.markdown('</div>', unsafe_allow_html=True)

                    # Document summaries
                    with st.expander("📋 Document Summaries", expanded=False):
                        for doc_summary in merged_word_insights.get('merged_insights', {}).get('document_summaries', []):
                            st.markdown(f"**{doc_summary['filename']}**")
                            col1, col2, col3 = st.columns(3)
                            with col1:
                                st.write(f"Words: {doc_summary['words']:,}")
                            with col2:
                                st.write(f"Tables: {doc_summary['tables']}")
                            with col3:
                                st.write(f"Author: {doc_summary['author']}")

                    # Combined insights
                    st.markdown("##### 🔍 Combined Content Analysis")
                    tab1, tab2, tab3 = st.tabs(["Key Points", "Decisions", "Metrics"])

                    with tab1:
                        if merged_word_insights.get('merged_insights', {}).get('all_key_points', []):
                            for i, point_data in enumerate(merged_word_insights.get('merged_insights', {}).get('all_key_points', [])[:15]):
                                st.write(f"• **[{point_data['source']}]** {point_data['content']}")
                        else:
                            st.info("No key points found across documents.")

                    with tab2:
                        if merged_word_insights.get('merged_insights', {}).get('all_decisions', []):
                            for i, decision_data in enumerate(merged_word_insights.get('merged_insights', {}).get('all_decisions', [])[:15]):
                                st.write(f"• **[{decision_data['source']}]** {decision_data['content']}")
                        else:
                            st.info("No decisions found across documents.")

                    with tab3:
                        if merged_word_insights.get('merged_insights', {}).get('all_metrics', []):
                            for i, metric_data in enumerate(merged_word_insights.get('merged_insights', {}).get('all_metrics', [])[:15]):
                                st.write(f"• **[{metric_data['source']}]** {metric_data['content']}")
                        else:
                            st.info("No metrics found across documents.")

                    # Combined data tables from Word documents
                    if merged_word_insights.get('merged_insights', {}).get('combined_data_tables', []):
                        st.markdown("##### 📊 Extracted Data Tables from Documents")
                        for i, table_data in enumerate(merged_word_insights.get('merged_insights', {}).get('combined_data_tables', [])[:5]):
                            st.markdown(f"**Table from {table_data['source']}** (Shape: {table_data['shape']})")
                            st.dataframe(table_data['dataframe'], use_container_width=True)

            # Process and merge PowerPoint data
            merged_pptx_data = None
            if processed_data['all_pptx_data']:
                merged_pptx_data = sap_processor.merge_powerpoint_data(processed_data['all_pptx_data'])

                if merged_pptx_data:
                    st.markdown("#### 📊 Combined PowerPoint Analysis")

                    # PowerPoint summary metrics
                    combined_stats = merged_pptx_data['combined_stats']
                    col1, col2, col3, col4 = st.columns(4)

                    with col1:
                        st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                        st.markdown(f'<div class="sap-metric-value">{combined_stats["total_presentations"]}</div>', unsafe_allow_html=True)
                        st.markdown('<div class="sap-metric-label">Presentations</div>', unsafe_allow_html=True)
                        st.markdown('</div>', unsafe_allow_html=True)

                    with col2:
                        st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                        st.markdown(f'<div class="sap-metric-value">{combined_stats["total_slides"]}</div>', unsafe_allow_html=True)
                        st.markdown('<div class="sap-metric-label">Total Slides</div>', unsafe_allow_html=True)
                        st.markdown('</div>', unsafe_allow_html=True)

                    with col3:
                        st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                        st.markdown(f'<div class="sap-metric-value">{combined_stats["total_words"]:,}</div>', unsafe_allow_html=True)
                        st.markdown('<div class="sap-metric-label">Total Words</div>', unsafe_allow_html=True)
                        st.markdown('</div>', unsafe_allow_html=True)

                    with col4:
                        st.markdown('<div class="sap-card">', unsafe_allow_html=True)
                        st.markdown(f'<div class="sap-metric-value">{combined_stats["total_tables"]}</div>', unsafe_allow_html=True)
                        st.markdown('<div class="sap-metric-label">Total Tables</div>', unsafe_allow_html=True)
                        st.markdown('</div>', unsafe_allow_html=True)

                    # Presentation summaries
                    with st.expander("📋 Presentation Summaries", expanded=False):
                        for pres_summary in merged_pptx_data['presentation_summaries']:
                            st.markdown(f"**{pres_summary['filename']}**")
                            col1, col2, col3 = st.columns(3)
                            with col1:
                                st.write(f"Slides: {pres_summary['slides']}")
                            with col2:
                                st.write(f"Text Blocks: {pres_summary['text_blocks']}")
                            with col3:
                                st.write(f"Tables: {pres_summary['tables']}")

            # Generate comprehensive report
            st.markdown("#### 📊 Generate Comprehensive Analysis Report")

            if st.button("🚀 Generate Multi-File SAP Report", type="primary", key="multi_file_report"):
                with st.spinner("Generating comprehensive multi-file SAP report..."):
                    # Create comprehensive insights object
                    comprehensive_insights = {
                        'file_summary': file_summary,
                        'combined_data_insights': insights if 'insights' in locals() else None,
                        'word_insights': merged_word_insights,
                        'pptx_insights': merged_pptx_data,
                        'authors': merged_word_insights.get('merged_insights', {}).get('authors', []) if merged_word_insights else [],
                        'total_files_processed': file_summary['total_files']
                    }

                    # Generate the multi-file report
                    report_buffer = generate_sap_powerpoint_report(
                        combined_df,
                        comprehensive_insights,
                        merged_pptx_data,
                        merged_word_insights
                    )

                    if report_buffer:
                        st.success("✅ Comprehensive multi-file SAP report generated successfully!")
                        st.download_button(
                            label="📎 Download Multi-File SAP Analysis Report",
                            data=report_buffer.getvalue(),
                            file_name=f"SAP_MultiFile_Analysis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )

        else:
            st.error("❌ No files were successfully processed. Please check your files and try again.")


# Footer
st.markdown("---")
st.markdown("""
<div style="text-align: center; color: #6a6d70; font-size: 0.9rem; padding: 1rem;">
    🚀 <strong>Universal Automation Platform - SAP Edition</strong> |
    Enterprise Multi-File Data Processing & Analytics |
    Powered by SAP Corporate Standards
</div>
""", unsafe_allow_html=True)
